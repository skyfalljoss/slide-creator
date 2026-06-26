import os
import time
from io import BytesIO
from pathlib import Path

from fastapi import HTTPException
from httpx import AsyncClient, ASGITransport
from pptx import Presentation
import pytest

from app.config import settings
from app.main import app, purge_local_temp_files
from app.routers import export as export_router, generate as generate_router, refine as refine_router, uploads
from app.models.schemas import ChartRecommendation, SlideData
from app import dependencies
from app.services.generation import providers
from app.routers.uploads import upload_file


@pytest.fixture
async def client():
    transport = ASGITransport(app=app, raise_app_exceptions=False)
    async with AsyncClient(transport=transport, base_url="http://test") as ac:
        yield ac
        app.dependency_overrides.clear()


def delete_upload(file_id: str):
    path = Path(settings.local_upload_dir) / file_id
    if path.exists():
        path.unlink()
    metadata_path = Path(settings.local_upload_dir) / f"{file_id}.json"
    if metadata_path.exists():
        metadata_path.unlink()


def _patch_generator(monkeypatch: pytest.MonkeyPatch, fake):
    dependencies.get_generator_service.cache_clear()
    app.dependency_overrides[dependencies.get_generator_service] = lambda: fake


def _patch_storage(monkeypatch: pytest.MonkeyPatch, fake):
    dependencies.get_storage_service.cache_clear()
    app.dependency_overrides[dependencies.get_storage_service] = lambda: fake


@pytest.mark.asyncio
async def test_health(client: AsyncClient):
    resp = await client.get("/api/v1/health")
    assert resp.status_code == 200
    assert resp.json()["status"] == "ok"


@pytest.mark.asyncio
async def test_generate_sales(client: AsyncClient):
    resp = await client.post(
        "/api/v1/generate",
        json={"prompt": "Pitch deck for Acme Corp", "deck_type": "sales_9"},
    )
    assert resp.status_code == 200
    data = resp.json()
    assert "session_id" in data
    assert len(data["slides"]) == 9
    assert data["slides"][-1]["title"] == "Thank You"


@pytest.mark.asyncio
async def test_generate_internal(client: AsyncClient):
    resp = await client.post(
        "/api/v1/generate",
        json={"prompt": "Internal strategy review", "deck_type": "internal_6"},
    )
    assert resp.status_code == 200
    data = resp.json()
    assert len(data["slides"]) == 6


@pytest.mark.asyncio
async def test_generate_records_user_id_from_header(client: AsyncClient):
    audit = providers.get_audit_service()
    audit.clear_events()

    try:
        resp = await client.post(
            "/api/v1/generate",
            headers={"x-user-id": "banker-123"},
            json={"prompt": "Pitch deck for Acme Corp", "deck_type": "sales_9"},
        )

        assert resp.status_code == 200
        assert audit.get_events()[-1].user_id == "banker-123"
    finally:
        audit.clear_events()


@pytest.mark.asyncio
async def test_generate_uses_provider_seam_at_request_time(
    client: AsyncClient,
    monkeypatch: pytest.MonkeyPatch,
):
    class FakeGenerator:
        def __init__(self):
            self.used = False

        async def generate(self, req, chart_data=None, upload_summary=None):
            self.used = True
            return [
                SlideData(
                    index=1,
                    title="Provider Seam",
                    bullets=["Safe generated content"],
                    notes="Speaker notes",
                    layout="title_content",
                    chart_data=chart_data,
                )
            ]

    fake = FakeGenerator()
    _patch_generator(monkeypatch, fake)

    resp = await client.post(
        "/api/v1/generate",
        json={"prompt": "Use fake generator", "deck_type": "sales_9"},
    )

    assert resp.status_code == 200
    assert fake.used
    assert resp.json()["slides"][0]["title"] == "Provider Seam"
    assert resp.json()["slides"][-1]["title"] == "Thank You"


@pytest.mark.asyncio
async def test_generate_keeps_normalized_deck_within_slide_count_window(
    client: AsyncClient,
    monkeypatch: pytest.MonkeyPatch,
):
    class MaxWindowGenerator:
        async def generate(self, req, chart_data=None, upload_summary=None):
            del req, chart_data, upload_summary
            return [
                SlideData(
                    index=i,
                    title=f"Slide {i}",
                    bullets=["Safe generated content"],
                    notes="Speaker notes",
                    layout="title" if i == 1 else "content",
                )
                for i in range(1, 13)
            ]

    _patch_generator(monkeypatch, MaxWindowGenerator())

    resp = await client.post(
        "/api/v1/generate",
        json={"prompt": "Use max allowed count", "deck_type": "sales_9"},
    )

    assert resp.status_code == 200
    slides = resp.json()["slides"]
    assert len(slides) == 12
    assert slides[-1]["title"] == "Thank You"
    assert slides[-1]["variant"] == "closing"


@pytest.mark.asyncio
async def test_generate_resolves_images_for_framework_visual_variants(
    client: AsyncClient,
    monkeypatch: pytest.MonkeyPatch,
):
    class FakeGenerator:
        async def generate(self, req, chart_data=None, upload_summary=None):
            del req, chart_data, upload_summary
            return [
                SlideData(index=1, title="Cover", bullets=[], notes="", layout="title", variant="cover"),
                SlideData(
                    index=2,
                    title="Problem",
                    bullets=["Legacy process"],
                    notes="",
                    layout="content",
                    variant="split_image",
                ),
                SlideData(
                    index=3,
                    title="Vision",
                    bullets=["Build the future"],
                    notes="",
                    layout="content",
                    variant="quote",
                    blocks=[{"type": "quote", "text": "Build the future"}],
                ),
                SlideData(index=4, title="Plain", bullets=["No image"], notes="", layout="content"),
                SlideData(index=5, title="Closing", bullets=["Questions"], notes="", layout="next_steps", variant="closing"),
            ]

    resolved: list[int] = []
    original_resolver = generate_router.image_resolver

    class Resolver:
        def needs_image(self, slide: SlideData) -> bool:
            return original_resolver.needs_image(slide)

        async def resolve(self, slide: SlideData) -> str:
            resolved.append(slide.index)
            return "IMG64"

    _patch_generator(monkeypatch, FakeGenerator())
    monkeypatch.setattr(generate_router, "image_resolver", Resolver())

    resp = await client.post(
        "/api/v1/generate",
        json={"prompt": "Use framework variants", "deck_type": "sales_9"},
    )

    assert resp.status_code == 200
    assert resolved == [1, 2, 3, 5, 6]
    slides = resp.json()["slides"]
    assert slides[1]["image_b64"] == "IMG64"
    assert slides[2]["image_b64"] == "IMG64"
    assert slides[3]["image_b64"] is None


@pytest.mark.asyncio
async def test_generate_with_uploaded_csv_adds_chart_data(client: AsyncClient):
    upload = await client.post(
        "/api/v1/uploads",
        files={"file": ("revenue.csv", b"quarter,revenue\nQ1,100\nQ2,125\n", "text/csv")},
    )
    file_id = upload.json()["file_id"]

    try:
        resp = await client.post(
            "/api/v1/generate",
            json={"prompt": "Internal revenue analysis", "deck_type": "internal_6", "file_id": file_id},
        )

        assert resp.status_code == 200
        chart_slides = [slide for slide in resp.json()["slides"] if slide["chart_data"]]
        assert chart_slides
        assert chart_slides[0]["chart_data"]["categories"] == ["Q1", "Q2"]
    finally:
        delete_upload(file_id)


@pytest.mark.asyncio
async def test_generate_accepts_ai_chart_recommendation_from_uploaded_columns(
    client: AsyncClient,
    monkeypatch: pytest.MonkeyPatch,
):
    class RecommendedChartService:
        async def generate(self, req, chart_data=None, upload_summary=None):
            assert chart_data is None
            assert upload_summary["filename"] == "revenue.csv"
            return [
                SlideData(index=1, title="Title", bullets=[], notes="", layout="title"),
                SlideData(
                    index=2,
                    title="Revenue Trend",
                    bullets=["Revenue increased across uploaded quarters"],
                    notes="Discuss uploaded data only.",
                    layout="chart",
                    visual_direction="Line chart on the right with source note.",
                    chart_recommendation=ChartRecommendation(
                        chart_type="line",
                        category_column="Quarter",
                        value_columns=["Revenue"],
                        rationale="Trend over time.",
                    ),
                ),
            ]

    _patch_generator(monkeypatch, RecommendedChartService())
    upload = await client.post(
        "/api/v1/uploads",
        files={"file": ("revenue.csv", b"Quarter,Revenue\nQ1,100\nQ2,125\n", "text/csv")},
    )
    file_id = upload.json()["file_id"]

    try:
        resp = await client.post(
            "/api/v1/generate",
            json={"prompt": "Build a revenue trend slide", "deck_type": "sales_9", "file_id": file_id},
        )

        assert resp.status_code == 200
        chart_slide = resp.json()["slides"][1]
        assert chart_slide["chart_data"]["type"] == "line"
        assert chart_slide["chart_data"]["categories"] == ["Q1", "Q2"]
        assert chart_slide["chart_data"]["series"] == [{"name": "Revenue", "values": [100.0, 125.0]}]
        assert chart_slide["chart_audit"]["source_filename"] == "revenue.csv"
        assert chart_slide["chart_audit"]["recommendation_status"] == "accepted"
    finally:
        delete_upload(file_id)


@pytest.mark.asyncio
async def test_generate_rejects_ai_chart_recommendation_for_missing_uploaded_columns(
    client: AsyncClient,
    monkeypatch: pytest.MonkeyPatch,
):
    class InvalidChartService:
        async def generate(self, req, chart_data=None, upload_summary=None):
            return [
                SlideData(index=1, title="Title", bullets=[], notes="", layout="title"),
                SlideData(
                    index=2,
                    title="Bookings Trend",
                    bullets=["Use uploaded facts only"],
                    notes="Missing chart columns should not create a fake chart.",
                    layout="chart",
                    visual_direction="Use a non-chart visual treatment if data is missing.",
                    chart_recommendation=ChartRecommendation(
                        chart_type="bar",
                        category_column="Month",
                        value_columns=["Bookings"],
                        rationale="Bookings by month.",
                    ),
                ),
            ]

    _patch_generator(monkeypatch, InvalidChartService())
    upload = await client.post(
        "/api/v1/uploads",
        files={"file": ("revenue.csv", b"Quarter,Revenue\nQ1,100\n", "text/csv")},
    )
    file_id = upload.json()["file_id"]

    try:
        resp = await client.post(
            "/api/v1/generate",
            json={"prompt": "Build a bookings trend slide", "deck_type": "sales_9", "file_id": file_id},
        )

        assert resp.status_code == 200
        chart_slide = resp.json()["slides"][1]
        assert chart_slide["chart_data"] is None
        assert chart_slide["chart_audit"]["recommendation_status"] == "rejected"
        assert "missing" in chart_slide["chart_audit"]["rejection_reason"].lower()
    finally:
        delete_upload(file_id)


@pytest.mark.asyncio
async def test_generate_with_prohibited_uploaded_csv_chart_data_returns_400(client: AsyncClient):
    upload = await client.post(
        "/api/v1/uploads",
        files={"file": ("revenue.csv", b"quarter,revenue\nrisk-free,100\n", "text/csv")},
    )
    file_id = upload.json()["file_id"]

    try:
        resp = await client.post(
            "/api/v1/generate",
            json={"prompt": "Internal revenue analysis", "deck_type": "internal_6", "file_id": file_id},
        )

        assert resp.status_code == 400
        assert "risk-free" in resp.json()["error"]["message"]
    finally:
        delete_upload(file_id)


@pytest.mark.asyncio
async def test_generate_with_corrupt_stored_xlsx_returns_400(client: AsyncClient):
    file_id = "corrupt.xlsx"
    path = Path(settings.local_upload_dir) / file_id
    path.parent.mkdir(parents=True, exist_ok=True)
    path.write_bytes(b"not a workbook")

    try:
        resp = await client.post(
            "/api/v1/generate",
            json={"prompt": "Internal revenue analysis", "deck_type": "internal_6", "file_id": file_id},
        )

        assert resp.status_code == 400
        assert resp.json()["detail"] == "Invalid uploaded file"
    finally:
        delete_upload(file_id)


@pytest.mark.asyncio
async def test_upload_csv(client: AsyncClient):
    resp = await client.post(
        "/api/v1/uploads",
        files={"file": ("revenue.csv", b"quarter,revenue\nQ1,100\n", "text/csv")},
    )

    assert resp.status_code == 200
    data = resp.json()
    try:
        assert data["file_id"].endswith(".csv")
        assert data["row_count"] == 1
        assert data["columns"] == ["quarter", "revenue"]
    finally:
        delete_upload(data["file_id"])


@pytest.mark.asyncio
async def test_upload_malformed_xlsx_returns_400(client: AsyncClient):
    resp = await client.post(
        "/api/v1/uploads",
        files={
            "file": (
                "bad.xlsx",
                b"not a workbook",
                "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
            )
        },
    )

    assert resp.status_code == 400


@pytest.mark.asyncio
async def test_upload_oversized_file_reads_only_limit_plus_one(monkeypatch: pytest.MonkeyPatch):
    class OversizedFile:
        filename = "large.csv"

        def __init__(self):
            self.read_size: int | None = None

        async def read(self, size: int | None = None):
            self.read_size = size
            return b"abcdef"

    monkeypatch.setattr(settings, "max_upload_bytes", 5)
    file = OversizedFile()

    with pytest.raises(HTTPException) as exc_info:
        await upload_file(file)  # type: ignore[arg-type]

    assert exc_info.value.status_code == 400
    assert file.read_size == settings.max_upload_bytes + 1


@pytest.mark.asyncio
async def test_upload_unexpected_service_error_is_not_400(
    client: AsyncClient,
    monkeypatch: pytest.MonkeyPatch,
):
    def raise_disk_full(*, filename: str, content: bytes):
        raise RuntimeError("disk full")

    monkeypatch.setattr(uploads.uploads, "save_upload", raise_disk_full)

    resp = await client.post(
        "/api/v1/uploads",
        files={"file": ("revenue.csv", b"quarter,revenue\nQ1,100\n", "text/csv")},
    )

    assert resp.status_code == 500
    assert resp.json()["error"]["code"] == "INTERNAL_ERROR"


@pytest.mark.asyncio
async def test_generate_dlp_block(client: AsyncClient):
    resp = await client.post(
        "/api/v1/generate",
        json={"prompt": "We guarantee returns of 20%", "deck_type": "sales_9"},
    )
    assert resp.status_code == 400
    assert "guarantee returns" in resp.json()["error"]["message"]


@pytest.mark.asyncio
async def test_refine_dlp_block(client: AsyncClient):
    gen = await client.post(
        "/api/v1/generate",
        json={"prompt": "Pitch for Acme Corp", "deck_type": "sales_9"},
    )
    session_id = gen.json()["session_id"]

    resp = await client.post(
        "/api/v1/refine",
        json={"session_id": session_id, "slide_index": 3, "instruction": "make this risk-free"},
    )

    assert resp.status_code == 400
    assert "risk-free" in resp.json()["error"]["message"]


@pytest.mark.asyncio
async def test_refine_slide(client: AsyncClient):
    gen = await client.post(
        "/api/v1/generate",
        json={"prompt": "Pitch for Acme Corp", "deck_type": "sales_9"},
    )
    session_id = gen.json()["session_id"]

    resp = await client.post(
        "/api/v1/refine",
        json={"session_id": session_id, "slide_index": 3, "instruction": "make shorter"},
    )
    assert resp.status_code == 200
    assert "Refined" in resp.json()["slide"]["title"]


@pytest.mark.asyncio
async def test_refine_uses_shared_image_resolver(client: AsyncClient, monkeypatch: pytest.MonkeyPatch):
    gen = await client.post(
        "/api/v1/generate",
        json={"prompt": "Pitch for Acme Corp", "deck_type": "sales_9"},
    )
    session_id = gen.json()["session_id"]
    resolved: list[str] = []

    class Resolver:
        def needs_image(self, slide: SlideData) -> bool:
            return slide.index == 3

        async def resolve(self, slide: SlideData) -> str:
            resolved.append(slide.title)
            return "REFINED64"

    monkeypatch.setattr(refine_router, "image_resolver", Resolver())

    resp = await client.post(
        "/api/v1/refine",
        json={"session_id": session_id, "slide_index": 3, "instruction": "make shorter"},
    )

    assert resp.status_code == 200
    assert resolved
    assert resp.json()["slide"]["image_b64"] == "REFINED64"


@pytest.mark.asyncio
async def test_refine_nonexistent_session(client: AsyncClient):
    resp = await client.post(
        "/api/v1/refine",
        json={"session_id": "does-not-exist", "slide_index": 1, "instruction": "shorter"},
    )
    assert resp.status_code == 404
    assert resp.json()["error"]["code"] == "SESSION_NOT_FOUND"


@pytest.mark.asyncio
async def test_export(client: AsyncClient):
    gen = await client.post(
        "/api/v1/generate",
        json={"prompt": "Export test", "deck_type": "sales_9"},
    )
    session_id = gen.json()["session_id"]

    resp = await client.post(
        "/api/v1/export",
        json={"session_id": session_id},
    )
    assert resp.status_code == 200
    data = resp.json()
    assert "download_url" in data
    assert "expires_at" in data
    # The download URL must be absolute (point at the backend host), not a
    # host-less relative path that the browser would resolve against the frontend.
    assert data["download_url"].startswith("http://test/api/v1/download/")
    download = await client.get(data["download_url"].replace("http://test", ""))
    assert download.status_code == 200
    assert download.content.startswith(b"PK")
    assert download.headers["content-type"].startswith(
        "application/vnd.openxmlformats-officedocument.presentationml.presentation"
    )
    Presentation(BytesIO(download.content))
    # Served as an attachment with a friendly filename, not the raw session UUID.
    disposition = download.headers.get("content-disposition", "")
    assert "attachment" in disposition
    assert "SlideForge-Presentation.pptx" in disposition


@pytest.mark.asyncio
async def test_export_rejects_invalid_pptx_before_upload(
    client: AsyncClient,
    monkeypatch: pytest.MonkeyPatch,
):
    gen = await client.post(
        "/api/v1/generate",
        json={"prompt": "Export validation test", "deck_type": "sales_9"},
    )
    session_id = gen.json()["session_id"]

    class BadEngine:
        def __init__(self, **kwargs):
            del kwargs

        def render(self, slides):
            del slides
            return b"not a pptx"

    uploaded: list[bytes] = []

    class RecordingStorage:
        async def upload_pptx(self, session_id, content, base_url=None):
            del session_id, base_url
            uploaded.append(content)
            return "http://test/api/v1/download/bad.pptx"

    monkeypatch.setattr(export_router, "PptxEngine", BadEngine)
    _patch_storage(monkeypatch, RecordingStorage())

    resp = await client.post(
        "/api/v1/export",
        json={"session_id": session_id},
    )

    assert resp.status_code == 502
    assert resp.json()["error"]["code"] == "GENERATION_ERROR"
    assert uploaded == []


@pytest.mark.asyncio
async def test_export_audit_slide_count_matches_rendered_pptx(
    client: AsyncClient,
    monkeypatch: pytest.MonkeyPatch,
):
    class ShortGenerator:
        async def generate(self, req, chart_data=None, upload_summary=None):
            del req, chart_data, upload_summary
            return [
                SlideData(index=1, title="Cover", bullets=[], notes="", layout="title"),
                SlideData(index=2, title="Action Plan", bullets=["Review"], notes="", layout="next_steps"),
            ]

    _patch_generator(monkeypatch, ShortGenerator())
    audit = providers.get_audit_service()
    audit.clear_events()

    gen = await client.post(
        "/api/v1/generate",
        json={"prompt": "Short deck", "deck_type": "sales_9"},
    )
    session_id = gen.json()["session_id"]

    export = await client.post("/api/v1/export", json={"session_id": session_id})
    download = await client.get(export.json()["download_url"].replace("http://test", ""))
    prs = Presentation(BytesIO(download.content))

    assert len(prs.slides) == 3
    assert audit.get_events()[-1].slide_count == len(prs.slides)


@pytest.mark.asyncio
async def test_download_expired_export_returns_404_and_removes_file(
    client: AsyncClient,
    tmp_path: Path,
    monkeypatch: pytest.MonkeyPatch,
):
    monkeypatch.setattr(settings, "local_export_dir", str(tmp_path))
    monkeypatch.setattr(settings, "signed_url_expiry_minutes", 30)
    filename = "expired.pptx"
    path = tmp_path / filename
    path.write_bytes(b"old pptx")
    old = time.time() - (31 * 60)
    os.utime(path, (old, old))

    resp = await client.get(f"/api/v1/download/{filename}")

    assert resp.status_code == 404
    assert not path.exists()


def test_purge_local_temp_files_invokes_storage_and_upload_purges(monkeypatch: pytest.MonkeyPatch):
    calls: list[tuple[str, int]] = []

    class FakeStorage:
        def purge_expired(self, max_age_seconds: int) -> int:
            calls.append(("storage", max_age_seconds))
            return 1

    class FakeUploads:
        def purge_expired(self, max_age_seconds: int) -> int:
            calls.append(("uploads", max_age_seconds))
            return 2

    monkeypatch.setattr("app.main.StorageService", FakeStorage)
    monkeypatch.setattr("app.main.UploadService", FakeUploads)
    monkeypatch.setattr(settings, "signed_url_expiry_minutes", 30)
    monkeypatch.setattr(settings, "session_ttl_minutes", 45)

    assert purge_local_temp_files() == {"exports": 1, "uploads": 2}
    assert calls == [("storage", 1800), ("uploads", 2700)]


@pytest.mark.asyncio
async def test_generate_script_mode_chunks_source(client: AsyncClient):
    script = (
        "Market Overview\nThe regional market is expanding steadily.\n\n"
        "Our Approach\nWe propose a phased rollout to manage risk.\n\n"
        "Next Steps\nReview the plan and confirm the timeline."
    )
    resp = await client.post(
        "/api/v1/generate",
        json={"prompt": script, "deck_type": "sales_9", "source_type": "script"},
    )
    assert resp.status_code == 200
    slides = resp.json()["slides"]
    # Content-derived count plus a normalized final thank-you slide.
    assert len(slides) == 4
    assert slides[0]["layout"] == "title"
    # Original source text preserved in notes.
    assert "regional market is expanding" in slides[0]["notes"]
    # Bullet cap enforced everywhere.
    assert all(len(s["bullets"]) <= 5 for s in slides)
    assert slides[-1]["title"] == "Thank You"
