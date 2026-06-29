import asyncio
import hashlib
from io import BytesIO

import httpx
import jwt
import pytest
from httpx import ASGITransport, AsyncClient
from pptx import Presentation
from structlog.testing import capture_logs

from app import dependencies
from app.main import app
from app.services.platform.database import Database
from app.services.platform.deck_files import LocalDeckFileStorage
from app.services.platform.deck_repository import DeckRepository
from app.services.platform.deck_versions import DeckVersionService
from app.services.platform.onlyoffice import OnlyOfficeService


SECRET = "test-secret-with-at-least-thirty-two-bytes"


class ChunkStream(httpx.AsyncByteStream):
    def __init__(self, *chunks: bytes) -> None:
        self._chunks = chunks

    async def __aiter__(self):
        for chunk in self._chunks:
            yield chunk


def _pptx_bytes() -> bytes:
    presentation = Presentation()
    presentation.slides.add_slide(presentation.slide_layouts[0])
    output = BytesIO()
    presentation.save(output)
    return output.getvalue()


class FakeVersionService:
    def __init__(self) -> None:
        self.calls: list[dict[str, object]] = []
        self.fail = False
        self.version_ids: set[str] = {"version-1"}

    async def save_edited_version(self, **kwargs):
        self.calls.append(kwargs)
        if self.fail:
            raise OSError("storage failed for http://onlyoffice/private?token=secret")
        self.version_ids.add(str(kwargs["callback_key"]))
        return object()


def _authorization(body: dict[str, object]) -> str:
    token = jwt.encode({"payload": body}, SECRET, algorithm="HS256")
    return f"Bearer {token}"


@pytest.fixture
async def callback_client():
    pptx = _pptx_bytes()

    async def handler(request: httpx.Request) -> httpx.Response:
        return httpx.Response(200, content=pptx, request=request)

    download_client = httpx.AsyncClient(transport=httpx.MockTransport(handler))
    service = OnlyOfficeService(
        public_url="http://localhost:8080",
        api_base_url="http://api:8000",
        internal_url="http://onlyoffice:80",
        jwt_secret=SECRET,
        file_token_ttl_seconds=300,
        max_file_bytes=len(pptx) + 100,
        enabled=True,
        download_client=download_client,
    )
    versions = FakeVersionService()
    app.dependency_overrides[dependencies.get_onlyoffice_service] = lambda: service
    app.dependency_overrides[dependencies.get_deck_version_service] = lambda: versions
    transport = ASGITransport(app=app, raise_app_exceptions=False)
    async with AsyncClient(transport=transport, base_url="http://test") as client:
        yield client, service, versions, download_client, pptx
    app.dependency_overrides.clear()
    await download_client.aclose()


async def _post(client, service, body, *, token=None, authorization=None):
    callback_token = token or service.create_scoped_token(
        subject="alice",
        deck_id="deck-1",
        version_id="version-1",
        purpose="callback",
    )
    headers = {"Authorization": authorization or _authorization(body)}
    return await client.post(
        f"/api/v1/decks/deck-1/callback?token={callback_token}",
        json=body,
        headers=headers,
    )


@pytest.mark.asyncio
@pytest.mark.parametrize("status", [1, 4])
async def test_non_save_status_is_successful_noop(callback_client, status):
    client, service, versions, *_ = callback_client
    body = {"key": "deck-1-version-1", "status": status, "users": ["attacker"]}

    response = await _post(client, service, body)

    assert response.status_code == 200
    assert response.json() == {"error": 0}
    assert versions.calls == []


@pytest.mark.asyncio
@pytest.mark.parametrize("status", [3, 7])
async def test_error_status_returns_callback_error(callback_client, status):
    client, service, versions, *_ = callback_client
    body = {"key": "deck-1-version-1", "status": status}

    response = await _post(client, service, body)

    assert response.status_code == 200
    assert response.json() == {"error": 1}
    assert versions.calls == []


@pytest.mark.asyncio
@pytest.mark.parametrize("status", [2, 6])
async def test_save_status_uses_only_trusted_token_identity(callback_client, status):
    client, service, versions, *_ = callback_client
    body = {
        "key": "deck-1-version-1",
        "status": status,
        "url": "http://onlyoffice/coauthoring/download/presentation.pptx",
        "users": ["mallory"],
        "userdata": "save-42",
    }

    response = await _post(client, service, body)

    assert response.status_code == 200
    assert response.json() == {"error": 0}
    assert versions.calls == [
        {
            "deck_id": "deck-1",
            "owner_id": "alice",
            "content": callback_client[4],
            "base_version_id": "version-1",
            "callback_key": f"deck-1-version-1:{status}:save-42",
            "created_by": "alice",
        }
    ]


@pytest.mark.asyncio
async def test_identical_callback_is_idempotent_at_version_service_boundary(callback_client):
    client, service, versions, *_ = callback_client
    body = {
        "key": "deck-1-version-1",
        "status": 6,
        "url": "http://onlyoffice/download.pptx",
    }

    first = await _post(client, service, body)
    second = await _post(client, service, body)

    assert first.json() == second.json() == {"error": 0}
    assert {str(item["callback_key"]) for item in versions.calls} == {
        "deck-1-version-1:6:"
    }
    assert len(versions.version_ids) == 2


@pytest.mark.asyncio
async def test_invalid_callback_token_is_rejected(callback_client):
    client, service, versions, *_ = callback_client
    body = {"key": "deck-1-version-1", "status": 1}

    response = await _post(client, service, body, token="invalid")

    assert response.status_code == 401
    assert versions.calls == []


@pytest.mark.asyncio
async def test_authorization_jwt_must_match_callback_body(callback_client):
    client, service, versions, *_ = callback_client
    body = {"key": "deck-1-version-1", "status": 1}
    mismatch = _authorization({"key": "other-key", "status": 1})

    response = await _post(client, service, body, authorization=mismatch)

    assert response.status_code == 401
    assert versions.calls == []


@pytest.mark.asyncio
@pytest.mark.parametrize(
    "url",
    [
        "http://evil.example/download.pptx",
        "https://onlyoffice/download.pptx",
        "http://onlyoffice:81/download.pptx",
        "http://onlyoffice.evil.example/download.pptx",
        "http://onlyoffice@evil.example/download.pptx",
    ],
)
async def test_download_url_must_match_configured_origin(callback_client, url):
    client, service, versions, *_ = callback_client
    body = {"key": "deck-1-version-1", "status": 6, "url": url}

    response = await _post(client, service, body)

    assert response.json() == {"error": 1}
    assert versions.calls == []


@pytest.mark.asyncio
@pytest.mark.parametrize(
    "failure", ["redirect", "timeout", "oversize", "stream_oversize", "invalid"]
)
async def test_unsafe_or_invalid_download_is_not_persisted(callback_client, failure):
    client, service, versions, download_client, pptx = callback_client

    async def handler(request: httpx.Request) -> httpx.Response:
        if failure == "redirect":
            return httpx.Response(302, headers={"Location": "http://evil.example/file"}, request=request)
        if failure == "timeout":
            raise httpx.ReadTimeout("timed out", request=request)
        if failure == "oversize":
            return httpx.Response(200, content=pptx + b"too-large" * 100, request=request)
        if failure == "stream_oversize":
            return httpx.Response(
                200,
                stream=ChunkStream(pptx, b"too-large" * 100),
                request=request,
            )
        return httpx.Response(200, content=b"not a pptx", request=request)

    download_client._transport = httpx.MockTransport(handler)
    body = {
        "key": "deck-1-version-1",
        "status": 6,
        "url": "http://onlyoffice/file.pptx",
    }

    response = await _post(client, service, body)

    assert response.json() == {"error": 1}
    assert versions.calls == []


@pytest.mark.asyncio
async def test_persistence_failure_returns_error_and_logs_no_secrets(callback_client):
    client, service, versions, *_ = callback_client
    versions.fail = True
    body = {
        "key": "deck-1-version-1",
        "status": 6,
        "url": "http://onlyoffice/private?token=download-secret",
    }
    callback_token = service.create_scoped_token(
        subject="alice", deck_id="deck-1", version_id="version-1", purpose="callback"
    )

    with capture_logs() as logs:
        response = await _post(client, service, body, token=callback_token)

    assert response.json() == {"error": 1}
    assert versions.version_ids == {"version-1"}
    rendered_logs = repr(logs)
    assert "onlyoffice_callback_failed" in rendered_logs
    assert "download-secret" not in rendered_logs
    assert callback_token not in rendered_logs
    assert "http://" not in rendered_logs


@pytest.mark.asyncio
async def test_mismatched_document_key_is_rejected_before_download(callback_client):
    client, service, versions, download_client, _pptx = callback_client

    async def unexpected_download(request: httpx.Request) -> httpx.Response:
        raise AssertionError(f"download must not occur: {request.url}")

    download_client._transport = httpx.MockTransport(unexpected_download)
    body = {
        "key": "deck-2-version-1",
        "status": 6,
        "url": "http://onlyoffice/private.pptx",
    }

    response = await _post(client, service, body)

    assert response.status_code == 200
    assert response.json() == {"error": 1}
    assert versions.calls == []


@pytest.mark.asyncio
async def test_callback_requires_authorization_when_enabled(callback_client):
    client, service, versions, *_ = callback_client
    callback_token = service.create_scoped_token(
        subject="alice",
        deck_id="deck-1",
        version_id="version-1",
        purpose="callback",
    )

    response = await client.post(
        f"/api/v1/decks/deck-1/callback?token={callback_token}",
        json={"key": "deck-1-version-1", "status": 1},
    )

    assert response.status_code == 401
    assert versions.calls == []


@pytest.mark.asyncio
async def test_drip_feed_download_hits_total_deadline(callback_client):
    client, _service, versions, download_client, pptx = callback_client

    class DripStream(httpx.AsyncByteStream):
        async def __aiter__(self):
            for byte in pptx:
                await asyncio.sleep(0.001)
                yield bytes([byte])

    async def handler(request: httpx.Request) -> httpx.Response:
        return httpx.Response(200, stream=DripStream(), request=request)

    download_client._transport = httpx.MockTransport(handler)
    deadline_service = OnlyOfficeService(
        public_url="http://localhost:8080",
        api_base_url="http://api:8000",
        internal_url="http://onlyoffice:80",
        jwt_secret=SECRET,
        file_token_ttl_seconds=300,
        max_file_bytes=len(pptx) + 100,
        enabled=True,
        download_client=download_client,
        download_deadline_seconds=0.01,
    )
    app.dependency_overrides[dependencies.get_onlyoffice_service] = (
        lambda: deadline_service
    )
    body = {
        "key": "deck-1-version-1",
        "status": 6,
        "url": "http://onlyoffice/drip.pptx",
    }

    response = await _post(client, deadline_service, body)

    assert response.status_code == 200
    assert response.json() == {"error": 1}
    assert versions.calls == []


@pytest.mark.asyncio
async def test_disabled_onlyoffice_rejects_all_routes_without_secret():
    class UnusedDependency:
        def __getattr__(self, name):
            raise AssertionError(f"disabled route used dependency method {name}")

    service = OnlyOfficeService(
        public_url="http://localhost:8080",
        api_base_url="http://api:8000",
        internal_url="http://onlyoffice",
        jwt_secret="",
        file_token_ttl_seconds=300,
        enabled=False,
    )
    unused = UnusedDependency()
    app.dependency_overrides[dependencies.get_onlyoffice_service] = lambda: service
    app.dependency_overrides[dependencies.get_deck_repository] = lambda: unused
    app.dependency_overrides[dependencies.get_deck_file_storage] = lambda: unused
    app.dependency_overrides[dependencies.get_deck_version_service] = lambda: unused
    transport = ASGITransport(app=app, raise_app_exceptions=False)
    async with AsyncClient(transport=transport, base_url="http://test") as client:
        responses = [
            await client.get(
                "/api/v1/decks/deck-1/editor-config",
                headers={"x-user-id": "alice"},
            ),
            await client.get("/api/v1/decks/deck-1/content?token=not-a-secret"),
            await client.post(
                "/api/v1/decks/deck-1/callback?token=not-a-secret",
                json={"key": "deck-1-version-1", "status": 1},
            ),
        ]
    app.dependency_overrides.clear()

    assert [response.status_code for response in responses] == [404, 404, 404]
    assert all("secret" not in response.text.lower() for response in responses)


@pytest.mark.asyncio
async def test_concurrent_identical_callbacks_create_one_real_version(tmp_path):
    pptx = _pptx_bytes()
    database = Database(f"sqlite+aiosqlite:///{tmp_path / 'callback.db'}")
    await database.create_schema()
    repository = DeckRepository(database)
    storage = LocalDeckFileStorage(tmp_path / "objects")
    await repository.create_with_initial_version(
        deck_id="deck-1",
        version_id="version-1",
        owner_id="alice",
        name="Deck",
        deck_type="sales_9",
        theme="minimalist",
        aspect_ratio="16:9",
        generation_payload=None,
        storage_key="decks/deck-1/versions/version-1.pptx",
        sha256=hashlib.sha256(pptx).hexdigest(),
        size_bytes=len(pptx),
    )
    versions = DeckVersionService(
        repository=repository,
        storage=storage,
        sample_template_path=None,
        max_file_bytes=len(pptx) + 100,
        retention=5,
    )

    async def handler(request: httpx.Request) -> httpx.Response:
        return httpx.Response(200, content=pptx, request=request)

    download_client = httpx.AsyncClient(transport=httpx.MockTransport(handler))
    service = OnlyOfficeService(
        public_url="http://localhost:8080",
        api_base_url="http://api:8000",
        internal_url="http://onlyoffice",
        jwt_secret=SECRET,
        file_token_ttl_seconds=300,
        max_file_bytes=len(pptx) + 100,
        enabled=True,
        download_client=download_client,
    )
    app.dependency_overrides[dependencies.get_onlyoffice_service] = lambda: service
    app.dependency_overrides[dependencies.get_deck_version_service] = lambda: versions
    transport = ASGITransport(app=app, raise_app_exceptions=False)
    body = {
        "key": "deck-1-version-1",
        "status": 6,
        "url": "http://onlyoffice/file.pptx",
        "userdata": "concurrent-save",
    }
    token = service.create_scoped_token(
        subject="alice",
        deck_id="deck-1",
        version_id="version-1",
        purpose="callback",
    )
    headers = {"Authorization": _authorization(body)}
    url = f"/api/v1/decks/deck-1/callback?token={token}"
    try:
        async with AsyncClient(transport=transport, base_url="http://test") as client:
            first, second = await asyncio.gather(
                client.post(url, json=body, headers=headers),
                client.post(url, json=body, headers=headers),
            )
            mismatched = {**body, "key": "deck-2-version-1"}
            mismatch_response = await client.post(
                url,
                json=mismatched,
                headers={"Authorization": _authorization(mismatched)},
            )

        persisted = await repository.get("deck-1", "alice")
        rows = await repository.list_versions("deck-1", "alice")
    finally:
        app.dependency_overrides.clear()
        await download_client.aclose()
        await database.dispose()

    assert first.json() == second.json() == {"error": 0}
    assert mismatch_response.json() == {"error": 1}
    assert len(rows) == 2
    assert persisted is not None
    assert persisted.current_version is not None
    assert persisted.current_version.version_number == 2
    assert persisted.current_version_id == rows[0].id
