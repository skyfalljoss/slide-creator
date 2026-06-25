import logging

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.enum.text import MSO_AUTO_SIZE, PP_ALIGN
from pptx.util import Inches, Pt

from app.services.presentation.pptx_canvas import PptxCanvas
from app.services.presentation.pptx_layout import CANVAS_DIMS, LAYOUT
from app.services.presentation.pptx_theme import DISPLAY_FONT, THEMES

_PNG_1PX = "iVBORw0KGgoAAAANSUhEUgAAAAEAAAABCAQAAAC1HAwCAAAAC0lEQVR42mP8/x8AAwMCAO+/p9sAAAAASUVORK5CYII="


def _canvas(theme_name: str = "minimalist") -> PptxCanvas:
    theme = THEMES[theme_name]
    canvas_width, canvas_height = CANVAS_DIMS["16:9"]
    return PptxCanvas(
        canvas_width=canvas_width,
        canvas_height=canvas_height,
        logical_width=LAYOUT.logical_width,
        theme=theme,
        active_theme=lambda: theme,
        map_color=lambda color: color,
    )


def _blank_slide():
    prs = Presentation()
    return prs.slides.add_slide(prs.slide_layouts[6])


def test_canvas_add_text_sets_font_and_margins():
    theme = THEMES["minimalist"]
    slide = _blank_slide()

    shape = _canvas().add_text(
        slide,
        LAYOUT.left_margin,
        1.0,
        4.0,
        0.5,
        "Revenue Outlook",
        20,
        theme.text,
        bold=True,
        align=PP_ALIGN.CENTER,
    )

    text_frame = shape.text_frame
    paragraph = text_frame.paragraphs[0]
    assert shape.left == Inches(LAYOUT.left_margin)
    assert text_frame.word_wrap is True
    assert text_frame.margin_left == 0
    assert text_frame.margin_right == 0
    assert text_frame.margin_top == 0
    assert text_frame.margin_bottom == 0
    assert paragraph.text == "Revenue Outlook"
    assert paragraph.font.name == DISPLAY_FONT
    assert paragraph.font.size == Pt(20)
    assert paragraph.font.bold is True
    assert paragraph.font.color.rgb == theme.text
    assert paragraph.alignment == PP_ALIGN.CENTER


def test_canvas_bullets_box_uses_shape_to_fit_text():
    slide = _blank_slide()

    shape = _canvas().add_bullets_box(slide, ["**Key:** detail", "Second point"], 1.0, 1.0, 5.0, 2.0)

    assert shape.text_frame.auto_size == MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT
    assert shape.text_frame.word_wrap is True
    assert "Key:" in shape.text
    assert "\u2022" in shape.text
    assert "**" not in shape.text


def test_canvas_add_slide_image_logs_failure(caplog):
    slide = _blank_slide()

    with caplog.at_level(logging.WARNING, logger="app.services.presentation.pptx_canvas"):
        inserted = _canvas().add_slide_image(slide, "not-an-image", 1.0, 1.0, 2.0, 2.0)

    assert inserted is False
    assert any("Failed to insert slide image" in record.message for record in caplog.records)
    assert any(record.exc_info for record in caplog.records)
    assert not any(shape.shape_type == MSO_SHAPE_TYPE.PICTURE for shape in slide.shapes)


def test_canvas_add_slide_image_inserts_picture():
    slide = _blank_slide()

    inserted = _canvas().add_slide_image(slide, _PNG_1PX, 1.0, 1.0, 2.0, 2.0)

    assert inserted is True
    pictures = [shape for shape in slide.shapes if shape.shape_type == MSO_SHAPE_TYPE.PICTURE]
    assert len(pictures) == 1
    assert pictures[0].left == Inches(1.0)
    assert pictures[0].width == Inches(2.0)
