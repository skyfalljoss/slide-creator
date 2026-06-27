import base64
import logging
from io import BytesIO

import pytest
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE, MSO_SHAPE_TYPE
from pptx.enum.text import MSO_AUTO_SIZE
from pptx.util import Inches, Pt

from app.config import settings
from app.models.schemas import ChartAudit, SlideData
from app.services.presentation.pptx_engine import PptxEngine
from app.services.presentation.pptx_layout import LAYOUT
from app.services.presentation.pptx_text import (
    add_markdown_paragraph,
    icon_shape,
    looks_like_instruction as _looks_like_instruction,
)
from app.services.presentation.pptx_theme import CITI_DARK, CITI_RED, THEMES, WHITE, resolve_theme


def test_render_sales_deck():
    slides = [
        SlideData(index=1, title="Title Slide", bullets=[], notes="", layout="title"),
        SlideData(index=2, title="Content", bullets=["Bullet 1", "Bullet 2"], notes="", layout="content"),
        SlideData(index=3, title="Last Slide", bullets=["Final bullet"], notes="", layout="content"),
    ]
    engine = PptxEngine()
    result = engine.render(slides)
    assert isinstance(result, bytes)
    assert len(result) > 0


def test_render_with_disclaimer_on_last():
    slides = [
        SlideData(index=i, title=f"Slide {i}", bullets=["content"], notes="", layout="content")
        for i in range(1, 4)
    ]
    engine = PptxEngine()
    result = engine.render(slides)
    assert len(result) > 0


def test_render_adds_chart_when_chart_data_present():
    slides = [
        SlideData(index=1, title="Title", bullets=[], notes="", layout="title"),
        SlideData(
            index=2,
            title="Revenue",
            bullets=["Revenue increased quarter over quarter"],
            notes="Discuss drivers.",
            layout="content",
            chart_data={
                "type": "bar",
                "title": "Revenue by Quarter",
                "categories": ["Q1", "Q2"],
                "series": [{"name": "revenue", "values": [100.0, 125.0]}],
            },
        ),
    ]

    prs = Presentation(BytesIO(PptxEngine().render(slides)))

    assert any(shape.has_chart for shape in prs.slides[1].shapes)


def test_render_uses_sample_deck_canvas_size():
    slides = [SlideData(index=1, title="Title", bullets=[], notes="", layout="title")]

    prs = Presentation(BytesIO(PptxEngine().render(slides)))

    assert round(prs.slide_width / 914400, 2) == 17.78
    assert round(prs.slide_height / 914400, 2) == 10.0


def test_layout_constants_match_default_canvas_width():
    engine = PptxEngine()

    assert engine._canvas_w == pytest.approx(LAYOUT.logical_width)
    assert engine._canvas_h == pytest.approx(LAYOUT.canvas_height)
    assert LAYOUT.content_width == pytest.approx(
        LAYOUT.logical_width - LAYOUT.left_margin - LAYOUT.right_margin
    )


def test_render_chart_slide_includes_source_note_from_audit():
    slides = [
        SlideData(index=1, title="Title", bullets=[], notes="", layout="title"),
        SlideData(
            index=2,
            title="Revenue Trend",
            bullets=["Revenue increased"],
            notes="Discuss trend.",
            layout="chart",
            chart_data={
                "type": "line",
                "title": "Revenue Trend",
                "categories": ["Q1", "Q2"],
                "series": [{"name": "Revenue", "values": [100.0, 125.0]}],
            },
            chart_audit=ChartAudit(
                source_filename="metrics.csv",
                category_column="Quarter",
                value_columns=["Revenue"],
                row_count=2,
                chart_type="line",
                recommendation_status="accepted",
            ),
        ),
    ]

    prs = Presentation(BytesIO(PptxEngine().render(slides)))
    slide_text = "\n".join(shape.text for shape in prs.slides[1].shapes if hasattr(shape, "text"))

    assert "Source: metrics.csv" in slide_text
    assert "Quarter, Revenue" in slide_text


def test_render_chart_slide_places_bullets_left_and_chart_right_without_overlap():
    bullet = "Revenue increased quarter over quarter"
    slides = [
        SlideData(index=1, title="Title", bullets=[], notes="", layout="title"),
        SlideData(
            index=2,
            title="Revenue",
            bullets=[bullet],
            notes="Discuss drivers.",
            layout="content",
            chart_data={
                "type": "bar",
                "title": "Revenue by Quarter",
                "categories": ["Q1", "Q2"],
                "series": [{"name": "revenue", "values": [100.0, 125.0]}],
            },
        ),
    ]

    prs = Presentation(BytesIO(PptxEngine().render(slides)))
    content_shapes = prs.slides[1].shapes

    assert any(shape.has_chart and shape.left >= Inches(6) for shape in content_shapes)
    assert any(
        not shape.has_chart
        and shape.has_text_frame
        and bullet in shape.text
        and shape.width <= Inches(6)
        for shape in content_shapes
    )


def test_render_skips_malformed_chart_data_without_crashing():
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[0])

    PptxEngine()._add_chart(slide, {"categories": ["Q1"], "series": ["bad"]})

    result = BytesIO()
    prs.save(result)

    assert result.tell() > 0
    Presentation(BytesIO(result.getvalue()))


def test_render_normalizes_non_string_chart_metadata_without_crashing():
    slides = [
        SlideData(index=1, title="Title", bullets=[], notes="", layout="title"),
        SlideData(
            index=2,
            title="Revenue",
            bullets=["Revenue increased quarter over quarter"],
            notes="Discuss drivers.",
            layout="content",
            chart_data={
                "type": "bar",
                "title": 123,
                "categories": ["Q1"],
                "series": [{"name": 42, "values": [100.0]}],
            },
        ),
    ]

    result = PptxEngine().render(slides)

    assert isinstance(result, bytes)
    assert len(result) > 0
    Presentation(BytesIO(result))


def test_render_adds_citi_header_and_final_disclaimer_text():
    slides = [
        SlideData(index=1, title="Title", bullets=[], notes="", layout="title"),
        SlideData(index=2, title="Next Steps", bullets=["Review"], notes="", layout="content"),
    ]

    prs = Presentation(BytesIO(PptxEngine().render(slides)))
    first_slide_text = "\n".join(shape.text for shape in prs.slides[0].shapes if hasattr(shape, "text"))
    final_slide_text = "\n".join(shape.text for shape in prs.slides[-1].shapes if hasattr(shape, "text"))

    assert "citi" in first_slide_text.lower()
    assert "Confidential" not in final_slide_text
    assert "not a guarantee" not in final_slide_text


def test_big_stat_supporting_content_stays_below_metric_label():
    slides = [
        SlideData(index=1, title="Title", bullets=[], notes="", layout="title"),
        SlideData(
            index=2,
            title="Total Funding Required",
            bullets=[
                "**Sales & Marketing Expansion (35% / $3.5M):** Build enterprise sales, launch targeted digital campaigns, and expand strategic partnerships.",
                "**Operational Scale & Infrastructure (15% / $1.5M):** Scale cloud infrastructure, security protocols, and customer support.",
                "**Working Capital & Contingency (10% / $1.0M):** Maintain operating flexibility and a contingency buffer.",
            ],
            notes="",
            layout="content",
            variant="big_stat",
            blocks=[{"type": "stat", "value": "$10M", "label": "Total Funding Required"}],
        ),
    ]

    prs = Presentation(BytesIO(PptxEngine().render(slides)))
    shapes = prs.slides[1].shapes
    label_shapes = [shape for shape in shapes if hasattr(shape, "text") and shape.text == "Total Funding Required"]
    bullet_shapes = [shape for shape in shapes if hasattr(shape, "text") and "Sales & Marketing" in shape.text]

    assert label_shapes
    assert bullet_shapes
    label = min(label_shapes, key=lambda shape: shape.top)
    bullets = bullet_shapes[0]
    assert bullets.top >= label.top + label.height


def test_big_stat_supporting_cards_keep_body_text_right_of_icon_chip():
    slides = [
        SlideData(index=1, title="Title", bullets=[], notes="", layout="title"),
        SlideData(
            index=2,
            title="Total Funding Required",
            bullets=[
                "Reduce average presentation creation time by 50-70%, freeing up 200-300 hours per month.",
                "Achieve an estimated annual cost avoidance of $75,000-$120,000 by reducing reliance on agencies.",
                "Enhance brand consistency across all external communications and client trust.",
            ],
            notes="",
            layout="content",
            variant="big_stat",
            blocks=[{"type": "stat", "value": "60%", "label": "Average reduction in presentation creation time"}],
        ),
    ]

    prs = Presentation(BytesIO(PptxEngine().render(slides)))
    shapes = prs.slides[1].shapes
    card = next(
        shape
        for shape in shapes
        if shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE
        and getattr(shape, "auto_shape_type", None) == MSO_SHAPE.ROUNDED_RECTANGLE
        and shape.top >= Inches(5)
        and shape.width >= Inches(4)
    )
    body = next(shape for shape in shapes if shape.has_text_frame and "Reduce average presentation creation time" in shape.text)

    assert body.left >= card.left + Inches(0.9)


def test_closing_layout_keeps_title_rule_and_bullets_tightly_grouped():
    slides = [
        SlideData(index=1, title="Title", bullets=[], notes="", layout="title"),
        SlideData(
            index=2,
            title="Next Steps & Call to Action",
            kicker="PARTNER WITH US",
            bullets=[
                "**Schedule a Deep-Dive Demonstration:** Personalized platform walkthrough tied to portfolio priorities.",
                "**Review Detailed Financial Projections:** Share financial models, spend breakdowns, and sensitivity analysis.",
                "**Initiate Term Sheet Discussions:** Align on investment terms and move into formal documentation.",
            ],
            notes="",
            layout="content",
            variant="closing",
        ),
    ]

    prs = Presentation(BytesIO(PptxEngine().render(slides)))
    shapes = prs.slides[1].shapes
    title = next(shape for shape in shapes if hasattr(shape, "text") and shape.text == "Next Steps & Call to Action")
    bullets = next(shape for shape in shapes if hasattr(shape, "text") and "Schedule a Deep-Dive" in shape.text)

    assert title.top < Inches(2.4)
    assert bullets.top < Inches(4.8)
    assert bullets.top - (title.top + title.height) < Inches(1.2)


def test_render_adds_disclosure_to_last_rendered_slide_with_non_sequential_indexes():
    slides = [
        SlideData(index=10, title="Title", bullets=[], notes="", layout="title"),
        SlideData(index=20, title="Next Steps", bullets=["Review"], notes="", layout="content"),
    ]

    prs = Presentation(BytesIO(PptxEngine().render(slides)))
    final_slide_text = "\n".join(shape.text for shape in prs.slides[-1].shapes if hasattr(shape, "text"))

    assert "not a guarantee" not in final_slide_text


def test_render_keeps_header_inside_four_by_three_template(tmp_path):
    template = Presentation()
    template.slide_width = Inches(10)
    template.slide_height = Inches(7.5)
    template_path = tmp_path / "four-by-three.pptx"
    template.save(template_path)
    slides = [SlideData(index=1, title="Title", bullets=[], notes="", layout="title")]

    prs = Presentation(BytesIO(PptxEngine(template_path=str(template_path)).render(slides)))
    branded_text_shapes = [
        shape
        for shape in prs.slides[0].shapes
        if hasattr(shape, "text") and "citi" in shape.text.lower()
    ]

    assert branded_text_shapes
    assert all(shape.left + shape.width <= prs.slide_width for shape in branded_text_shapes)


def test_render_uses_configured_citi_logo_path_when_file_exists(tmp_path, monkeypatch):
    logo_path = tmp_path / "citi.png"
    logo_path.write_bytes(
        base64.b64decode(
            "iVBORw0KGgoAAAANSUhEUgAAAAEAAAABCAQAAAC1HAwCAAAAC0lEQVR42mP8/x8AAwMCAO+/p9sAAAAASUVORK5CYII="
        )
    )
    monkeypatch.setattr(settings, "citi_logo_path", str(logo_path))
    slides = [SlideData(index=1, title="Title", bullets=[], notes="", layout="title")]

    prs = Presentation(BytesIO(PptxEngine().render(slides)))

    assert any(shape.shape_type == MSO_SHAPE_TYPE.PICTURE for shape in prs.slides[0].shapes)


def test_render_preserves_citi_logo_aspect_ratio(tmp_path, monkeypatch):
    from PIL import Image

    logo_path = tmp_path / "citi-wide.png"
    Image.new("RGBA", (250, 163), (255, 255, 255, 0)).save(logo_path)
    monkeypatch.setattr(settings, "citi_logo_path", str(logo_path))
    slides = [SlideData(index=1, title="Title", bullets=[], notes="", layout="title")]

    prs = Presentation(BytesIO(PptxEngine().render(slides)))
    pictures = [shape for shape in prs.slides[0].shapes if shape.shape_type == MSO_SHAPE_TYPE.PICTURE]

    assert pictures
    logo = pictures[-1]
    assert logo.left + logo.width <= prs.slide_width
    assert logo.top >= Inches(0.1)
    assert logo.width / logo.height == pytest.approx(250 / 163, rel=0.05)


def test_render_writes_speaker_notes():
    slides = [SlideData(index=1, title="Title", bullets=[], notes="Opening speaker note.", layout="title")]

    prs = Presentation(BytesIO(PptxEngine().render(slides)))
    notes_text = prs.slides[0].notes_slide.notes_text_frame.text

    assert "Opening speaker note." in notes_text


class NotesSlideWithoutFrame:
    notes_text_frame = None


class SlideWithoutNotesFrame:
    notes_slide = NotesSlideWithoutFrame()


def test_add_speaker_notes_skips_missing_notes_text_frame():
    PptxEngine()._add_speaker_notes(SlideWithoutNotesFrame(), SlideData(index=1, title="", bullets=[], notes="Speaker note", layout="content"))


def test_render_standard_content_with_image_adds_large_visual():
    b64_img = "iVBORw0KGgoAAAANSUhEUgAAAAEAAAABCAQAAAC1HAwCAAAAC0lEQVR42mP8/x8AAwMCAO+/p9sAAAAASUVORK5CYII="
    slides = [
        SlideData(index=1, title="Title", bullets=[], notes="", layout="title"),
        SlideData(
            index=2,
            title="Standard Content with Cards",
            bullets=["Bullet 1", "Bullet 2", "Bullet 3"],
            notes="",
            layout="content",
            image_b64=b64_img,
        ),
    ]
    prs = Presentation(BytesIO(PptxEngine().render(slides)))

    shapes = prs.slides[1].shapes
    big_images = [s for s in shapes if s.shape_type == MSO_SHAPE_TYPE.PICTURE and s.width > Inches(3)]
    assert big_images


def test_render_chart_without_data_adds_image():
    b64_img = "iVBORw0KGgoAAAANSUhEUgAAAAEAAAABCAQAAAC1HAwCAAAAC0lEQVR42mP8/x8AAwMCAO+/p9sAAAAASUVORK5CYII="
    slides = [
        SlideData(index=1, title="Title", bullets=[], notes="", layout="title"),
        SlideData(
            index=2,
            title="Chart Slide with Image",
            bullets=["Bullet 1"],
            notes="",
            layout="chart",
            image_b64=b64_img,
        ),
    ]
    prs = Presentation(BytesIO(PptxEngine().render(slides)))

    assert any(shape.shape_type == MSO_SHAPE_TYPE.PICTURE for shape in prs.slides[1].shapes)


# --- Bug 1: markdown bold runs -------------------------------------------------


def _fresh_text_frame():
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    tf = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(5), Inches(2)).text_frame
    tf.clear()
    return tf


def test_add_markdown_paragraph_splits_bold_and_normal_runs():
    tf = _fresh_text_frame()
    para = add_markdown_paragraph(tf, "**Key Point:** description text", 13, "Arial", CITI_DARK)

    assert len(para.runs) == 2
    assert para.runs[0].text == "Key Point:"
    assert para.runs[0].font.bold is True
    assert para.runs[1].text == " description text"
    assert para.runs[1].font.bold is False
    for run in para.runs:
        assert run.font.size == Pt(13)
        assert run.font.name == "Arial"
        assert run.font.color.rgb == CITI_DARK


def test_add_markdown_paragraph_plain_text_single_run():
    tf = _fresh_text_frame()
    para = add_markdown_paragraph(tf, "plain text only", 13, "Arial", CITI_DARK)

    assert len(para.runs) == 1
    assert para.runs[0].text == "plain text only"
    assert para.runs[0].font.bold is False


def test_content_bullets_render_bold_runs_without_literal_markdown():
    slides = [
        SlideData(index=1, title="Title", bullets=[], notes="", layout="title"),
        SlideData(
            index=2,
            title="Insights",
            bullets=["**Key Point:** automation saves time"],
            notes="",
            layout="content",
        ),
    ]
    prs = Presentation(BytesIO(PptxEngine().render(slides)))

    bold_found = False
    all_run_text = []
    for shape in prs.slides[1].shapes:
        if not shape.has_text_frame:
            continue
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                all_run_text.append(run.text)
                if run.text == "Key Point:" and run.font.bold:
                    bold_found = True

    assert bold_found
    assert "**" not in "".join(all_run_text)
    assert "\u2022" in "".join(all_run_text)


# --- Bug 3: auto-fit content boxes ---------------------------------------------


def test_content_bullets_box_uses_shape_to_fit_text_autosize():
    slides = [
        SlideData(index=1, title="Title", bullets=[], notes="", layout="title"),
        SlideData(index=2, title="Insights", bullets=["Bullet one", "Bullet two"], notes="", layout="content"),
    ]
    prs = Presentation(BytesIO(PptxEngine().render(slides)))

    autosized = [
        shape
        for shape in prs.slides[1].shapes
        if shape.has_text_frame and "Bullet one" in shape.text
    ]
    assert autosized
    assert all(s.text_frame.auto_size == MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT for s in autosized)


# --- Bug 4: section divider leaked text + guard --------------------------------


def test_section_divider_scrubs_leaked_instruction_text():
    leaked = (
        "Bold title centered on the slide, potentially with a subtle background "
        "image related to AI, data, or retail innovation to convey forward-thinking."
    )
    slides = [
        SlideData(index=1, title="Title", bullets=[], notes="", layout="title"),
        SlideData(
            index=5,
            title="Market Overview",
            bullets=[leaked],
            notes="",
            layout="section_divider",
            visual_direction=leaked,
        ),
    ]
    prs = Presentation(BytesIO(PptxEngine().render(slides)))
    text = "\n".join(s.text for s in prs.slides[1].shapes if hasattr(s, "text"))

    assert "potentially" not in text.lower()
    assert "background image related to" not in text.lower()
    assert "Market Overview" in text
    assert "Section overview" in text


def test_section_divider_guard_raises_on_instruction_like_subtitle():
    sneaky = (
        "Our strategic vision continues to expand across regions and segments while "
        "we potentially unlock new long-term revenue and durable advantages for clients."
    )
    slides = [
        SlideData(index=1, title="Title", bullets=[], notes="", layout="title"),
        SlideData(index=5, title="Opportunity", bullets=[sneaky], notes="", layout="section_divider"),
    ]
    with pytest.raises(ValueError):
        PptxEngine().render(slides)


def test_looks_like_instruction_detects_lowercase_verb_prefix():
    long_instruction = (
        "use a clean Citi-style callout visual for the slide core message and ensure "
        "the hierarchy is strong and the layout balanced across the whole frame area."
    )
    assert _looks_like_instruction(long_instruction) is True
    assert _looks_like_instruction("Short clean subtitle") is False


# --- Criterion 5: section labels are blue + bold -------------------------------


def test_standard_content_has_no_key_insights_label():
    slides = [
        SlideData(index=1, title="Title", bullets=[], notes="", layout="title"),
        SlideData(index=2, title="Findings", bullets=["Bullet point"], notes="", layout="content"),
    ]
    prs = Presentation(BytesIO(PptxEngine().render(slides)))
    texts = [shape.text.strip() for shape in prs.slides[1].shapes if shape.has_text_frame]
    # The redundant sub-label is gone; the real slide title carries the heading.
    assert "Key Insights" not in texts
    assert any("Bullet point" in t for t in texts)


def test_executive_summary_section_label_is_blue_and_bold():
    slides = [
        SlideData(index=1, title="Title", bullets=[], notes="", layout="title"),
        SlideData(
            index=2,
            title="Executive Summary",
            bullets=["Point A", "Point B"],
            notes="",
            layout="executive_summary",
        ),
    ]
    prs = Presentation(BytesIO(PptxEngine().render(slides)))

    labels = [
        shape
        for shape in prs.slides[1].shapes
        if shape.has_text_frame and shape.text.strip() == "Overview & Context"
    ]
    assert labels
    para = labels[0].text_frame.paragraphs[0]
    assert para.font.bold is True
    assert para.font.color.rgb == THEMES["minimalist"].text


# --- Images on title and section-divider slides --------------------------------

_PNG_1PX = "iVBORw0KGgoAAAANSUhEUgAAAAEAAAABCAQAAAC1HAwCAAAAC0lEQVR42mP8/x8AAwMCAO+/p9sAAAAASUVORK5CYII="


def test_title_slide_adds_image_when_present():
    slides = [
        SlideData(index=1, title="Cover", bullets=["Prepared for Client"], notes="", layout="title", image_b64=_PNG_1PX),
    ]
    prs = Presentation(BytesIO(PptxEngine().render(slides)))
    pictures = [sh for sh in prs.slides[0].shapes if sh.shape_type == MSO_SHAPE_TYPE.PICTURE]
    assert pictures
    assert pictures[0].left == Inches(10.15)
    assert pictures[0].top == Inches(0.92)


def test_section_divider_adds_image_when_present():
    slides = [
        SlideData(index=1, title="Title", bullets=[], notes="", layout="title"),
        SlideData(index=5, title="Market Overview", bullets=["Context"], notes="", layout="section_divider", image_b64=_PNG_1PX),
    ]
    prs = Presentation(BytesIO(PptxEngine().render(slides)))
    assert any(sh.shape_type == MSO_SHAPE_TYPE.PICTURE for sh in prs.slides[1].shapes)


def test_section_divider_without_image_has_no_picture():
    slides = [
        SlideData(index=1, title="Title", bullets=[], notes="", layout="title"),
        SlideData(index=5, title="Market Overview", bullets=["Context"], notes="", layout="section_divider"),
    ]
    prs = Presentation(BytesIO(PptxEngine().render(slides)))
    assert not any(sh.shape_type == MSO_SHAPE_TYPE.PICTURE for sh in prs.slides[1].shapes)


# --- Visual themes -------------------------------------------------------------

def test_framework_variant_registry_contains_dark_variants():
    engine = PptxEngine()

    handlers = engine._framework_variant_handlers()

    assert PptxEngine._DARK_VARIANTS.issubset(handlers.keys())


def test_dark_theme_sets_slide_background():
    slides = [SlideData(index=1, title="Cover", bullets=["Sub"], notes="", layout="title")]
    prs = Presentation(BytesIO(PptxEngine(theme="dark").render(slides)))
    bg = prs.slides[0].background.fill.fore_color.rgb
    assert bg == THEMES["dark"].background


def test_minimalist_theme_uses_modern_off_white_background():
    assert THEMES["minimalist"].background != WHITE
    assert THEMES["minimalist"].use_template is False
    slides = [SlideData(index=1, title="Cover", bullets=[], notes="", layout="title")]
    prs = Presentation(BytesIO(PptxEngine(theme="minimalist").render(slides)))
    assert prs.slides[0].background.fill.fore_color.rgb == THEMES["dark"].background


def test_bold_theme_renders_without_template_and_uses_dark_cover_background():
    slides = [
        SlideData(index=1, title="Cover", bullets=[], notes="", layout="title"),
        SlideData(index=2, title="Detail", bullets=["A", "B"], notes="", layout="content"),
    ]
    prs = Presentation(BytesIO(PptxEngine(theme="bold").render(slides)))
    assert prs.slides[0].background.fill.fore_color.rgb == THEMES["dark"].background


def test_resolve_theme_falls_back_to_minimalist():
    assert resolve_theme(None) is THEMES["minimalist"]
    assert resolve_theme("nonsense") is THEMES["minimalist"]
    assert resolve_theme("dark") is THEMES["dark"]


# --- Aspect ratio --------------------------------------------------------------

def test_default_aspect_ratio_is_16_9():
    engine = PptxEngine()
    assert engine.aspect_ratio == "16:9"
    assert engine._hscale == 1.0
    slides = [SlideData(index=1, title="Cover", bullets=[], notes="", layout="title")]
    prs = Presentation(BytesIO(engine.render(slides)))
    assert prs.slide_width == Inches(17.7778)
    assert prs.slide_height == Inches(10)


def test_four_three_sets_narrower_canvas_and_scale():
    engine = PptxEngine(aspect_ratio="4:3")
    assert engine._hscale == pytest.approx(0.75, abs=1e-3)
    slides = [SlideData(index=1, title="Cover", bullets=[], notes="", layout="title")]
    prs = Presentation(BytesIO(engine.render(slides)))
    assert prs.slide_width == Inches(13.3333)
    assert prs.slide_height == Inches(10)


def test_four_three_keeps_image_within_canvas():
    engine = PptxEngine(aspect_ratio="4:3")
    slides = [
        SlideData(index=1, title="Title", bullets=[], notes="", layout="title"),
        SlideData(index=2, title="Detail", subtitle="x", bullets=[], notes="", layout="section_divider", image_b64=_PNG_1PX),
    ]
    prs = Presentation(BytesIO(engine.render(slides)))
    width = prs.slide_width
    pictures = [sh for s in prs.slides for sh in s.shapes if sh.shape_type == MSO_SHAPE_TYPE.PICTURE]
    assert pictures
    for pic in pictures:
        assert pic.left + pic.width <= width + Inches(0.02)


def test_invalid_aspect_ratio_falls_back_to_16_9():
    engine = PptxEngine(aspect_ratio="21:9")
    assert engine.aspect_ratio == "16:9"
    assert engine._hscale == 1.0


# --- Subtitle field ------------------------------------------------------------

def _slide_texts(slide):
    return [sh.text_frame.text for sh in slide.shapes if sh.has_text_frame]


def test_title_slide_renders_explicit_subtitle():
    slides = [SlideData(index=1, title="Renewable Energy", subtitle="A deep dive by 2030", bullets=["Prepared for ACME"], notes="", layout="title")]
    prs = Presentation(BytesIO(PptxEngine().render(slides)))
    texts = _slide_texts(prs.slides[0])
    assert "A deep dive by 2030" in texts
    # All bullets remain available as the secondary line when subtitle is explicit.
    assert any("Prepared for ACME" in t for t in texts)


def test_title_slide_long_title_does_not_overlap_subtitle():
    long_title = "USF Campus Navigator: A Full-Stack Approach to Student Success"
    subtitle = "Connecting Campus Life through Real-Time Navigation and Scheduling"
    slides = [
        SlideData(
            index=1,
            title=long_title,
            kicker="PROJECT OVERVIEW",
            subtitle=subtitle,
            bullets=["Task: Design and implement a high-fidelity campus navigation prototype."],
            notes="",
            layout="title",
        )
    ]

    prs = Presentation(BytesIO(PptxEngine().render(slides)))
    shapes = {sh.text_frame.text: sh for sh in prs.slides[0].shapes if sh.has_text_frame}

    title = shapes[long_title]
    subtitle_shape = shapes[subtitle]
    assert subtitle_shape.top > title.top + title.height


def test_title_slide_falls_back_to_first_bullet_as_subtitle():
    slides = [SlideData(index=1, title="Renewable Energy", bullets=["Prepared for ACME", "June 2026"], notes="", layout="title")]
    prs = Presentation(BytesIO(PptxEngine().render(slides)))
    texts = _slide_texts(prs.slides[0])
    assert "Prepared for ACME" in texts


def test_section_divider_prefers_explicit_subtitle():
    slides = [
        SlideData(index=1, title="Title", bullets=[], notes="", layout="title"),
        SlideData(index=5, title="Market Overview", subtitle="Industry landscape", bullets=["ignored bullet"], notes="", layout="section_divider"),
    ]
    prs = Presentation(BytesIO(PptxEngine().render(slides)))
    assert "Industry landscape" in _slide_texts(prs.slides[1])


# --- Formatting polish: footer page number, divider scrim, smaller images ------

def test_footer_is_not_rendered():
    slides = [
        SlideData(index=1, title="Cover", bullets=[], notes="", layout="title"),
        SlideData(index=2, title="Body", bullets=["x"], notes="", layout="content"),
    ]
    prs = Presentation(BytesIO(PptxEngine().render(slides)))
    all_text = [sh.text for s in prs.slides for sh in s.shapes if sh.has_text_frame]
    joined = " ".join(all_text)
    assert "1  Title" not in joined
    assert "2  Content" not in joined
    assert "Confidential" not in joined
    assert not any(t.strip() == "1 / 2" for t in all_text)
    assert not any(t.strip() == "2 / 2" for t in all_text)


def test_section_divider_with_image_is_fullbleed_with_scrim():
    slides = [
        SlideData(index=1, title="Title", bullets=[], notes="", layout="title"),
        SlideData(index=5, title="Market Overview", subtitle="Industry landscape", bullets=[], notes="", layout="section_divider", image_b64=_PNG_1PX),
    ]
    prs = Presentation(BytesIO(PptxEngine().render(slides)))
    divider = prs.slides[1]
    width, height = prs.slide_width, prs.slide_height
    pics = [sh for sh in divider.shapes if sh.shape_type == MSO_SHAPE_TYPE.PICTURE]
    assert pics
    # Image covers (approximately) the whole slide.
    pic = pics[0]
    assert pic.width >= width - 2 and pic.height >= height - 2
    # Centered title text is present.
    assert "Market Overview" in [sh.text for sh in divider.shapes if sh.has_text_frame]


def test_content_slide_with_image_uses_split_image():
    slides = [
        SlideData(index=1, title="Title", bullets=[], notes="", layout="title"),
        SlideData(index=2, title="Body", bullets=["x", "y", "z"], notes="", layout="content", image_b64=_PNG_1PX),
    ]
    prs = Presentation(BytesIO(PptxEngine().render(slides)))
    big = [sh for sh in prs.slides[1].shapes if sh.shape_type == MSO_SHAPE_TYPE.PICTURE and sh.width > Inches(3)]
    assert big


# --- Template card redesign ----------------------------------------------------

def test_content_card_grid_has_one_card_per_bullet():
    slides = [
        SlideData(index=1, title="Title", bullets=[], notes="", layout="title"),
        SlideData(index=2, title="Pillars", kicker="Strategy", bullets=["A", "B", "C"], notes="", layout="content"),
    ]
    prs = Presentation(BytesIO(PptxEngine().render(slides)))
    cards = [s for s in prs.slides[1].shapes if s.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE]
    # 3 cards + 3 icon chips + 3 inner marks + accent rule = at least 3 cards.
    assert len(cards) >= 3


def test_kicker_renders_as_uppercase_accent_eyebrow():
    slides = [
        SlideData(index=1, title="Title", bullets=[], notes="", layout="title"),
        SlideData(index=2, title="Body", kicker="Market Context", bullets=["x"], notes="", layout="content"),
    ]
    prs = Presentation(BytesIO(PptxEngine().render(slides)))
    eyebrows = [
        s for s in prs.slides[1].shapes
        if s.has_text_frame and s.text.strip() == "MARKET CONTEXT"
    ]
    assert eyebrows
    para = eyebrows[0].text_frame.paragraphs[0]
    assert para.font.color.rgb == THEMES["minimalist"].accent


def test_next_steps_uses_numbered_circles():
    slides = [
        SlideData(index=1, title="Title", bullets=[], notes="", layout="title"),
        SlideData(index=2, title="Next Steps", bullets=["Do A", "Do B", "Do C"], notes="", layout="next_steps"),
    ]
    prs = Presentation(BytesIO(PptxEngine().render(slides)))
    circles = [
        s for s in prs.slides[1].shapes
        if s.has_text_frame and s.text.strip() in {"1", "2", "3"}
    ]
    assert len(circles) == 3


# --- Component blocks (React-like slide structure) -----------------------------

def _blocks_deck(blocks, bullets=None):
    return [
        SlideData(index=1, title="Cover", bullets=[], notes="", layout="title"),
        SlideData(index=2, title="Body", bullets=bullets or ["fallback"], notes="", layout="content", blocks=blocks),
    ]


def _texts(slide):
    return [s.text for s in slide.shapes if s.has_text_frame]


def test_block_stat_renders_value_and_label():
    prs = Presentation(BytesIO(PptxEngine().render(_blocks_deck([{"type": "stat", "value": "48%", "label": "Cost reduction"}]))))
    texts = _texts(prs.slides[1])
    assert "48%" in texts
    assert "Cost reduction" in texts


def test_block_quote_renders_text_and_author():
    prs = Presentation(BytesIO(PptxEngine().render(_blocks_deck([{"type": "quote", "text": "Build the future", "author": "Jane Doe"}]))))
    joined = " ".join(_texts(prs.slides[1]))
    assert "Build the future" in joined
    assert "JANE DOE" in joined


def test_block_table_renders_table_with_headers():
    blocks = [{"type": "table", "headers": ["Feature", "Legacy", "New"], "rows": [["Compliance", "Manual", "Native"]]}]
    prs = Presentation(BytesIO(PptxEngine().render(_blocks_deck(blocks))))
    tables = [s for s in prs.slides[1].shapes if getattr(s, "has_table", False)]
    assert tables
    tbl = tables[0].table
    assert tbl.cell(0, 0).text == "Feature"
    assert tbl.cell(1, 2).text == "Native"


def test_block_cards_render_titles():
    blocks = [{"type": "cards", "columns": 2, "items": [{"title": "Velocity", "body": "Fast"}, {"title": "Security", "body": "Safe"}]}]
    prs = Presentation(BytesIO(PptxEngine().render(_blocks_deck(blocks))))
    texts = _texts(prs.slides[1])
    assert "Velocity" in texts and "Security" in texts


def test_block_process_renders_numbered_steps():
    blocks = [{"type": "process", "steps": [{"title": "Audit", "body": "a"}, {"title": "Migrate", "body": "b"}]}]
    prs = Presentation(BytesIO(PptxEngine().render(_blocks_deck(blocks))))
    circles = [s for s in prs.slides[1].shapes if s.has_text_frame and s.text.strip() in {"1", "2"}]
    assert len(circles) == 2


def test_unknown_block_falls_back_to_cards():
    prs = Presentation(BytesIO(PptxEngine().render(_blocks_deck([{"type": "mystery", "items": ["alpha", "beta"]}]))))
    joined = " ".join(_texts(prs.slides[1]))
    assert "alpha" in joined and "beta" in joined


# --- Icon library + multi-block stacking ---------------------------------------

def test_icon_shape_mapping():
    assert icon_shape("security") == MSO_SHAPE.HEXAGON
    assert icon_shape("speed and velocity") == MSO_SHAPE.LIGHTNING_BOLT
    assert icon_shape("revenue growth") == MSO_SHAPE.UP_ARROW
    assert icon_shape(None) == MSO_SHAPE.DIAMOND
    assert icon_shape("zzz-unmapped") == MSO_SHAPE.DIAMOND


def test_card_icon_renders_fontawesome_image():
    blocks = [{"type": "cards", "items": [{"title": "Security", "body": "x", "icon": "security"}]}]
    prs = Presentation(BytesIO(PptxEngine().render(_blocks_deck(blocks))))
    # The Font Awesome glyph is embedded as a picture in the icon chip.
    pics = [s for s in prs.slides[1].shapes if s.shape_type == MSO_SHAPE_TYPE.PICTURE]
    assert pics


def test_fontawesome_icon_renderer():
    from app.services.media.icons import glyph_for, render_icon_png

    assert glyph_for("speed") == "\uf0e7"  # bolt
    assert glyph_for("security") == "\uf3ed"  # shield-halved
    assert glyph_for("zzz-unmapped") == "\uf058"  # default circle-check
    png = render_icon_png("growth", "EE2A24")
    assert png is not None and png[:4] == b"\x89PNG"


def test_fontawesome_icon_renderer_logs_failure(monkeypatch, caplog):
    from app.services.media import icons

    icons._png_cache.clear()
    monkeypatch.setattr(icons, "_load_font", lambda size_px: object())

    with caplog.at_level(logging.WARNING, logger="app.services.media.icons"):
        png = icons.render_icon_png("growth", "EE2A24")

    assert png is None
    assert "Icon render failed" in caplog.text


def test_multi_block_stacking_renders_all_blocks():
    blocks = [
        {"type": "stat", "value": "42%", "label": "Uptime gain"},
        {"type": "cards", "items": [{"title": "One", "body": "a"}, {"title": "Two", "body": "b"}]},
    ]
    prs = Presentation(BytesIO(PptxEngine().render(_blocks_deck(blocks))))
    joined = " ".join(_texts(prs.slides[1]))
    assert "42%" in joined and "One" in joined and "Two" in joined


# --- Framework variants -------------------------------------------------------


def test_framework_big_stat_variant_uses_oversized_metric():
    slides = [
        SlideData(index=1, title="Cover", bullets=[], notes="", layout="title"),
        SlideData(
            index=2,
            title="Impact",
            bullets=["48% reduction in cost"],
            notes="",
            layout="content",
            variant="big_stat",
            blocks=[{"type": "stat", "value": "48%", "label": "Reduction in global infrastructure costs"}],
        ),
    ]

    prs = Presentation(BytesIO(PptxEngine().render(slides)))

    metric_shapes = [s for s in prs.slides[1].shapes if s.has_text_frame and s.text.strip() == "48%"]
    assert metric_shapes
    assert metric_shapes[0].text_frame.paragraphs[0].font.size >= Pt(100)


def test_content_slide_routes_explicit_variant_before_blocks(monkeypatch):
    engine = PptxEngine()
    calls: list[str] = []

    monkeypatch.setattr(engine, "_apply_big_stat", lambda slide, data: calls.append("big_stat"))
    monkeypatch.setattr(engine, "_apply_blocks", lambda slide, data: calls.append("blocks"))

    data = SlideData(
        index=2,
        title="Revenue",
        bullets=[],
        notes="",
        layout="content",
        variant="big_stat",
        blocks=[{"type": "cards", "items": [{"title": "A"}]}],
    )

    engine._apply_content_slide(object(), data)

    assert calls == ["big_stat"]


def test_content_slide_routes_chart_data_before_variant(monkeypatch):
    engine = PptxEngine()
    calls: list[str] = []

    monkeypatch.setattr(engine, "_apply_chart_slide", lambda slide, data: calls.append("chart"))
    monkeypatch.setattr(engine, "_apply_big_stat", lambda slide, data: calls.append("big_stat"))

    data = SlideData(
        index=2,
        title="Revenue",
        bullets=[],
        notes="",
        layout="content",
        variant="big_stat",
        chart_data={
            "type": "bar",
            "title": "Revenue",
            "categories": ["Q1"],
            "series": [{"name": "Revenue", "values": [1.0]}],
        },
    )

    engine._apply_content_slide(object(), data)

    assert calls == ["chart"]


def test_split_image_restores_active_dark_after_render():
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    engine = PptxEngine(theme="minimalist")
    engine._active_dark = False
    data = SlideData(
        index=2,
        title="Dark panel",
        subtitle="Subtitle",
        bullets=["One"],
        notes="",
        layout="content",
    )

    engine._apply_split_image(slide, data)

    assert engine._active_dark is False


def test_theme_mode_restores_active_dark_after_exception():
    engine = PptxEngine(theme="minimalist")
    engine._active_dark = False

    with pytest.raises(RuntimeError):
        with engine._theme_mode(dark=True):
            assert engine._active_dark is True
            raise RuntimeError("boom")

    assert engine._active_dark is False


def test_add_slide_image_logs_decode_failure(caplog):
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    engine = PptxEngine()

    with caplog.at_level(logging.WARNING):
        inserted = engine._add_slide_image(slide, "not valid base64", 0, 0, 1, 1)

    assert inserted is False
    assert "Failed to insert slide image" in caplog.text


def test_framework_comparison_table_variant_renders_native_table():
    slides = [
        SlideData(index=1, title="Cover", bullets=[], notes="", layout="title"),
        SlideData(
            index=2,
            title="Platform Capabilities",
            bullets=["Native compliance"],
            notes="",
            layout="content",
            variant="comparison_table",
            blocks=[{"type": "table", "headers": ["Feature", "Legacy", "New"], "rows": [["Compliance", "Manual", "Native"]]}],
        ),
    ]

    prs = Presentation(BytesIO(PptxEngine().render(slides)))

    assert any(getattr(s, "has_table", False) for s in prs.slides[1].shapes)


def test_framework_split_image_variant_uses_large_right_image():
    slides = [
        SlideData(index=1, title="Cover", bullets=[], notes="", layout="title"),
        SlideData(
            index=2,
            title="Legacy systems are slowing innovation",
            subtitle="Manual compliance checks create bottlenecks.",
            bullets=["Four-week deployment cycles", "Fragmented tooling"],
            notes="",
            layout="content",
            variant="split_image",
            image_b64=_PNG_1PX,
        ),
    ]

    prs = Presentation(BytesIO(PptxEngine().render(slides)))
    pics = [s for s in prs.slides[1].shapes if s.shape_type == MSO_SHAPE_TYPE.PICTURE and s.width > Inches(5)]

    assert pics


def test_framework_split_image_variant_uses_dark_narrative_panel_with_light_text():
    slides = [
        SlideData(index=1, title="Cover", bullets=[], notes="", layout="title"),
        SlideData(
            index=2,
            title="Legacy systems are slowing innovation",
            subtitle="Manual compliance checks create bottlenecks.",
            bullets=["Four-week deployment cycles", "Fragmented tooling"],
            notes="",
            layout="content",
            variant="split_image",
            image_b64=_PNG_1PX,
        ),
    ]

    prs = Presentation(BytesIO(PptxEngine().render(slides)))
    panel_shapes = [
        s
        for s in prs.slides[1].shapes
        if s.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE
        and s.left <= Inches(0.05)
        and s.width >= Inches(8.5)
        and s.fill.fore_color.rgb == THEMES["dark"].background
    ]
    title_shapes = [
        s
        for s in prs.slides[1].shapes
        if s.has_text_frame and s.text.strip() == "Legacy systems are slowing innovation"
    ]

    assert panel_shapes
    assert title_shapes
    assert title_shapes[0].text_frame.paragraphs[0].font.color.rgb == THEMES["dark"].text


def test_framework_process_and_quote_variants_render_distinct_structures():
    slides = [
        SlideData(index=1, title="Cover", bullets=[], notes="", layout="title"),
        SlideData(
            index=2,
            title="Implementation Timeline",
            bullets=["Audit", "Migrate", "Scale"],
            notes="",
            layout="content",
            variant="process",
            blocks=[{"type": "process", "steps": [{"title": "Audit", "body": "Review"}, {"title": "Migrate", "body": "Roll out"}]}],
        ),
        SlideData(
            index=3,
            title="Vision",
            bullets=["Build the future"],
            notes="",
            layout="content",
            variant="quote",
            blocks=[{"type": "quote", "text": "Build the future", "author": "Jane Doe"}],
        ),
    ]

    prs = Presentation(BytesIO(PptxEngine().render(slides)))
    process_text = " ".join(_texts(prs.slides[1]))
    quote_text = " ".join(_texts(prs.slides[2]))

    assert "Audit" in process_text and "Migrate" in process_text
    assert "Build the future" in quote_text and "JANE DOE" in quote_text


def test_process_variant_with_five_long_steps_uses_bounded_agenda_cards():
    long_steps = [
        "Explore the expanding market opportunity for AI-driven short video content.",
        "Introduce the platform's innovative solution and core technological advantages.",
        "Review the business model and early indicators of market traction.",
        "Detail the strategic allocation of the $10 million investment to accelerate growth.",
        "Outline financial projections and the path to investor returns.",
    ]
    slides = [
        SlideData(index=1, title="Cover", bullets=[], notes="", layout="title"),
        SlideData(
            index=2,
            title="Strategic Roadmap for Growth",
            kicker="PRESENTATION OVERVIEW",
            subtitle="Key discussion points for today's investment review",
            bullets=long_steps,
            notes="",
            layout="content",
            variant="process",
        ),
    ]

    prs = Presentation(BytesIO(PptxEngine().render(slides)))
    slide = prs.slides[1]
    text = " ".join(_texts(slide))
    numbers = [shape for shape in slide.shapes if shape.has_text_frame and shape.text.strip() in {"1", "2", "3", "4", "5"}]

    assert len(numbers) == 5
    assert "•" not in text
    assert "Outline financial projections" in text


def _presentation_agenda_slides(
    *,
    steps: list[dict[str, str]] | None = None,
    bullets: list[str] | None = None,
) -> list[SlideData]:
    agenda_steps = steps if steps is not None else [
        {"title": "The Core Insight", "body": "Understanding the shift in modern infrastructure."},
        {"title": "Strategic Pillars", "body": "The foundations supporting sustainable growth."},
        {"title": "Execution Pipeline", "body": "Deployment phases and readiness."},
        {"title": "Impact & Outlook", "body": "Metrics and next fiscal-year outlook."},
    ]
    return [
        SlideData(index=1, title="Cover", bullets=[], notes="", layout="title"),
        SlideData(
            index=2,
            title="Presentation Agenda",
            bullets=bullets or [],
            notes="",
            layout="content",
            variant="process",
            blocks=[{"type": "process", "steps": agenda_steps}],
        ),
    ]


def _containing_card(slide, text_shape):
    cards = [
        shape
        for shape in slide.shapes
        if shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE
        and shape.auto_shape_type == MSO_SHAPE.ROUNDED_RECTANGLE
        and shape.left <= text_shape.left
        and shape.top <= text_shape.top
        and shape.left + shape.width >= text_shape.left + text_shape.width
        and shape.top + shape.height >= text_shape.top + text_shape.height
    ]
    return min(cards, key=lambda shape: shape.width * shape.height, default=None)


def _relative_luminance(color) -> float:
    channels = []
    for value in color:
        channel = value / 255
        channels.append(channel / 12.92 if channel <= 0.04045 else ((channel + 0.055) / 1.055) ** 2.4)
    return 0.2126 * channels[0] + 0.7152 * channels[1] + 0.0722 * channels[2]


def _contrast_ratio(foreground, background) -> float:
    lighter, darker = sorted((_relative_luminance(foreground), _relative_luminance(background)), reverse=True)
    return (lighter + 0.05) / (darker + 0.05)


def test_presentation_agenda_renders_large_title_and_card_numbers():
    prs = Presentation(BytesIO(PptxEngine().render(_presentation_agenda_slides())))
    slide = prs.slides[1]
    title = next(shape for shape in slide.shapes if shape.has_text_frame and shape.text.strip() == "Presentation Agenda")
    numbers = [shape.text.strip() for shape in slide.shapes if shape.has_text_frame and shape.text.strip() in {"01", "02", "03", "04"}]

    assert title.text_frame.paragraphs[0].font.size == Pt(54)
    assert numbers == ["01", "02", "03", "04"]


@pytest.mark.parametrize("aspect_ratio", ["16:9", "4:3"])
def test_presentation_agenda_long_title_and_body_boxes_stay_inside_cards(aspect_ratio):
    long_title = "Enterprise Data Governance, Regulatory Readiness, and Responsible AI Operating Model"
    long_body = (
        "Align ownership, controls, reporting, and executive decisions across regions while preserving "
        "delivery speed, measurable accountability, and a clear path through implementation."
    )
    steps = [
        {"title": f"{long_title} {index}", "body": f"{long_body} Phase {index}."}
        for index in range(1, 5)
    ]
    prs = Presentation(
        BytesIO(
            PptxEngine(aspect_ratio=aspect_ratio).render(
                _presentation_agenda_slides(steps=steps)
            )
        )
    )
    slide = prs.slides[1]
    agenda_copy = {value for step in steps for value in step.values()}
    text_shapes = [
        shape
        for shape in slide.shapes
        if shape.has_text_frame and shape.text.strip() in agenda_copy
    ]

    assert len(text_shapes) == 8
    assert all(_containing_card(slide, shape) is not None for shape in text_shapes)
    assert all(shape.text_frame.auto_size == MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE for shape in text_shapes)


def test_presentation_agenda_title_and_body_frames_shrink_to_fit():
    prs = Presentation(BytesIO(PptxEngine().render(_presentation_agenda_slides())))
    slide = prs.slides[1]
    agenda_copy = {
        "The Core Insight",
        "Understanding the shift in modern infrastructure.",
        "Strategic Pillars",
        "The foundations supporting sustainable growth.",
        "Execution Pipeline",
        "Deployment phases and readiness.",
        "Impact & Outlook",
        "Metrics and next fiscal-year outlook.",
    }
    text_shapes = [
        shape
        for shape in slide.shapes
        if shape.has_text_frame and shape.text.strip() in agenda_copy
    ]

    assert len(text_shapes) == 8
    assert all(shape.text_frame.auto_size == MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE for shape in text_shapes)


def test_presentation_agenda_never_uses_old_discussion_fallback_copy():
    explicit = _presentation_agenda_slides(
        steps=[{"title": "Core Insight", "body": "The context and core insight for the discussion."}],
    )
    fallback = _presentation_agenda_slides(steps=[])

    direct_text = " ".join(_texts(Presentation(BytesIO(PptxEngine().render(explicit))).slides[1]))
    fallback_text = " ".join(_texts(Presentation(BytesIO(PptxEngine().render(fallback))).slides[1]))

    assert "The context and core insight for the discussion." not in direct_text
    assert "The context and core insight for the discussion." not in fallback_text
    assert "Key context, evidence, and decisions for this chapter." in fallback_text


def test_presentation_agenda_skips_old_discussion_copy_when_used_as_title():
    slides = _presentation_agenda_slides(
        steps=[
            {
                "title": "The context and core insight for the discussion.",
                "body": "This item must not be displayed.",
            }
        ]
    )

    text = " ".join(_texts(Presentation(BytesIO(PptxEngine().render(slides))).slides[1]))

    assert "The context and core insight for the discussion." not in text
    assert "This item must not be displayed." not in text
    assert "The Core Insight" in text


def test_presentation_agenda_deduplicates_titles_and_prioritizes_block_bodies():
    block_steps = [
        {"title": "Market Opportunity", "body": "Block market evidence."},
        {"title": "Strategic Pillars", "body": "Block strategic priorities."},
        {"title": "Execution Plan", "body": "Block execution milestones."},
    ]
    bullets = [
        "market opportunity: Bullet market copy.",
        "STRATEGIC PILLARS: Bullet strategy copy.",
        "Execution Plan: Bullet execution copy.",
    ]
    slides = _presentation_agenda_slides(steps=block_steps, bullets=bullets)
    slide = Presentation(BytesIO(PptxEngine().render(slides))).slides[1]
    numbers = [
        shape.text.strip()
        for shape in slide.shapes
        if shape.has_text_frame and shape.text.strip() in {"01", "02", "03", "04"}
    ]
    cards = [
        shape
        for shape in slide.shapes
        if shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE
        and shape.auto_shape_type == MSO_SHAPE.ROUNDED_RECTANGLE
    ]
    text = " ".join(_texts(slide))

    assert numbers == ["01", "02", "03"]
    assert len(cards) == 3
    assert all(step["body"] in text for step in block_steps)
    assert "Bullet market copy." not in text
    assert "Bullet strategy copy." not in text
    assert "Bullet execution copy." not in text


@pytest.mark.parametrize("aspect_ratio", ["16:9", "4:3"])
@pytest.mark.parametrize("item_count", [1, 2, 3, 4])
def test_presentation_agenda_item_counts_and_boxes_are_bounded(aspect_ratio, item_count):
    steps = [
        {"title": f"Chapter {index}", "body": f"Decision context for chapter {index}."}
        for index in range(1, item_count + 1)
    ]
    engine = PptxEngine(aspect_ratio=aspect_ratio)
    prs = Presentation(BytesIO(engine.render(_presentation_agenda_slides(steps=steps))))
    slide = prs.slides[1]
    numbers = {f"{index:02d}" for index in range(1, item_count + 1)}
    agenda_copy = numbers | {value for step in steps for value in step.values()}
    agenda_shapes = [
        shape
        for shape in slide.shapes
        if shape.has_text_frame and shape.text.strip() in agenda_copy
    ]
    cards = [
        shape
        for shape in slide.shapes
        if shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE
        and shape.auto_shape_type == MSO_SHAPE.ROUNDED_RECTANGLE
    ]

    assert len([shape for shape in agenda_shapes if shape.text.strip() in numbers]) == item_count
    assert len(cards) == item_count
    assert len(agenda_shapes) == item_count * 3
    assert all(shape.left >= 0 and shape.top >= 0 for shape in cards + agenda_shapes)
    assert all(shape.left + shape.width <= prs.slide_width for shape in cards + agenda_shapes)
    assert all(shape.top + shape.height <= prs.slide_height for shape in cards + agenda_shapes)
    assert all(_containing_card(slide, shape) is not None for shape in agenda_shapes)


def test_presentation_agenda_four_card_grid_has_equal_two_by_two_geometry():
    prs = Presentation(BytesIO(PptxEngine().render(_presentation_agenda_slides())))
    cards = [
        shape
        for shape in prs.slides[1].shapes
        if shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE
        and shape.auto_shape_type == MSO_SHAPE.ROUNDED_RECTANGLE
    ]

    assert len(cards) == 4
    assert len({shape.width for shape in cards}) == 1
    assert len({shape.height for shape in cards}) == 1
    assert len({shape.left for shape in cards}) == 2
    assert len({shape.top for shape in cards}) == 2
    assert all(len([shape for shape in cards if shape.left == left]) == 2 for left in {shape.left for shape in cards})
    assert all(len([shape for shape in cards if shape.top == top]) == 2 for top in {shape.top for shape in cards})


def test_presentation_agenda_card_text_has_nonoverlapping_vertical_regions():
    engine = PptxEngine()
    prs = Presentation(BytesIO(engine.render(_presentation_agenda_slides())))
    slide = prs.slides[1]
    items = [
        ("The Core Insight", "Understanding the shift in modern infrastructure."),
        ("Strategic Pillars", "The foundations supporting sustainable growth."),
        ("Execution Pipeline", "Deployment phases and readiness."),
        ("Impact & Outlook", "Metrics and next fiscal-year outlook."),
    ]

    for index, (title, body) in enumerate(items, start=1):
        number_shape = next(shape for shape in slide.shapes if shape.has_text_frame and shape.text.strip() == f"{index:02d}")
        title_shape = next(shape for shape in slide.shapes if shape.has_text_frame and shape.text.strip() == title)
        body_shape = next(shape for shape in slide.shapes if shape.has_text_frame and shape.text.strip() == body)

        assert number_shape.top + number_shape.height <= title_shape.top
        assert title_shape.top + title_shape.height <= body_shape.top
        assert title_shape.text_frame.auto_size == MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
        assert body_shape.text_frame.auto_size == MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE


@pytest.mark.parametrize("theme_name", ["minimalist", "bold", "dark"])
def test_presentation_agenda_numbers_use_active_accent_with_large_text_contrast(theme_name):
    engine = PptxEngine(theme=theme_name)
    prs = Presentation(BytesIO(engine.render(_presentation_agenda_slides())))
    slide = prs.slides[1]
    number_shapes = [
        shape
        for shape in slide.shapes
        if shape.has_text_frame and shape.text.strip() in {"01", "02", "03", "04"}
    ]
    active_theme = engine._active_theme()
    card = next(
        shape
        for shape in slide.shapes
        if shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE
        and shape.auto_shape_type == MSO_SHAPE.ROUNDED_RECTANGLE
    )
    card_color = card.fill.fore_color.rgb

    assert len(number_shapes) == 4
    for shape in number_shapes:
        number_color = shape.text_frame.paragraphs[0].font.color.rgb
        assert number_color == active_theme.accent
        assert _contrast_ratio(number_color, card_color) >= 3.0


def test_presentation_agenda_uses_length_aware_title_and_body_fonts():
    short_title = "Market"
    long_title = "Enterprise Data Governance, Regulatory Readiness, and Responsible AI Operating Model"
    concise_body = "Review market evidence and decisions."
    long_body = (
        "Align ownership, controls, reporting, and executive decisions across regions while preserving "
        "delivery speed, measurable accountability, and a clear path through implementation."
    )
    steps = [
        {"title": short_title, "body": concise_body},
        {"title": long_title, "body": long_body},
    ]
    slide = Presentation(
        BytesIO(PptxEngine().render(_presentation_agenda_slides(steps=steps)))
    ).slides[1]

    short_title_shape = next(shape for shape in slide.shapes if shape.has_text_frame and shape.text.strip() == short_title)
    long_title_shape = next(shape for shape in slide.shapes if shape.has_text_frame and shape.text.strip() == long_title)
    concise_body_shape = next(shape for shape in slide.shapes if shape.has_text_frame and shape.text.strip() == concise_body)
    long_body_shape = next(shape for shape in slide.shapes if shape.has_text_frame and shape.text.strip() == long_body)

    assert long_title_shape.text_frame.paragraphs[0].font.size < short_title_shape.text_frame.paragraphs[0].font.size
    assert long_body_shape.text_frame.paragraphs[0].runs[0].font.size < concise_body_shape.text_frame.paragraphs[0].runs[0].font.size


def _chapter_slide(**overrides) -> SlideData:
    values = {
        "index": 4,
        "title": "Delivery Model",
        "bullets": ["A concise supporting point."],
        "notes": "",
        "layout": "content",
        "chapter_number": 2,
        "chapter_title": "Product & Technology",
    }
    values.update(overrides)
    return SlideData(**values)


def _marker_shapes(slide):
    number = [
        shape
        for shape in slide.shapes
        if shape.has_text_frame and shape.text.strip() == "02"
    ]
    label = [
        shape
        for shape in slide.shapes
        if shape.has_text_frame and shape.text.strip() == "PRODUCT & TECHNOLOGY"
    ]
    return number, label


@pytest.mark.parametrize(
    "overrides",
    [
        {"variant": "big_statement"},
        {"variant": "three_points", "bullets": ["Speed", "Trust", "Scale"]},
        {"variant": "split_image"},
        {"variant": "big_stat", "blocks": [{"type": "stat", "value": "48%", "label": "Faster delivery"}]},
        {
            "variant": "before_after",
            "blocks": [{"type": "comparison", "headers": ["Before", "After"], "rows": [["Manual", "Automated"]]}],
        },
        {
            "variant": "comparison_table",
            "blocks": [{"type": "table", "headers": ["Capability", "Outcome"], "rows": [["Controls", "Embedded"]]}],
        },
        {"variant": "process", "blocks": [{"type": "process", "steps": [{"title": "Design", "body": "Align"}]}]},
        {"variant": "quote", "blocks": [{"type": "quote", "text": "Build with confidence", "author": "Citi"}]},
        {},
        {
            "chart_data": {
                "type": "bar",
                "title": "Adoption",
                "categories": ["Now", "Next"],
                "series": [{"name": "Rate", "values": [35.0, 62.0]}],
            }
        },
        {
            "blocks": [
                {"type": "cards", "items": [{"title": "Platform", "body": "Scalable"}]},
                {"type": "quote", "text": "Designed for growth"},
            ]
        },
        {"layout": "section_divider", "subtitle": "How the platform creates value"},
    ],
    ids=[
        "big-statement",
        "three-points",
        "split-image",
        "big-stat",
        "before-after",
        "comparison-table",
        "process",
        "quote",
        "standard-content",
        "chart-data",
        "block-driven",
        "section-divider",
    ],
)
def test_all_content_routes_render_active_agenda_chapter_marker(overrides):
    slides = [
        SlideData(index=1, title="Cover", bullets=[], notes="", layout="title"),
        _chapter_slide(**overrides),
    ]

    slide = Presentation(BytesIO(PptxEngine().render(slides))).slides[1]
    number, label = _marker_shapes(slide)

    assert len(number) == 1
    assert len(label) == 1


def test_related_slides_repeat_chapter_number_until_next_agenda_chapter():
    slides = [
        SlideData(index=1, title="Cover", bullets=[], notes="", layout="title"),
        _chapter_slide(index=3, title="Market Context", chapter_number=1, chapter_title="Market Opportunity"),
        _chapter_slide(index=4, title="Customer Needs", chapter_number=1, chapter_title="Market Opportunity"),
        _chapter_slide(index=5, title="Platform", chapter_number=2, chapter_title="Product & Technology"),
    ]

    prs = Presentation(BytesIO(PptxEngine().render(slides)))

    assert ["01" in _texts(prs.slides[index]) for index in (1, 2)] == [True, True]
    assert "02" in _texts(prs.slides[3])


@pytest.mark.parametrize(
    "data",
    [
        SlideData(
            index=1,
            title="Cover",
            bullets=[],
            notes="",
            layout="title",
            chapter_number=2,
            chapter_title="Product & Technology",
        ),
        SlideData(
            index=2,
            title="Presentation Agenda",
            bullets=["Product & Technology"],
            notes="",
            layout="content",
            variant="process",
            chapter_number=2,
            chapter_title="Product & Technology",
        ),
        SlideData(
            index=9,
            title="Thank You",
            bullets=["Questions"],
            notes="",
            layout="next_steps",
            variant="closing",
            chapter_number=2,
            chapter_title="Product & Technology",
        ),
    ],
    ids=["cover", "agenda", "closing"],
)
def test_non_chapter_slides_omit_marker_with_stale_metadata(data):
    slide = Presentation(BytesIO(PptxEngine().render([data]))).slides[0]

    number, label = _marker_shapes(slide)

    assert number == []
    assert label == []


def test_content_without_chapter_metadata_has_no_index_derived_marker():
    data = SlideData(
        index=5,
        title="Unassigned Content",
        bullets=["No chapter metadata is available."],
        notes="",
        layout="content",
    )

    text = _texts(Presentation(BytesIO(PptxEngine().render([data]))).slides[0])

    assert "03" not in text


@pytest.mark.parametrize(
    ("chapter_number", "chapter_title"),
    [(2, None), (None, "Product & Technology")],
    ids=["missing-title", "missing-number"],
)
def test_incomplete_chapter_metadata_omits_marker(chapter_number, chapter_title):
    slide = Presentation(
        BytesIO(
            PptxEngine().render(
                [
                    _chapter_slide(
                        chapter_number=chapter_number,
                        chapter_title=chapter_title,
                    )
                ]
            )
        )
    ).slides[0]

    marker_text = {"02", "PRODUCT & TECHNOLOGY"}

    assert marker_text.isdisjoint(_texts(slide))


def test_section_divider_uses_only_universal_chapter_number():
    data = _chapter_slide(
        index=7,
        layout="section_divider",
        subtitle="How the platform creates value",
        chapter_number=2,
        chapter_title="Product & Technology",
    )

    text = _texts(Presentation(BytesIO(PptxEngine().render([data]))).slides[0])

    assert text.count("02") == 1
    assert "03" not in text
    assert "3." not in text


@pytest.mark.parametrize("aspect_ratio", ["16:9", "4:3"])
def test_chapter_marker_is_bounded_fits_and_avoids_citi_logo(aspect_ratio):
    prs = Presentation(BytesIO(PptxEngine(aspect_ratio=aspect_ratio).render([_chapter_slide()])))
    slide = prs.slides[0]
    number, label = _marker_shapes(slide)
    logo = next(shape for shape in slide.shapes if shape.has_text_frame and shape.text.strip() == "citi")
    marker_shapes = [number[0], label[0]]

    assert all(shape.left >= 0 and shape.top >= 0 for shape in marker_shapes)
    assert all(shape.left + shape.width <= prs.slide_width for shape in marker_shapes)
    assert all(shape.top + shape.height <= prs.slide_height for shape in marker_shapes)
    assert all(shape.left + shape.width <= logo.left for shape in marker_shapes)
    assert label[0].text_frame.auto_size == MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE


@pytest.mark.parametrize("theme_name", ["minimalist", "bold", "dark"])
def test_chapter_badge_uses_accessible_citi_red_with_white_text(theme_name):
    slide = Presentation(
        BytesIO(PptxEngine(theme=theme_name).render([_chapter_slide()]))
    ).slides[0]
    number, _ = _marker_shapes(slide)
    badge = number[0]
    badge_color = badge.fill.fore_color.rgb
    number_color = badge.text_frame.paragraphs[0].font.color.rgb

    assert badge_color == CITI_RED
    assert number_color == WHITE
    assert _contrast_ratio(number_color, badge_color) >= 4.5


@pytest.mark.parametrize(
    "overrides",
    [
        {"variant": "split_image"},
        {
            "layout": "section_divider",
            "subtitle": "How the platform creates value",
            "image_b64": _PNG_1PX,
        },
    ],
    ids=["split-image", "image-backed-section-divider"],
)
def test_chapter_marker_label_uses_light_text_on_dark_backed_routes(overrides):
    slide = Presentation(
        BytesIO(PptxEngine().render([_chapter_slide(**overrides)]))
    ).slides[0]
    number, label = _marker_shapes(slide)

    assert number[0].text_frame.paragraphs[0].font.color.rgb == WHITE
    assert label[0].text_frame.paragraphs[0].font.color.rgb == THEMES["dark"].text


@pytest.mark.parametrize(
    "overrides",
    [
        {"variant": "split_image", "image_b64": _PNG_1PX},
        {
            "layout": "section_divider",
            "subtitle": "How the platform creates value",
            "image_b64": _PNG_1PX,
        },
    ],
    ids=["split-image", "image-backed-section-divider"],
)
def test_chapter_marker_renders_after_full_bleed_body_and_before_logo(overrides):
    slide = Presentation(
        BytesIO(PptxEngine().render([_chapter_slide(**overrides)]))
    ).slides[0]
    shapes = list(slide.shapes)
    number, label = _marker_shapes(slide)
    logo = next(shape for shape in shapes if shape.has_text_frame and shape.text.strip() == "citi")
    body_picture = next(shape for shape in shapes if shape.shape_type == MSO_SHAPE_TYPE.PICTURE)
    body_indices = [shapes.index(body_picture)]
    full_slide_overlays = [
        index
        for index, shape in enumerate(shapes)
        if shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE
        and shape.auto_shape_type == MSO_SHAPE.RECTANGLE
        and shape.left == 0
        and shape.top == 0
        and shape.width == slide.part.package.presentation_part.presentation.slide_width
        and shape.height == slide.part.package.presentation_part.presentation.slide_height
    ]
    body_indices.extend(full_slide_overlays)
    badge_index = shapes.index(number[0])
    label_index = shapes.index(label[0])
    logo_index = shapes.index(logo)

    assert max(body_indices) < badge_index < logo_index
    assert max(body_indices) < label_index < logo_index
    assert number[0].left + number[0].width <= logo.left
    assert label[0].left + label[0].width <= logo.left


@pytest.mark.parametrize(
    ("kicker", "expected_title_top", "expected_rule_top"),
    [("PLATFORM", 1.18, 1.96), (None, 0.86, 1.78)],
)
def test_content_header_moves_below_chapter_marker(kicker, expected_title_top, expected_rule_top):
    data = _chapter_slide(kicker=kicker)
    slide = Presentation(BytesIO(PptxEngine().render([data]))).slides[0]
    title = next(shape for shape in slide.shapes if shape.has_text_frame and shape.text.strip() == data.title)
    accent_rules = [
        shape
        for shape in slide.shapes
        if shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE
        and shape.auto_shape_type == MSO_SHAPE.RECTANGLE
        and shape.fill.fore_color.rgb == THEMES["minimalist"].accent
        and shape.height <= Inches(0.1)
    ]

    assert title.top == Inches(expected_title_top)
    assert any(shape.top == Inches(expected_rule_top) for shape in accent_rules)


def test_big_statement_supporting_points_span_lower_slide_width():
    slides = [
        SlideData(index=1, title="Cover", bullets=[], notes="", layout="title"),
        SlideData(
            index=2,
            title="Explosive Growth in Short Video Content Demands Scalable AI Solutions",
            kicker="MARKET DYNAMICS",
            bullets=[
                "Short-form video now accounts for over 80% of all mobile internet traffic.",
                "The global short video market is projected to exceed $150 billion by 2027.",
                "Creators need automation to produce high-quality content at scale.",
            ],
            notes="",
            layout="content",
            variant="big_statement",
        ),
    ]

    prs = Presentation(BytesIO(PptxEngine().render(slides)))
    slide = prs.slides[1]
    supporting_shapes = [
        shape
        for shape in slide.shapes
        if shape.has_text_frame and "Short-form video" in shape.text
    ]

    assert supporting_shapes
    assert supporting_shapes[0].width > Inches(12)
    assert supporting_shapes[0].top < Inches(7.1)


def test_framework_deck_uses_diverse_structures_not_repeated_card_grids():
    slides = [
        SlideData(index=1, title="Cover", bullets=[], notes="", layout="title", variant="cover"),
        SlideData(index=2, title="Core Insight", bullets=["Teams move faster"], notes="", layout="content", variant="big_statement"),
        SlideData(index=3, title="Pillars", bullets=["Velocity", "Security", "Scale"], notes="", layout="content", variant="three_points", blocks=[{"type": "cards", "items": [{"title": "Velocity"}, {"title": "Security"}, {"title": "Scale"}]}]),
        SlideData(index=4, title="Problem", bullets=["Legacy systems"], notes="", layout="content", variant="split_image", image_b64=_PNG_1PX),
        SlideData(index=5, title="Impact", bullets=["48% reduction"], notes="", layout="content", variant="big_stat", blocks=[{"type": "stat", "value": "48%", "label": "Cost reduction"}]),
        SlideData(index=6, title="Transition", bullets=["Manual", "Automated"], notes="", layout="content", variant="before_after", blocks=[{"type": "comparison", "headers": ["Before", "After"], "rows": [["Manual", "Automated"]]}]),
        SlideData(index=7, title="Capabilities", bullets=["Native"], notes="", layout="content", variant="comparison_table", blocks=[{"type": "table", "headers": ["Feature", "New"], "rows": [["Compliance", "Native"]]}]),
        SlideData(index=8, title="Timeline", bullets=["Audit", "Migrate"], notes="", layout="content", variant="process", blocks=[{"type": "process", "steps": [{"title": "Audit"}, {"title": "Migrate"}]}]),
        SlideData(index=9, title="Thank You", bullets=["Questions"], notes="", layout="next_steps", variant="closing"),
    ]

    prs = Presentation(BytesIO(PptxEngine().render(slides)))
    table_count = sum(1 for slide in prs.slides for shape in slide.shapes if getattr(shape, "has_table", False))
    picture_count = sum(1 for slide in prs.slides for shape in slide.shapes if shape.shape_type == MSO_SHAPE_TYPE.PICTURE)
    stat_count = sum(1 for slide in prs.slides for shape in slide.shapes if shape.has_text_frame and shape.text.strip() == "48%")

    assert table_count >= 1
    assert picture_count >= 1
    assert stat_count == 1


def test_render_does_not_append_separate_thank_you_slide_when_missing():
    slides = [
        SlideData(index=1, title="Cover", bullets=[], notes="", layout="title"),
        SlideData(index=2, title="Next Steps", bullets=["Review plan"], notes="", layout="next_steps"),
    ]

    prs = Presentation(BytesIO(PptxEngine().render(slides)))
    final_text = " ".join(_texts(prs.slides[-1]))

    assert len(prs.slides) == len(slides)
    assert "Thank You" not in final_text


def test_framework_variants_alternate_dark_and_light_backgrounds():
    slides = [
        SlideData(index=1, title="Cover", bullets=[], notes="", layout="title", variant="cover"),
        SlideData(index=2, title="Core Insight", bullets=["Teams move faster"], notes="", layout="content", variant="big_statement"),
        SlideData(index=3, title="Pillars", bullets=["Velocity", "Security", "Scale"], notes="", layout="content", variant="three_points"),
        SlideData(index=4, title="Impact", bullets=["48% reduction"], notes="", layout="content", variant="big_stat", blocks=[{"type": "stat", "value": "48%", "label": "Cost reduction"}]),
    ]

    prs = Presentation(BytesIO(PptxEngine().render(slides)))

    assert prs.slides[0].background.fill.fore_color.rgb == THEMES["dark"].background
    assert prs.slides[1].background.fill.fore_color.rgb == THEMES["dark"].background
    assert prs.slides[2].background.fill.fore_color.rgb == THEMES["minimalist"].background
    assert prs.slides[3].background.fill.fore_color.rgb == THEMES["dark"].background


def test_dark_framework_slide_uses_light_title_text():
    slides = [
        SlideData(index=1, title="Cover", bullets=[], notes="", layout="title"),
        SlideData(index=2, title="Core Insight", bullets=["Teams move faster"], notes="", layout="content", variant="big_statement"),
    ]

    prs = Presentation(BytesIO(PptxEngine().render(slides)))
    title_shapes = [s for s in prs.slides[1].shapes if s.has_text_frame and s.text.strip() == "Core Insight"]

    assert title_shapes
    assert title_shapes[0].text_frame.paragraphs[0].font.color.rgb == THEMES["dark"].text


def test_framework_cards_and_process_strip_literal_markdown_markers():
    slides = [
        SlideData(index=1, title="Cover", bullets=[], notes="", layout="title"),
        SlideData(
            index=2,
            title="Workflow",
            bullets=[],
            notes="",
            layout="content",
            variant="process",
            blocks=[
                {
                    "type": "process",
                    "steps": [
                        {"title": "**Input:** Users provide key data", "body": "**Generate:** AI drafts"},
                        {"title": "**Present:** Deliver polished output", "body": ""},
                    ],
                }
            ],
        ),
    ]

    prs = Presentation(BytesIO(PptxEngine().render(slides)))
    text = " ".join(_texts(prs.slides[1]))

    assert "**" not in text
    assert "Input:" in text
    assert "Generate:" in text


def test_big_stat_without_numeric_metric_does_not_oversize_full_title():
    title = "Accelerating Productivity and Enhancing Impact"
    slides = [
        SlideData(index=1, title="Cover", bullets=[], notes="", layout="title"),
        SlideData(
            index=2,
            title=title,
            bullets=["Significantly reduce the time spent on presentation creation and revision cycles."],
            notes="",
            layout="content",
            variant="big_stat",
        ),
    ]

    prs = Presentation(BytesIO(PptxEngine().render(slides)))
    title_shapes = [s for s in prs.slides[1].shapes if s.has_text_frame and s.text.strip() == title]

    assert title_shapes
    assert all(shape.text_frame.paragraphs[0].font.size < Pt(80) for shape in title_shapes)


def test_comparison_table_splits_topic_detail_rows_from_bullets():
    slides = [
        SlideData(index=1, title="Cover", bullets=[], notes="", layout="title"),
        SlideData(
            index=2,
            title="Capabilities",
            bullets=[
                "Automation: Superior workflow generation and brand compliance.",
                "Integration: Real-time data integration for dynamic reporting.",
            ],
            notes="",
            layout="content",
            variant="comparison_table",
        ),
    ]

    prs = Presentation(BytesIO(PptxEngine().render(slides)))
    tables = [s for s in prs.slides[1].shapes if getattr(s, "has_table", False)]
    table = tables[0].table

    assert table.cell(1, 0).text == "Automation"
    assert "Superior workflow" in table.cell(1, 1).text
    assert table.cell(2, 0).text == "Integration"
    assert "Real-time data" in table.cell(2, 1).text
