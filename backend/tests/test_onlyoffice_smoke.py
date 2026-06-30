import hashlib
import json
import os
import threading
from contextlib import contextmanager
from http.server import BaseHTTPRequestHandler, ThreadingHTTPServer
from io import BytesIO
from urllib.parse import urlsplit
from uuid import uuid4

import httpx
import jwt
import pytest
from httpx import ASGITransport, AsyncClient
from pptx import Presentation

from app import dependencies
from app.main import app
from app.services.platform.database import Database
from app.services.platform.deck_files import LocalDeckFileStorage
from app.services.platform.deck_repository import DeckRepository
from app.services.platform.deck_versions import DeckVersionService
from app.services.platform.onlyoffice import OnlyOfficeService


pytestmark = pytest.mark.onlyoffice


def _pptx_bytes(slides: int) -> bytes:
    presentation = Presentation()
    for _ in range(slides):
        presentation.slides.add_slide(presentation.slide_layouts[0])
    output = BytesIO()
    presentation.save(output)
    return output.getvalue()


async def _convert(
    client: httpx.AsyncClient,
    *,
    docs_url: str,
    secret: str,
    source_url: str,
    source_type: str,
    output_type: str,
) -> str:
    conversion_key = uuid4().hex
    payload = {
        "async": False,
        "filetype": source_type,
        "key": conversion_key,
        "outputtype": output_type,
        "title": f"onlyoffice-smoke.{source_type}",
        "url": source_url,
    }
    token = jwt.encode(payload, secret, algorithm="HS256")
    response = await client.post(
        f"{docs_url}/converter?shardkey={conversion_key}",
        json={**payload, "token": token},
        headers={
            "Accept": "application/json",
            "Authorization": f"Bearer {token}",
        },
    )
    response.raise_for_status()
    result = response.json()
    assert result.get("endConvert") is True, result
    converted_url = result.get("fileUrl")
    assert isinstance(converted_url, str) and converted_url.startswith(
        ("http://", "https://")
    )
    return converted_url


@contextmanager
def _fixture_server(content: bytes):
    class Handler(BaseHTTPRequestHandler):
        def do_GET(self):
            if self.path != "/edited.pptx":
                self.send_error(404)
                return
            self.send_response(200)
            self.send_header("Content-Type", "application/vnd.openxmlformats-officedocument.presentationml.presentation")
            self.send_header("Content-Length", str(len(content)))
            self.end_headers()
            self.wfile.write(content)

        def log_message(self, format, *args):
            return

    server = ThreadingHTTPServer(("0.0.0.0", 0), Handler)
    thread = threading.Thread(target=server.serve_forever, daemon=True)
    started = False
    try:
        thread.start()
        started = True
        fixture_host = os.getenv(
            "ONLYOFFICE_SMOKE_FIXTURE_HOST", "host.docker.internal"
        )
        yield (
            f"http://127.0.0.1:{server.server_port}",
            f"http://{fixture_host}:{server.server_port}",
        )
    finally:
        if started:
            server.shutdown()
        server.server_close()
        if started:
            thread.join(timeout=5)


@pytest.mark.skipif(not os.getenv("ONLYOFFICE_SMOKE_URL"), reason="ONLYOFFICE_SMOKE_URL is not set")
async def test_real_onlyoffice_docs_api_and_force_save_callback(tmp_path):
    docs_url = os.environ["ONLYOFFICE_SMOKE_URL"].rstrip("/")
    secret = os.getenv("ONLYOFFICE_SMOKE_JWT_SECRET") or os.getenv(
        "ONLYOFFICE_JWT_SECRET"
    )
    if not secret:
        pytest.fail(
            "ONLYOFFICE_SMOKE_JWT_SECRET or ONLYOFFICE_JWT_SECRET is required"
        )
    database = Database(f"sqlite+aiosqlite:///{tmp_path / 'smoke.db'}")
    repository = DeckRepository(database)
    storage = LocalDeckFileStorage(tmp_path / "files")
    deck_id, version_id = str(uuid4()), str(uuid4())
    created = False
    download_client: httpx.AsyncClient | None = None
    try:
        await database.create_schema()
        initial = _pptx_bytes(1)
        key = f"decks/{deck_id}/versions/{version_id}.pptx"
        await storage.put(key, initial)
        await repository.create_with_initial_version(
            deck_id=deck_id,
            version_id=version_id,
            owner_id="local-user",
            name="ONLYOFFICE smoke",
            deck_type="sales",
            theme="minimalist",
            aspect_ratio="16:9",
            generation_payload={"slides": []},
            storage_key=key,
            sha256=hashlib.sha256(initial).hexdigest(),
            size_bytes=len(initial),
        )
        created = True
        source = _pptx_bytes(2)
        with _fixture_server(source) as (_local_fixture, document_fixture):
            async with httpx.AsyncClient(timeout=60) as probe:
                api_response = await probe.get(
                    f"{docs_url}/web-apps/apps/api/documents/api.js"
                )
                api_response.raise_for_status()
                assert b"DocsAPI" in api_response.content
                # A two-hop conversion avoids a no-op native-format request:
                # ONLYOFFICE must parse the fixture PPTX, emit ODP, then parse
                # that result and emit the PPTX used by the callback.
                odp_url = await _convert(
                    probe,
                    docs_url=docs_url,
                    secret=secret,
                    source_url=f"{document_fixture}/edited.pptx",
                    source_type="pptx",
                    output_type="odp",
                )
                converted_url = await _convert(
                    probe,
                    docs_url=docs_url,
                    secret=secret,
                    source_url=odp_url,
                    source_type="odp",
                    output_type="pptx",
                )
            converted_origin = urlsplit(converted_url)
            download_client = httpx.AsyncClient()
            onlyoffice = OnlyOfficeService(
                public_url=docs_url,
                api_base_url="http://test",
                internal_url=f"{converted_origin.scheme}://{converted_origin.netloc}",
                jwt_secret=secret,
                file_token_ttl_seconds=300,
                max_file_bytes=len(source) + 1_000_000,
                enabled=True,
                download_client=download_client,
            )
            versions = DeckVersionService(
                repository, storage, None, len(source) + 1_000_000, 5
            )
            editor_config = onlyoffice.build_editor_config(
                deck=await repository.get(deck_id, "local-user"),
                user_id="local-user",
                user_name="Local User",
            )
            assert editor_config.document_server_url == docs_url
            assert editor_config.config["documentType"] == "slide"
            app.dependency_overrides[dependencies.get_onlyoffice_service] = (
                lambda: onlyoffice
            )
            app.dependency_overrides[dependencies.get_deck_version_service] = (
                lambda: versions
            )
            app.dependency_overrides[dependencies.get_deck_repository] = (
                lambda: repository
            )
            app.dependency_overrides[dependencies.get_deck_file_storage] = (
                lambda: storage
            )
            body = {
                "key": f"{deck_id}-{version_id}",
                "status": 6,
                "url": converted_url,
                "userdata": "smoke-force-save",
            }
            authorization = jwt.encode(
                {"payload": body}, secret, algorithm="HS256"
            )
            callback_token = onlyoffice.create_scoped_token(
                subject="local-user",
                deck_id=deck_id,
                version_id=version_id,
                purpose="callback",
            )
            async with AsyncClient(
                transport=ASGITransport(app=app, raise_app_exceptions=False),
                base_url="http://test",
            ) as client:
                callback = await client.post(
                    f"/api/v1/decks/{deck_id}/callback?token={callback_token}",
                    content=json.dumps(body),
                    headers={
                        "Authorization": f"Bearer {authorization}",
                        "Content-Type": "application/json",
                    },
                )
                downloaded = await client.get(f"/api/v1/decks/{deck_id}/download")
            assert callback.status_code == 200 and callback.json() == {"error": 0}
            assert downloaded.status_code == 200
            deck = await repository.get(deck_id, "local-user")
            assert deck is not None and deck.current_version is not None
            assert deck.current_version.version_number == 2
            assert len(Presentation(BytesIO(downloaded.content)).slides) == 2
    finally:
        app.dependency_overrides.clear()
        try:
            if download_client is not None:
                await download_client.aclose()
        finally:
            try:
                if created:
                    keys = await repository.delete(deck_id, "local-user")
                    for storage_key in keys:
                        await storage.delete(storage_key)
            finally:
                await database.dispose()
