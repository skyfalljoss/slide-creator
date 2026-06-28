import asyncio
import hashlib
import threading
from datetime import datetime, timezone
from io import BytesIO
from unittest.mock import AsyncMock
from uuid import NAMESPACE_URL, uuid5
import pytest

from app.models.schemas import SlideData
from app.services.platform.deck_repository import DeckRecord, DeckVersionRecord
from app.services.platform.deck_versions import DeckVersionService, VersionStorageConflictError
from app.services.presentation.pptx_validation import InvalidPptxError, validate_pptx


NOW = datetime.now(timezone.utc)


def slide(title: str = "A title") -> SlideData:
    return SlideData(
        index=1,
        title=title,
        bullets=["A point"],
        notes="Notes",
        layout="title",
    )


def pptx_bytes(title: str = "A title") -> bytes:
    from pptx import Presentation

    presentation = Presentation()
    page = presentation.slides.add_slide(presentation.slide_layouts[6])
    page.shapes.title if page.shapes.title else None
    page.notes_slide.notes_text_frame.text = title
    output = BytesIO()
    presentation.save(output)
    return output.getvalue()


def version_record(
    *,
    deck_id: str = "deck-1",
    version_id: str = "version-1",
    version_number: int = 1,
    storage_key: str | None = None,
    sha256: str = "checksum",
    size_bytes: int = 123,
    source: str = "generated",
) -> DeckVersionRecord:
    return DeckVersionRecord(
        id=version_id,
        deck_id=deck_id,
        version_number=version_number,
        storage_key=storage_key or f"decks/{deck_id}/versions/{version_id}.pptx",
        sha256=sha256,
        size_bytes=size_bytes,
        source=source,
        created_by="owner-1",
        created_at=NOW,
    )


def deck_record(version: DeckVersionRecord) -> DeckRecord:
    return DeckRecord(
        id=version.deck_id,
        owner_id="owner-1",
        name="Deck",
        deck_type="pitch",
        theme="citi",
        aspect_ratio="16:9",
        generation_payload={"slides": [slide().model_dump(mode="json")]},
        current_version_id=version.id,
        created_at=NOW,
        updated_at=NOW,
        current_version=version,
    )


class FakeStorage:
    def __init__(self) -> None:
        self.objects: dict[str, bytes] = {}
        self.events: list[str] = []
        self.delete_error: Exception | None = None
        self.delete_errors: set[str] = set()

    async def put(self, key: str, content: bytes) -> None:
        self.events.append("upload")
        if key in self.objects:
            raise FileExistsError(key)
        self.objects[key] = content

    async def read(self, key: str) -> bytes:
        return self.objects[key]

    async def delete(self, key: str) -> None:
        self.events.append("delete")
        if self.delete_error or key in self.delete_errors:
            raise self.delete_error or OSError("delete failed")
        self.objects.pop(key, None)

    async def exists(self, key: str) -> bool:
        return key in self.objects

    async def list_keys(self, prefix: str) -> list[str]:
        return sorted(key for key in self.objects if key.startswith(prefix))


class BlockingStorage(FakeStorage):
    def __init__(self) -> None:
        super().__init__()
        self.entered: list[asyncio.Event] = [asyncio.Event(), asyncio.Event()]
        self.release = asyncio.Event()

    async def put(self, key: str, content: bytes) -> None:
        index = len(self.events)
        self.events.append("upload")
        self.entered[index].set()
        await self.release.wait()
        if key in self.objects:
            raise FileExistsError(key)
        self.objects[key] = content


class InitialRepository:
    def __init__(self, storage: FakeStorage) -> None:
        self.storage = storage
        self.kwargs: dict | None = None
        self.error: Exception | None = None
        self.committed: DeckRecord | None = None
        self.reconciliation_error: Exception | None = None

    async def create_with_initial_version(self, **kwargs) -> DeckRecord:
        assert self.storage.events == ["upload"]
        self.kwargs = kwargs
        version = version_record(
            deck_id=kwargs["deck_id"],
            version_id=kwargs["version_id"],
            storage_key=kwargs["storage_key"],
            sha256=kwargs["sha256"],
            size_bytes=kwargs["size_bytes"],
        )
        result = deck_record(version)
        if self.error:
            raise self.error
        return result

    async def get(self, deck_id: str, owner_id: str) -> DeckRecord | None:
        if self.reconciliation_error:
            raise self.reconciliation_error
        return self.committed


class VersionRepository:
    def __init__(self, current: DeckRecord | None = None) -> None:
        self.current = current
        self.versions: dict[str, DeckVersionRecord] = {}
        if current is not None and current.current_version is not None:
            self.versions[current.current_version.id] = current.current_version
        self.version_results: list[DeckVersionRecord | None] = []
        self.append_calls: list[dict] = []
        self.deleted_row_batches: list[list[str]] = []
        self.delete_rows_error: Exception | None = None
        self.append_error: Exception | None = None
        self.commit_before_append_error = False
        self.version_error: Exception | None = None

    async def get(self, deck_id: str, owner_id: str) -> DeckRecord | None:
        if self.current is None or self.current.id != deck_id or self.current.owner_id != owner_id:
            return None
        return self.current

    async def version(
        self, deck_id: str, version_id: str, owner_id: str
    ) -> DeckVersionRecord | None:
        if self.version_error and self.append_calls:
            raise self.version_error
        if self.version_results:
            return self.version_results.pop(0)
        result = self.versions.get(version_id)
        if result is None or result.deck_id != deck_id or owner_id != "owner-1":
            return None
        return result

    async def append_version(self, **kwargs) -> DeckVersionRecord:
        self.append_calls.append(kwargs)
        result = version_record(
            deck_id=kwargs["deck_id"],
            version_id=kwargs["version_id"],
            version_number=(self.current.current_version.version_number + 1),
            storage_key=kwargs["storage_key"],
            sha256=kwargs["sha256"],
            source=kwargs["source"],
        )
        self.versions[result.id] = result
        assert self.current is not None
        self.current = DeckRecord(
            **{
                **self.current.__dict__,
                "generation_payload": kwargs.get("generation_payload")
                or self.current.generation_payload,
                "current_version_id": result.id,
                "current_version": result,
            }
        )
        if self.append_error:
            if not self.commit_before_append_error:
                self.versions.pop(result.id, None)
            raise self.append_error
        return result

    async def stale_versions(self, deck_id: str, keep: int) -> list[DeckVersionRecord]:
        ordered = sorted(
            (version for version in self.versions.values() if version.deck_id == deck_id),
            key=lambda version: version.version_number,
            reverse=True,
        )
        current_id = self.current.current_version_id if self.current is not None else None
        return [version for version in ordered[keep:] if version.id != current_id]

    async def delete_version_rows(self, version_ids: list[str]) -> None:
        self.deleted_row_batches.append(version_ids)
        if self.delete_rows_error:
            raise self.delete_rows_error
        for version_id in version_ids:
            self.versions.pop(version_id, None)


@pytest.mark.asyncio
async def test_create_generated_deck_renders_uploads_then_creates_metadata():
    storage = FakeStorage()
    repository = InitialRepository(storage)
    service = DeckVersionService(
        repository=repository,
        storage=storage,
        sample_template_path=None,
        max_file_bytes=50_000_000,
        retention=5,
    )

    result = await service.create_generated_deck(
        owner_id="owner-1",
        name="Deck",
        deck_type="pitch",
        theme="minimalist",
        aspect_ratio="16:9",
        slides=[slide()],
    )

    assert result.current_version is not None
    assert result.current_version.version_number == 1
    assert repository.kwargs is not None
    assert repository.kwargs["generation_payload"] == {
        "slides": [slide().model_dump(mode="json")]
    }
    key = repository.kwargs["storage_key"]
    assert key == f"decks/{result.id}/versions/{result.current_version.id}.pptx"
    content = await storage.read(key)
    validate_pptx(content, max_bytes=50_000_000)
    assert repository.kwargs["size_bytes"] == len(content)
    assert result.current_version.size_bytes == len(content)
    assert repository.kwargs["sha256"] == hashlib.sha256(content).hexdigest()
    assert result.current_version.sha256 == repository.kwargs["sha256"]


@pytest.mark.asyncio
async def test_create_generated_deck_deletes_upload_when_metadata_fails():
    storage = FakeStorage()
    repository = InitialRepository(storage)
    repository.error = RuntimeError("metadata failed")
    service = DeckVersionService(repository, storage, None, 50_000_000, 5)

    with pytest.raises(RuntimeError, match="metadata failed"):
        await service.create_generated_deck(
            owner_id="owner-1",
            name="Deck",
            deck_type="pitch",
            theme="minimalist",
            aspect_ratio="16:9",
            slides=[slide()],
        )

    assert storage.objects == {}
    assert storage.events == ["upload", "delete"]


@pytest.mark.asyncio
async def test_create_generated_deck_preserves_metadata_error_when_cleanup_fails():
    storage = FakeStorage()
    storage.delete_error = OSError("cleanup failed")
    repository = InitialRepository(storage)
    repository.error = RuntimeError("metadata failed")
    service = DeckVersionService(repository, storage, None, 50_000_000, 5)

    with pytest.raises(RuntimeError, match="metadata failed"):
        await service.create_generated_deck(
            owner_id="owner-1",
            name="Deck",
            deck_type="pitch",
            theme="minimalist",
            aspect_ratio="16:9",
            slides=[slide()],
        )


@pytest.mark.asyncio
async def test_create_generated_deck_returns_commit_after_repository_raises():
    storage = FakeStorage()
    repository = InitialRepository(storage)
    repository.error = RuntimeError("commit acknowledgement failed")
    original_create = repository.create_with_initial_version

    async def commit_then_raise(**kwargs):
        repository.error = None
        committed = await original_create(**kwargs)
        repository.committed = committed
        repository.error = RuntimeError("commit acknowledgement failed")
        raise repository.error

    repository.create_with_initial_version = commit_then_raise
    service = DeckVersionService(repository, storage, None, 50_000_000, 5)

    result = await service.create_generated_deck(
        owner_id="owner-1",
        name="Deck",
        deck_type="pitch",
        theme="minimalist",
        aspect_ratio="16:9",
        slides=[slide()],
    )

    assert result == repository.committed
    assert len(storage.objects) == 1
    assert storage.events == ["upload"]


@pytest.mark.asyncio
async def test_create_generated_deck_preserves_object_when_reconciliation_fails():
    storage = FakeStorage()
    repository = InitialRepository(storage)
    original_error = RuntimeError("metadata failed")
    repository.error = original_error
    repository.reconciliation_error = OSError("database unavailable")
    service = DeckVersionService(repository, storage, None, 50_000_000, 5)

    with pytest.raises(RuntimeError, match="metadata failed") as raised:
        await service.create_generated_deck(
            owner_id="owner-1",
            name="Deck",
            deck_type="pitch",
            theme="minimalist",
            aspect_ratio="16:9",
            slides=[slide()],
        )

    assert raised.value is original_error
    assert len(storage.objects) == 1
    assert storage.events == ["upload"]


@pytest.mark.asyncio
async def test_render_runs_off_event_loop(monkeypatch):
    event_loop_thread = threading.get_ident()
    observed_threads: list[int] = []
    real_validate = validate_pptx
    real_sha256 = hashlib.sha256

    def tracking_validate(content: bytes, max_bytes: int) -> None:
        observed_threads.append(threading.get_ident())
        real_validate(content, max_bytes)

    def tracking_sha256(content: bytes):
        observed_threads.append(threading.get_ident())
        return real_sha256(content)

    monkeypatch.setattr(
        "app.services.platform.deck_versions.validate_pptx", tracking_validate
    )
    monkeypatch.setattr(
        "app.services.platform.deck_versions.hashlib.sha256", tracking_sha256
    )
    storage = FakeStorage()
    repository = InitialRepository(storage)
    service = DeckVersionService(repository, storage, None, 50_000_000, 5)

    await service.create_generated_deck(
        owner_id="owner-1",
        name="Deck",
        deck_type="pitch",
        theme="minimalist",
        aspect_ratio="16:9",
        slides=[slide()],
    )

    assert len(observed_threads) == 2
    assert all(thread_id != event_loop_thread for thread_id in observed_threads)


@pytest.mark.asyncio
async def test_save_edited_version_rejects_invalid_bytes_before_upload():
    current = version_record()
    repository = VersionRepository(deck_record(current))
    storage = FakeStorage()
    service = DeckVersionService(repository, storage, None, 50_000_000, 5)

    with pytest.raises(InvalidPptxError):
        await service.save_edited_version(
            deck_id="deck-1",
            owner_id="owner-1",
            content=b"not a pptx",
            base_version_id=current.id,
            callback_key="callback-1",
            created_by="editor@example.com",
        )

    assert storage.events == []
    assert repository.append_calls == []


@pytest.mark.asyncio
async def test_save_edited_version_appends_onlyoffice_version_and_is_idempotent():
    current = version_record()
    repository = VersionRepository(deck_record(current))
    storage = FakeStorage()
    service = DeckVersionService(repository, storage, None, 50_000_000, 5)
    content = pptx_bytes("edited")

    first = await service.save_edited_version(
        deck_id="deck-1",
        owner_id="owner-1",
        content=content,
        base_version_id=current.id,
        callback_key="callback-1",
        created_by="editor@example.com",
    )
    second = await service.save_edited_version(
        deck_id="deck-1",
        owner_id="owner-1",
        content=content,
        base_version_id=current.id,
        callback_key="callback-1",
        created_by="editor@example.com",
    )

    assert second == first
    assert len(storage.objects) == 1
    assert len(repository.append_calls) == 1
    call = repository.append_calls[0]
    assert call["source"] == "onlyoffice_save"
    assert call["base_version_id"] == current.id
    assert call["created_by"] == "editor@example.com"
    assert first.storage_key == f"decks/deck-1/versions/{first.id}.pptx"


@pytest.mark.asyncio
async def test_save_edited_version_same_callback_different_content_has_distinct_id():
    current = version_record()
    repository = VersionRepository(deck_record(current))
    storage = FakeStorage()
    service = DeckVersionService(repository, storage, None, 50_000_000, 5)

    first = await service.save_edited_version(
        deck_id="deck-1",
        owner_id="owner-1",
        content=pptx_bytes("first"),
        base_version_id=current.id,
        callback_key="callback-1",
        created_by="editor",
    )
    second = await service.save_edited_version(
        deck_id="deck-1",
        owner_id="owner-1",
        content=pptx_bytes("second"),
        base_version_id=current.id,
        callback_key="callback-1",
        created_by="editor",
    )

    assert first.id != second.id
    assert len(repository.append_calls) == 2


@pytest.mark.asyncio
async def test_save_edited_version_returns_commit_after_append_raises():
    current = version_record()
    repository = VersionRepository(deck_record(current))
    repository.append_error = RuntimeError("commit acknowledgement failed")
    repository.commit_before_append_error = True
    storage = FakeStorage()
    service = DeckVersionService(repository, storage, None, 50_000_000, 5)

    result = await service.save_edited_version(
        deck_id="deck-1",
        owner_id="owner-1",
        content=pptx_bytes("committed"),
        base_version_id=current.id,
        callback_key="callback-commit",
        created_by="editor",
    )

    assert repository.versions[result.id] == result
    assert result.storage_key in storage.objects
    assert "delete" not in storage.events


@pytest.mark.asyncio
async def test_save_edited_version_preserves_object_when_reconciliation_fails():
    current = version_record()
    repository = VersionRepository(deck_record(current))
    original_error = RuntimeError("append failed")
    repository.append_error = original_error
    repository.version_error = OSError("database unavailable")
    storage = FakeStorage()
    service = DeckVersionService(repository, storage, None, 50_000_000, 5)

    with pytest.raises(RuntimeError, match="append failed") as raised:
        await service.save_edited_version(
            deck_id="deck-1",
            owner_id="owner-1",
            content=pptx_bytes("ambiguous"),
            base_version_id=current.id,
            callback_key="callback-ambiguous",
            created_by="editor",
        )

    assert raised.value is original_error
    assert len(storage.objects) == 1
    assert "delete" not in storage.events


@pytest.mark.asyncio
async def test_save_edited_version_deletes_object_after_confirmed_rollback():
    current = version_record()
    repository = VersionRepository(deck_record(current))
    original_error = RuntimeError("append rolled back")
    repository.append_error = original_error
    storage = FakeStorage()
    service = DeckVersionService(repository, storage, None, 50_000_000, 5)

    with pytest.raises(RuntimeError, match="append rolled back"):
        await service.save_edited_version(
            deck_id="deck-1",
            owner_id="owner-1",
            content=pptx_bytes("rolled back"),
            base_version_id=current.id,
            callback_key="callback-rollback",
            created_by="editor",
        )

    assert storage.objects == {}
    assert storage.events[-1] == "delete"


@pytest.mark.asyncio
async def test_save_edited_version_file_exists_race_returns_committed_match():
    current = version_record()
    repository = VersionRepository(deck_record(current))
    storage = FakeStorage()
    service = DeckVersionService(repository, storage, None, 50_000_000, 5)
    content = pptx_bytes("race")
    checksum = hashlib.sha256(content).hexdigest()

    # Determine the stable ID without exposing the namespace as public API.
    initial = await service.save_edited_version(
        deck_id="deck-1",
        owner_id="owner-1",
        content=content,
        base_version_id=current.id,
        callback_key="callback-1",
        created_by="editor",
    )
    repository.versions.clear()
    repository.append_calls.clear()
    repository.version_results = [None, None, initial]
    storage.events.clear()

    result = await service.save_edited_version(
        deck_id="deck-1",
        owner_id="owner-1",
        content=content,
        base_version_id=current.id,
        callback_key="callback-1",
        created_by="editor",
    )

    assert result == initial
    assert result.sha256 == checksum
    assert repository.append_calls == []


@pytest.mark.asyncio
async def test_save_edited_version_repairs_matching_upload_orphan():
    current = version_record()
    repository = VersionRepository(deck_record(current))
    storage = FakeStorage()
    service = DeckVersionService(repository, storage, None, 50_000_000, 5)
    content = pptx_bytes("orphan")
    first = await service.save_edited_version(
        deck_id="deck-1",
        owner_id="owner-1",
        content=content,
        base_version_id=current.id,
        callback_key="callback-1",
        created_by="editor",
    )
    repository.versions.pop(first.id)
    repository.append_calls.clear()

    repaired = await service.save_edited_version(
        deck_id="deck-1",
        owner_id="owner-1",
        content=content,
        base_version_id=current.id,
        callback_key="callback-1",
        created_by="editor",
    )

    assert repaired.id == first.id
    assert len(repository.append_calls) == 1
    assert storage.events.count("upload") == 2


@pytest.mark.asyncio
async def test_save_edited_version_rejects_mismatched_upload_orphan():
    current = version_record()
    repository = VersionRepository(deck_record(current))
    storage = FakeStorage()
    service = DeckVersionService(repository, storage, None, 50_000_000, 5)
    requested = pptx_bytes("requested")
    checksum = hashlib.sha256(requested).hexdigest()
    version_id = str(
        uuid5(NAMESPACE_URL, f"slideforge:deck-1:callback-1:{checksum}")
    )
    key = f"decks/deck-1/versions/{version_id}.pptx"
    storage.objects[key] = pptx_bytes("different")

    with pytest.raises(VersionStorageConflictError, match="checksum"):
        await service.save_edited_version(
            deck_id="deck-1",
            owner_id="owner-1",
            content=requested,
            base_version_id=current.id,
            callback_key="callback-1",
            created_by="editor",
        )

    assert repository.append_calls == []
    assert key in storage.objects


@pytest.mark.asyncio
async def test_two_service_instances_observe_delayed_winner_commit(monkeypatch):
    current = version_record()
    repository = VersionRepository(deck_record(current))
    storage = FakeStorage()
    winner_entered = asyncio.Event()
    allow_winner = asyncio.Event()
    winner_committed = asyncio.Event()
    original_append = repository.append_version

    async def delayed_append(**kwargs):
        if not winner_entered.is_set():
            winner_entered.set()
            await allow_winner.wait()
            result = await original_append(**kwargs)
            winner_committed.set()
            return result
        return await original_append(**kwargs)

    repository.append_version = delayed_append
    first_service = DeckVersionService(repository, storage, None, 50_000_000, 5)
    second_service = DeckVersionService(repository, storage, None, 50_000_000, 5)
    content = pptx_bytes("race")

    first_task = asyncio.create_task(
        first_service.save_edited_version(
            deck_id="deck-1",
            owner_id="owner-1",
            content=content,
            base_version_id=current.id,
            callback_key="callback-race",
            created_by="editor",
        )
    )
    await winner_entered.wait()

    async def release_winner(_delay: float) -> None:
        allow_winner.set()
        await winner_committed.wait()

    second_service._sleep = release_winner
    second_task = asyncio.create_task(
        second_service.save_edited_version(
            deck_id="deck-1",
            owner_id="owner-1",
            content=content,
            base_version_id=current.id,
            callback_key="callback-race",
            created_by="editor",
        )
    )

    first, second = await asyncio.gather(first_task, second_task)
    assert first == second
    assert len(repository.append_calls) == 1


@pytest.mark.asyncio
async def test_file_exists_queries_after_final_retry_delay():
    current = version_record()
    repository = VersionRepository(deck_record(current))
    storage = FakeStorage()
    service = DeckVersionService(repository, storage, None, 50_000_000, 5)
    service._retry_delays = (0, 0.01)
    content = pptx_bytes("final observation")
    checksum = hashlib.sha256(content).hexdigest()
    version_id = str(
        uuid5(NAMESPACE_URL, f"slideforge:deck-1:callback-final:{checksum}")
    )
    key = f"decks/deck-1/versions/{version_id}.pptx"
    storage.objects[key] = content
    committed = version_record(
        version_id=version_id,
        version_number=2,
        storage_key=key,
        sha256=checksum,
        source="onlyoffice_save",
    )

    async def commit_during_final_delay(delay: float) -> None:
        if delay == 0.01:
            repository.versions[version_id] = committed

    service._sleep = commit_during_final_delay

    result = await service.save_edited_version(
        deck_id="deck-1",
        owner_id="owner-1",
        content=content,
        base_version_id=current.id,
        callback_key="callback-final",
        created_by="editor",
    )

    assert result == committed
    assert repository.append_calls == []


@pytest.mark.asyncio
async def test_different_editor_identities_do_not_share_save_lock(monkeypatch):
    current = version_record()
    repository = VersionRepository(deck_record(current))
    storage = BlockingStorage()
    service = DeckVersionService(repository, storage, None, 50_000_000, 5)
    content = pptx_bytes("concurrent")
    checksum = hashlib.sha256(content).hexdigest()
    monkeypatch.setattr(
        service, "_validate_and_hash", AsyncMock(return_value=checksum), raising=False
    )

    first = asyncio.create_task(
        service.save_edited_version(
            deck_id="deck-1",
            owner_id="owner-1",
            content=content,
            base_version_id=current.id,
            callback_key="callback-1",
            created_by="editor",
        )
    )
    await storage.entered[0].wait()
    second = asyncio.create_task(
        service.save_edited_version(
            deck_id="deck-1",
            owner_id="owner-1",
            content=content,
            base_version_id=current.id,
            callback_key="callback-2",
            created_by="editor",
        )
    )

    await asyncio.wait_for(storage.entered[1].wait(), timeout=1)
    storage.release.set()
    await asyncio.gather(first, second)
    assert storage.events == ["upload", "upload"]


@pytest.mark.asyncio
async def test_exact_duplicate_saves_serialize_and_release_keyed_lock(monkeypatch):
    current = version_record()
    repository = VersionRepository(deck_record(current))
    storage = BlockingStorage()
    service = DeckVersionService(repository, storage, None, 50_000_000, 5)
    content = pptx_bytes("duplicate")
    checksum = hashlib.sha256(content).hexdigest()
    monkeypatch.setattr(service, "_validate_and_hash", AsyncMock(return_value=checksum))
    kwargs = {
        "deck_id": "deck-1",
        "owner_id": "owner-1",
        "content": content,
        "base_version_id": current.id,
        "callback_key": "callback-duplicate",
        "created_by": "editor",
    }

    first = asyncio.create_task(service.save_edited_version(**kwargs))
    await storage.entered[0].wait()
    second = asyncio.create_task(service.save_edited_version(**kwargs))
    for _ in range(10):
        entry = next(iter(service._lock_entries.values()))
        if entry.users == 2:
            break
        await asyncio.sleep(0)

    assert entry.users == 2
    assert not storage.entered[1].is_set()
    storage.release.set()
    first_result, second_result = await asyncio.gather(first, second)
    assert first_result == second_result
    assert storage.events == ["upload"]
    assert service._lock_entries == {}


@pytest.mark.asyncio
async def test_cancelled_duplicate_waiter_is_removed_from_keyed_lock(monkeypatch):
    current = version_record()
    repository = VersionRepository(deck_record(current))
    storage = BlockingStorage()
    service = DeckVersionService(repository, storage, None, 50_000_000, 5)
    content = pptx_bytes("cancelled duplicate")
    checksum = hashlib.sha256(content).hexdigest()
    monkeypatch.setattr(service, "_validate_and_hash", AsyncMock(return_value=checksum))
    kwargs = {
        "deck_id": "deck-1",
        "owner_id": "owner-1",
        "content": content,
        "base_version_id": current.id,
        "callback_key": "callback-cancel",
        "created_by": "editor",
    }

    first = asyncio.create_task(service.save_edited_version(**kwargs))
    await storage.entered[0].wait()
    waiter = asyncio.create_task(service.save_edited_version(**kwargs))
    for _ in range(10):
        entry = next(iter(service._lock_entries.values()))
        if entry.users == 2:
            break
        await asyncio.sleep(0)
    waiter.cancel()
    with pytest.raises(asyncio.CancelledError):
        await waiter

    assert entry.users == 1
    storage.release.set()
    await first
    assert service._lock_entries == {}


@pytest.mark.asyncio
async def test_restore_version_copies_owned_source_to_new_current_version():
    current = version_record(version_id="current", version_number=2)
    selected = version_record(version_id="selected", version_number=1)
    repository = VersionRepository(deck_record(current))
    repository.versions[selected.id] = selected
    storage = FakeStorage()
    source_content = pptx_bytes("restore")
    storage.objects[selected.storage_key] = source_content
    service = DeckVersionService(repository, storage, None, 50_000_000, 5)

    restored = await service.restore_version(
        deck_id="deck-1",
        version_id="selected",
        owner_id="owner-1",
        created_by="restorer",
    )

    assert restored.id != selected.id
    assert storage.objects[restored.storage_key] == source_content
    call = repository.append_calls[0]
    assert call["source"] == "restore"
    assert call["base_version_id"] == current.id
    assert call["created_by"] == "restorer"


@pytest.mark.asyncio
@pytest.mark.parametrize("owner_id,version_id", [("other-owner", "selected"), ("owner-1", "missing")])
async def test_restore_version_rejects_missing_or_cross_owner(owner_id: str, version_id: str):
    current = version_record(version_id="current")
    selected = version_record(version_id="selected")
    repository = VersionRepository(deck_record(current))
    repository.versions[selected.id] = selected
    storage = FakeStorage()
    service = DeckVersionService(repository, storage, None, 50_000_000, 5)

    with pytest.raises(LookupError, match="not found"):
        await service.restore_version(
            deck_id="deck-1",
            version_id=version_id,
            owner_id=owner_id,
            created_by="restorer",
        )

    assert storage.events == []


@pytest.mark.asyncio
async def test_save_slides_as_version_renders_and_updates_generation_provenance():
    current = version_record(version_id="current")
    repository = VersionRepository(deck_record(current))
    storage = FakeStorage()
    service = DeckVersionService(repository, storage, None, 50_000_000, 5)
    slides = [slide("New generated slide")]

    result = await service.save_slides_as_version(
        deck_id="deck-1",
        owner_id="owner-1",
        slides=slides,
        theme="minimalist",
        aspect_ratio="16:9",
        created_by="generator",
    )

    validate_pptx(storage.objects[result.storage_key], max_bytes=50_000_000)
    call = repository.append_calls[0]
    assert call["source"] == "generated"
    assert call["base_version_id"] == current.id
    assert call["generation_payload"] == {
        "slides": [item.model_dump(mode="json") for item in slides]
    }


async def save_editor_revision(
    service: DeckVersionService, number: int
) -> DeckVersionRecord:
    return await service.save_edited_version(
        deck_id="deck-1",
        owner_id="owner-1",
        content=pptx_bytes(f"revision-{number}"),
        base_version_id=f"base-{number}",
        callback_key=f"callback-{number}",
        created_by="editor",
    )


@pytest.mark.asyncio
async def test_retention_keeps_newest_five_after_sixth_version():
    initial = version_record(version_id="initial")
    repository = VersionRepository(deck_record(initial))
    storage = FakeStorage()
    storage.objects[initial.storage_key] = pptx_bytes("initial")
    service = DeckVersionService(repository, storage, None, 50_000_000, 5)

    for number in range(1, 5):
        await save_editor_revision(service, number)

    assert len(repository.versions) == 5
    assert initial.id in repository.versions

    await save_editor_revision(service, 5)

    assert len(repository.versions) == 5
    assert len(storage.objects) == 5
    assert "initial" not in repository.versions
    assert initial.storage_key not in storage.objects
    assert repository.current is not None
    assert repository.current.current_version_id in repository.versions


@pytest.mark.asyncio
async def test_retention_leaves_row_when_object_delete_fails():
    initial = version_record(version_id="initial")
    repository = VersionRepository(deck_record(initial))
    storage = FakeStorage()
    storage.objects[initial.storage_key] = pptx_bytes("initial")
    storage.delete_errors.add(initial.storage_key)
    service = DeckVersionService(repository, storage, None, 50_000_000, 5)

    for number in range(1, 6):
        await save_editor_revision(service, number)

    assert initial.id in repository.versions
    assert initial.storage_key in storage.objects
    assert all(initial.id not in batch for batch in repository.deleted_row_batches)


@pytest.mark.asyncio
async def test_retention_row_cleanup_failure_does_not_fail_committed_save():
    initial = version_record(version_id="initial")
    repository = VersionRepository(deck_record(initial))
    repository.delete_rows_error = RuntimeError("row cleanup failed")
    storage = FakeStorage()
    storage.objects[initial.storage_key] = pptx_bytes("initial")
    service = DeckVersionService(repository, storage, None, 50_000_000, 5)

    for number in range(1, 6):
        result = await save_editor_revision(service, number)

    assert repository.current is not None
    assert repository.current.current_version_id == result.id
    assert initial.id in repository.versions
    assert initial.storage_key not in storage.objects
