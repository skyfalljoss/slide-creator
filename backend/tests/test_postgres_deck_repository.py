import asyncio
import hashlib
import os
from io import BytesIO
from uuid import uuid4

import pytest
from pptx import Presentation

from app.services.platform.database import Database
from app.services.platform.deck_files import LocalDeckFileStorage
from app.services.platform.deck_repository import DeckRepository
from app.services.platform.deck_versions import DeckVersionService


pytestmark = pytest.mark.postgres


def _pptx_bytes(slides: int = 1) -> bytes:
    presentation = Presentation()
    for _ in range(slides):
        presentation.slides.add_slide(presentation.slide_layouts[0])
    output = BytesIO()
    presentation.save(output)
    return output.getvalue()


@pytest.mark.skipif(not os.getenv("TEST_DATABASE_URL"), reason="TEST_DATABASE_URL is not set")
async def test_postgres_concurrent_versions_and_callback_idempotency(tmp_path):
    url = os.environ["TEST_DATABASE_URL"]
    if not url.startswith("postgresql+asyncpg://"):
        pytest.fail("TEST_DATABASE_URL must use postgresql+asyncpg")
    first_database = Database(url)
    second_database = Database(url)
    first = DeckRepository(first_database)
    second = DeckRepository(second_database)
    deck_id = str(uuid4())
    initial_id = str(uuid4())
    prefix = f"decks/{deck_id}/versions/"
    storage = LocalDeckFileStorage(tmp_path / "files")
    created = False
    try:
        await first_database.create_schema()
        second_entered = asyncio.Event()

        async def enter_second_guard():
            async with second.storage_key_guard("decks/advisory-lock-test.pptx"):
                second_entered.set()

        async with first.storage_key_guard("decks/advisory-lock-test.pptx"):
            guard_waiter = asyncio.create_task(enter_second_guard())
            await asyncio.sleep(0.05)
            assert second_entered.is_set() is False
        await asyncio.wait_for(guard_waiter, timeout=2)
        assert second_entered.is_set() is True

        initial = _pptx_bytes()
        initial_key = f"{prefix}{initial_id}.pptx"
        await storage.put(initial_key, initial)
        await first.create_with_initial_version(
            deck_id=deck_id,
            version_id=initial_id,
            owner_id="postgres-test-owner",
            name="Concurrency",
            deck_type="sales",
            theme="minimalist",
            aspect_ratio="16:9",
            generation_payload={"slides": []},
            storage_key=initial_key,
            sha256=hashlib.sha256(initial).hexdigest(),
            size_bytes=len(initial),
        )
        created = True
        ids = [str(uuid4()), str(uuid4())]
        concurrent_contents = [_pptx_bytes(2), _pptx_bytes(3)]
        concurrent_keys = [f"{prefix}{ids[0]}.pptx", f"{prefix}{ids[1]}.pptx"]
        await storage.put(concurrent_keys[0], concurrent_contents[0])
        await storage.put(concurrent_keys[1], concurrent_contents[1])
        records = await asyncio.gather(
            first.append_version(
                deck_id=deck_id,
                owner_id="postgres-test-owner",
                version_id=ids[0],
                storage_key=concurrent_keys[0],
                sha256=hashlib.sha256(concurrent_contents[0]).hexdigest(),
                size_bytes=len(concurrent_contents[0]),
                source="test",
                created_by="postgres-test-owner",
            ),
            second.append_version(
                deck_id=deck_id,
                owner_id="postgres-test-owner",
                version_id=ids[1],
                storage_key=concurrent_keys[1],
                sha256=hashlib.sha256(concurrent_contents[1]).hexdigest(),
                size_bytes=len(concurrent_contents[1]),
                source="test",
                created_by="postgres-test-owner",
            ),
        )
        assert {record.version_number for record in records} == {2, 3}
        current = await first.get(deck_id, "postgres-test-owner")
        assert current.current_version_id in ids
        assert await storage.exists(current.current_version.storage_key)

        edited = _pptx_bytes(2)
        service_one = DeckVersionService(first, storage, None, len(edited) + 100, 10)
        service_two = DeckVersionService(second, storage, None, len(edited) + 100, 10)
        base_id = (await first.get(deck_id, "postgres-test-owner")).current_version_id
        repeated = await asyncio.gather(
            service_one.save_edited_version(
                deck_id=deck_id,
                owner_id="postgres-test-owner",
                content=edited,
                base_version_id=base_id,
                callback_key="same-callback",
                created_by="postgres-test-owner",
            ),
            service_two.save_edited_version(
                deck_id=deck_id,
                owner_id="postgres-test-owner",
                content=edited,
                base_version_id=base_id,
                callback_key="same-callback",
                created_by="postgres-test-owner",
            ),
        )
        assert repeated[0].id == repeated[1].id
        assert sum(
            version.id == repeated[0].id
            for version in await first.list_versions(deck_id, "postgres-test-owner")
        ) == 1
    finally:
        try:
            if created:
                keys = await first.delete(deck_id, "postgres-test-owner")
                for key in keys:
                    await storage.delete(key)
        finally:
            try:
                await first_database.dispose()
            finally:
                await second_database.dispose()
