from datetime import datetime, timedelta, timezone
from contextlib import asynccontextmanager
import asyncio
import hashlib
from io import BytesIO
from uuid import NAMESPACE_URL, uuid4, uuid5
from argparse import Namespace

from pptx import Presentation
import pytest

from app.services.platform.deck_files import DeckFileObject, LocalDeckFileStorage
from app.services.platform.database import Database
from app.services.platform.deck_repository import DeckRepository
from app.services.platform.deck_versions import DeckVersionService
from scripts.cleanup_orphan_deck_files import cleanup_orphan_deck_files
import scripts.cleanup_orphan_deck_files as cleanup_script


class FakeRepository:
    def __init__(self, keys=None, error=None):
        self.keys = set(keys or ())
        self.error = error

    async def all_storage_keys(self):
        if self.error:
            raise self.error
        return self.keys

    async def storage_key_referenced(self, key, *, session=None):
        return key in self.keys

    @asynccontextmanager
    async def storage_key_guard(self, key):
        yield


class FakeStorage:
    def __init__(self, objects=None, *, list_error=None, delete_error=None):
        self.objects = list(objects or ())
        self.list_error = list_error
        self.delete_error = delete_error
        self.deleted = []

    async def list_objects(self, prefix):
        assert prefix == "decks/"
        if self.list_error:
            raise self.list_error
        return self.objects

    async def delete(self, key):
        if key == self.delete_error:
            raise OSError("delete failed")
        self.deleted.append(key)


async def test_cleanup_cli_uses_configured_lock_dir_for_cross_process_coordination(
    tmp_path, monkeypatch
):
    database_url = f"sqlite+aiosqlite:///{tmp_path / 'cli.db'}"
    lock_dir = tmp_path / "custom-locks"
    app_database = Database(database_url)
    app_repository = DeckRepository(app_database, lock_dir=lock_dir)

    class Storage:
        async def close(self):
            return None

    async def coordinated_cleanup(script_repository, _storage, *, apply):
        assert apply is False
        entered = asyncio.Event()

        async def waiter():
            async with script_repository.storage_key_guard("decks/shared.pptx"):
                entered.set()

        async with app_repository.storage_key_guard("decks/shared.pptx"):
            task = asyncio.create_task(waiter())
            await asyncio.sleep(0.05)
            assert entered.is_set() is False
        await asyncio.wait_for(task, 1)
        assert entered.is_set()
        return cleanup_script.CleanupResult()

    monkeypatch.setattr(cleanup_script.settings, "database_url", database_url)
    monkeypatch.setattr(cleanup_script.settings, "deck_lock_dir", str(lock_dir))
    monkeypatch.setattr(cleanup_script, "_storage", lambda: Storage())
    monkeypatch.setattr(
        cleanup_script, "cleanup_orphan_deck_files", coordinated_cleanup
    )
    try:
        assert await cleanup_script._main(Namespace(apply=False)) == 0
    finally:
        await app_database.dispose()


async def test_cleanup_defaults_to_dry_run_and_retains_exact_grace_boundary():
    now = datetime(2026, 1, 2, tzinfo=timezone.utc)
    storage = FakeStorage(
        [
            DeckFileObject("decks/d/v1.pptx", now - timedelta(days=10)),
            DeckFileObject("decks/d/v2.pptx", now - timedelta(days=2)),
            DeckFileObject("decks/d/orphan.pptx", now - timedelta(hours=24)),
            DeckFileObject("decks/d/fresh.pptx", now - timedelta(hours=23)),
        ]
    )
    repository = FakeRepository({"decks/d/v1.pptx", "decks/d/v2.pptx"})

    result = await cleanup_orphan_deck_files(repository, storage, now=now)

    assert (result.examined, result.retained, result.candidates, result.deleted) == (4, 4, 0, 0)
    assert result.failed == 0
    assert storage.deleted == []


async def test_cleanup_apply_deletes_candidates_and_reports_each_failure():
    now = datetime(2026, 1, 2, tzinfo=timezone.utc)
    storage = FakeStorage(
        [
            DeckFileObject("decks/d/a.pptx", now - timedelta(days=2)),
            DeckFileObject("decks/d/b.pptx", now - timedelta(days=2)),
        ],
        delete_error="decks/d/b.pptx",
    )

    result = await cleanup_orphan_deck_files(
        FakeRepository(), storage, apply=True, now=now
    )

    assert (result.candidates, result.deleted, result.failed) == (2, 1, 1)
    assert storage.deleted == ["decks/d/a.pptx"]
    assert result.failure_messages == (
        "storage delete failed (OSError): decks/d/b.pptx",
    )


async def test_cleanup_rechecks_candidate_that_becomes_referenced_before_delete():
    now = datetime(2026, 1, 2, tzinfo=timezone.utc)
    key = "decks/d/newly-referenced.pptx"
    storage = FakeStorage([DeckFileObject(key, now - timedelta(days=2))])

    class RacingRepository(FakeRepository):
        async def storage_key_referenced(self, candidate, *, session=None):
            self.keys.add(candidate)
            return True

    result = await cleanup_orphan_deck_files(
        RacingRepository(), storage, apply=True, now=now
    )

    assert result.deleted == 0
    assert result.retained == 1
    assert storage.deleted == []


async def test_cleanup_recheck_failure_is_fail_closed():
    now = datetime(2026, 1, 2, tzinfo=timezone.utc)
    key = "decks/d/uncertain.pptx"
    storage = FakeStorage([DeckFileObject(key, now - timedelta(days=2))])

    class UnavailableRepository(FakeRepository):
        async def storage_key_referenced(self, candidate, *, session=None):
            raise OSError("database unavailable")

    result = await cleanup_orphan_deck_files(
        UnavailableRepository(), storage, apply=True, now=now
    )

    assert result.failed == 1 and result.retained == 1
    assert result.candidates == 0
    assert result.failure_messages == (
        "storage guard/recheck failed (OSError): decks/d/uncertain.pptx",
    )
    assert storage.deleted == []


async def test_cleanup_never_deletes_when_listing_or_repository_is_uncertain():
    now = datetime.now(timezone.utc)
    old = [DeckFileObject("decks/d/orphan.pptx", now - timedelta(days=2))]
    metadata_failure = FakeStorage(old)
    metadata_failure.objects[0] = DeckFileObject("decks/d/orphan.pptx", None)
    repository_failure = FakeStorage(old)
    listing_failure = FakeStorage(old, list_error=OSError("listing incomplete"))

    missing_time = await cleanup_orphan_deck_files(
        FakeRepository(), metadata_failure, apply=True, now=now
    )
    missing_keys = await cleanup_orphan_deck_files(
        FakeRepository(error=OSError("database unavailable")),
        repository_failure,
        apply=True,
        now=now,
    )
    missing_list = await cleanup_orphan_deck_files(
        FakeRepository(), listing_failure, apply=True, now=now
    )

    assert missing_time.failed == 1 and metadata_failure.deleted == []
    assert missing_keys.failed == 1 and repository_failure.deleted == []
    assert missing_list.failed == 1 and listing_failure.deleted == []


def _pptx(slides: int) -> bytes:
    presentation = Presentation()
    for _ in range(slides):
        presentation.slides.add_slide(presentation.slide_layouts[0])
    output = BytesIO()
    presentation.save(output)
    return output.getvalue()


class CoordinatedStorage(LocalDeckFileStorage):
    def __init__(self, root):
        super().__init__(root)
        self.target = ""
        self.pause_put = False
        self.pause_delete = False
        self.put_entered = asyncio.Event()
        self.delete_entered = asyncio.Event()
        self.release = asyncio.Event()

    async def put(self, key, content):
        await super().put(key, content)
        if self.pause_put and key == self.target:
            self.put_entered.set()
            await self.release.wait()

    async def delete(self, key):
        if self.pause_delete and key == self.target:
            self.delete_entered.set()
            await self.release.wait()
        await super().delete(key)

    async def list_objects(self, prefix):
        objects = await super().list_objects(prefix)
        old = datetime.now(timezone.utc) - timedelta(days=2)
        return [
            DeckFileObject(item.key, old if item.key == self.target else item.updated_at)
            for item in objects
        ]


async def _version_race_fixture(tmp_path):
    url = f"sqlite+aiosqlite:///{tmp_path / 'race.db'}"
    first_database = Database(url)
    second_database = Database(url)
    await first_database.create_schema()
    lock_dir = tmp_path / "locks"
    first = DeckRepository(first_database, lock_dir=lock_dir)
    second = DeckRepository(second_database, lock_dir=lock_dir)
    storage = CoordinatedStorage(tmp_path / "files")
    deck_id, initial_id = str(uuid4()), str(uuid4())
    initial = _pptx(1)
    initial_key = f"decks/{deck_id}/versions/{initial_id}.pptx"
    await storage.put(initial_key, initial)
    await first.create_with_initial_version(
        deck_id=deck_id,
        version_id=initial_id,
        owner_id="owner",
        name="Race",
        deck_type="sales",
        theme="minimalist",
        aspect_ratio="16:9",
        generation_payload={"slides": []},
        storage_key=initial_key,
        sha256=hashlib.sha256(initial).hexdigest(),
        size_bytes=len(initial),
    )
    edited = _pptx(2)
    callback_key = "cleanup-race"
    checksum = hashlib.sha256(edited).hexdigest()
    version_id = str(
        uuid5(NAMESPACE_URL, f"slideforge:{deck_id}:{callback_key}:{checksum}")
    )
    target = f"decks/{deck_id}/versions/{version_id}.pptx"
    storage.target = target
    service = DeckVersionService(second, storage, None, len(edited) + 100, 5)
    return (
        first_database,
        second_database,
        first,
        service,
        storage,
        deck_id,
        initial_id,
        edited,
        callback_key,
        target,
    )


async def test_cleanup_guard_deletes_first_then_writer_uploads_and_commits(tmp_path):
    fixture = await _version_race_fixture(tmp_path)
    db1, db2, repository, service, storage, deck_id, base_id, edited, callback_key, target = fixture
    await storage.put(target, edited)
    storage.pause_delete = True
    try:
        cleanup = asyncio.create_task(
            cleanup_orphan_deck_files(
                repository, storage, apply=True, now=datetime.now(timezone.utc)
            )
        )
        await asyncio.wait_for(storage.delete_entered.wait(), 1)
        writer = asyncio.create_task(
            service.save_edited_version(
                deck_id=deck_id,
                owner_id="owner",
                content=edited,
                base_version_id=base_id,
                callback_key=callback_key,
                created_by="owner",
            )
        )
        await asyncio.sleep(0.05)
        assert writer.done() is False
        storage.release.set()
        cleanup_result, version = await asyncio.gather(cleanup, writer)
        assert cleanup_result.deleted == 1
        assert version.storage_key == target
        assert await storage.exists(target)
    finally:
        await db1.dispose()
        await db2.dispose()


async def test_writer_guard_commits_before_cleanup_rechecks_and_retains(tmp_path):
    fixture = await _version_race_fixture(tmp_path)
    db1, db2, repository, service, storage, deck_id, base_id, edited, callback_key, target = fixture
    storage.pause_put = True
    try:
        writer = asyncio.create_task(
            service.save_edited_version(
                deck_id=deck_id,
                owner_id="owner",
                content=edited,
                base_version_id=base_id,
                callback_key=callback_key,
                created_by="owner",
            )
        )
        await asyncio.wait_for(storage.put_entered.wait(), 1)
        cleanup = asyncio.create_task(
            cleanup_orphan_deck_files(
                repository, storage, apply=True, now=datetime.now(timezone.utc)
            )
        )
        await asyncio.sleep(0.05)
        assert cleanup.done() is False
        storage.release.set()
        version, cleanup_result = await asyncio.gather(writer, cleanup)
        assert version.storage_key == target
        assert cleanup_result.deleted == 0
        assert await storage.exists(target)
    finally:
        await db1.dispose()
        await db2.dispose()


async def test_cancelled_cleanup_holds_guard_until_delete_finishes_then_writer_commits(tmp_path):
    fixture = await _version_race_fixture(tmp_path)
    db1, db2, repository, service, storage, deck_id, base_id, edited, callback_key, target = fixture
    await storage.put(target, edited)
    storage.pause_delete = True
    try:
        cleanup = asyncio.create_task(
            cleanup_orphan_deck_files(
                repository, storage, apply=True, now=datetime.now(timezone.utc)
            )
        )
        await asyncio.wait_for(storage.delete_entered.wait(), 1)
        writer = asyncio.create_task(
            service.save_edited_version(
                deck_id=deck_id,
                owner_id="owner",
                content=edited,
                base_version_id=base_id,
                callback_key=callback_key,
                created_by="owner",
            )
        )
        cleanup.cancel()
        await asyncio.sleep(0.05)
        assert cleanup.done() is False
        assert writer.done() is False
        storage.release.set()
        with pytest.raises(asyncio.CancelledError):
            await cleanup
        version = await writer
        assert version.storage_key == target
        assert await repository.storage_key_referenced(target)
        assert await storage.exists(target)
    finally:
        await db1.dispose()
        await db2.dispose()


async def test_cancelled_retention_holds_guard_until_delete_finishes(tmp_path):
    fixture = await _version_race_fixture(tmp_path)
    db1, db2, repository, _service, storage, deck_id, initial_id, _edited, _callback, _target = fixture
    current_id = str(uuid4())
    current_key = f"decks/{deck_id}/versions/{current_id}.pptx"
    current_content = _pptx(2)
    await storage.put(current_key, current_content)
    await repository.append_version(
        deck_id=deck_id,
        owner_id="owner",
        version_id=current_id,
        storage_key=current_key,
        sha256=hashlib.sha256(current_content).hexdigest(),
        size_bytes=len(current_content),
        source="test",
        created_by="owner",
    )
    initial_key = f"decks/{deck_id}/versions/{initial_id}.pptx"
    storage.target = initial_key
    storage.pause_delete = True
    retention_service = DeckVersionService(repository, storage, None, 1_000_000, 1)
    waiter_entered = asyncio.Event()
    other = DeckRepository(db2, lock_dir=tmp_path / "locks")

    async def wait_for_guard():
        async with other.storage_key_guard(initial_key):
            waiter_entered.set()

    try:
        retention = asyncio.create_task(retention_service._enforce_retention(deck_id))
        await asyncio.wait_for(storage.delete_entered.wait(), 1)
        waiter = asyncio.create_task(wait_for_guard())
        retention.cancel()
        await asyncio.sleep(0.05)
        assert retention.done() is False
        assert waiter_entered.is_set() is False
        storage.release.set()
        with pytest.raises(asyncio.CancelledError):
            await retention
        await asyncio.wait_for(waiter, 1)
        current = await repository.get(deck_id, "owner")
        assert current.current_version_id == current_id
        assert await storage.exists(current_key)
    finally:
        await db1.dispose()
        await db2.dispose()
