from io import BytesIO
from unittest.mock import patch
from zipfile import ZIP_DEFLATED, ZipFile

from pptx import Presentation
import pytest

from app.services.presentation import pptx_validation
from app.services.presentation.pptx_validation import InvalidPptxError, validate_pptx


def _valid_pptx() -> bytes:
    presentation = Presentation()
    presentation.slides.add_slide(presentation.slide_layouts[0])
    output = BytesIO()
    presentation.save(output)
    return output.getvalue()


def _zip_with(*names: str) -> bytes:
    output = BytesIO()
    with ZipFile(output, "w", ZIP_DEFLATED) as package:
        for name in names:
            package.writestr(name, b"not-sensitive")
    return output.getvalue()


def _zip_with_payload(payload: bytes) -> bytes:
    output = BytesIO()
    with ZipFile(output, "w", ZIP_DEFLATED) as package:
        package.writestr("[Content_Types].xml", b"types")
        package.writestr("ppt/presentation.xml", b"presentation")
        package.writestr("ppt/slides/slide1.xml", payload)
    return output.getvalue()


def test_validate_pptx_accepts_valid_presentation():
    content = _valid_pptx()

    validate_pptx(content, max_bytes=len(content))


def test_validate_pptx_rejects_oversized_content():
    with pytest.raises(InvalidPptxError, match="exceeds maximum size"):
        validate_pptx(b"PK-too-large", max_bytes=2)


def test_validate_pptx_rejects_non_zip_content():
    with pytest.raises(InvalidPptxError, match="not a ZIP/OpenXML package"):
        validate_pptx(b"not-a-pptx", max_bytes=100)


def test_validate_pptx_rejects_bad_zip():
    with pytest.raises(InvalidPptxError, match="not readable"):
        validate_pptx(b"PKbad-zip", max_bytes=100)


@pytest.mark.parametrize(
    ("names", "message"),
    [
        (("ppt/presentation.xml", "ppt/slides/slide1.xml"), r"missing \[Content_Types\].xml"),
        (("[Content_Types].xml", "ppt/slides/slide1.xml"), "missing ppt/presentation.xml"),
        (("[Content_Types].xml", "ppt/presentation.xml"), "contains no slide XML parts"),
    ],
)
def test_validate_pptx_rejects_missing_required_parts(names: tuple[str, ...], message: str):
    with pytest.raises(InvalidPptxError, match=message):
        validate_pptx(_zip_with(*names), max_bytes=10_000)


def test_validate_pptx_rejects_package_python_pptx_cannot_open():
    content = _valid_pptx()

    with patch(
        "app.services.presentation.pptx_validation.Presentation",
        side_effect=ValueError("secret file content"),
    ):
        with pytest.raises(InvalidPptxError, match="cannot be opened as a presentation") as exc_info:
            validate_pptx(content, max_bytes=len(content))

    assert "secret file content" not in str(exc_info.value)


def test_validate_pptx_rejects_zero_slide_presentation_with_orphan_slide_member():
    presentation = Presentation()
    output = BytesIO()
    presentation.save(output)
    with ZipFile(output, "a", ZIP_DEFLATED) as package:
        package.writestr("ppt/slides/slide999.xml", b"<orphan/>")
    content = output.getvalue()

    with pytest.raises(InvalidPptxError, match="contains no presentation slides"):
        validate_pptx(content, max_bytes=len(content))


def test_validate_pptx_rejects_excessive_zip_member_count(monkeypatch):
    monkeypatch.setattr(pptx_validation, "MAX_ZIP_MEMBERS", 2)
    content = _valid_pptx()

    with pytest.raises(InvalidPptxError, match="too many ZIP members"):
        validate_pptx(content, max_bytes=len(content))


def test_validate_pptx_rejects_oversized_uncompressed_member(monkeypatch):
    monkeypatch.setattr(pptx_validation, "ZIP_MEMBER_SIZE_FLOOR", 1)
    payload = bytes(range(256)) * 4
    content = _zip_with_payload(payload)

    with pytest.raises(InvalidPptxError, match="ZIP member exceeds size limit"):
        validate_pptx(content, max_bytes=len(content))


def test_validate_pptx_rejects_excessive_total_uncompressed_size(monkeypatch):
    monkeypatch.setattr(pptx_validation, "ZIP_TOTAL_SIZE_FLOOR", 1)
    monkeypatch.setattr(pptx_validation, "ZIP_TOTAL_SIZE_MULTIPLIER", 1)
    payload = bytes(range(256)) * 4
    content = _zip_with_payload(payload)

    with pytest.raises(InvalidPptxError, match="ZIP expands beyond size limit"):
        validate_pptx(content, max_bytes=len(content))


def test_validate_pptx_rejects_extreme_compression_ratio():
    content = _zip_with_payload(b"A" * 10_000)

    with pytest.raises(InvalidPptxError, match="ZIP member compression ratio exceeds limit"):
        validate_pptx(content, max_bytes=len(content))
