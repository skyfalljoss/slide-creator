from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.slide import Slide
from pptx.util import Pt

from app.models.schemas import SlideData
from app.services.pptx_layout import LAYOUT
from app.services.pptx_text import clean_inline_text, split_label_body, table_row_from_text
from app.services.pptx_theme import WHITE


class PptxBlockMixin:
    """Renderer methods for AI-selected slide component blocks."""

    def _apply_blocks(self, slide: Slide, data: SlideData) -> None:
        top = self._add_content_header(
            slide, data.title, getattr(data, "kicker", None), getattr(data, "subtitle", None)
        )
        blocks = [b for b in (data.blocks or []) if isinstance(b, dict)]
        x = LAYOUT.left_margin
        w = LAYOUT.content_width
        bottom = LAYOUT.content_bottom
        if not blocks:
            self._add_card_grid(slide, data.bullets, top)
            return
        n = len(blocks)
        gap = LAYOUT.content_gap
        bh = (bottom - top - gap * (n - 1)) / n
        for i, block in enumerate(blocks):
            y = top + i * (bh + gap)
            self._render_block(slide, block, data, x, y, w, bh)

    def _render_block(self, slide: Slide, block: dict, data: SlideData, x: float, y: float, w: float, h: float) -> None:
        btype = str(block.get("type", "")).lower()
        if btype == "stat":
            self._block_stat(slide, block, x, y, w, h)
        elif btype == "quote":
            self._block_quote(slide, block, x, y, w, h)
        elif btype in ("table", "comparison"):
            self._block_table(slide, block, x, y, w, h)
        elif btype in ("process", "steps", "timeline"):
            self._block_process(slide, block, x, y, w, h)
        elif btype in ("cards", "card_grid", "grid", "columns"):
            self._block_cards(slide, block, x, y, w, h)
        elif btype == "image" and self._add_slide_image(
            slide,
            block.get("image_b64") or getattr(data, "image_b64", None),
            x,
            y,
            w,
            h,
        ):
            return
        else:
            items = block.get("items") or block.get("bullets") or data.bullets
            self._add_card_grid(slide, [str(i) for i in items], y)

    def _block_stat(self, slide: Slide, block: dict, x: float, y: float, w: float, h: float) -> None:
        value = str(block.get("value") or block.get("number") or "")
        label = str(block.get("label") or block.get("caption") or "")
        self._add_text(
            slide,
            x,
            y + h * 0.12,
            w,
            h * 0.55,
            value,
            130,
            self.theme.accent,
            bold=True,
            align=PP_ALIGN.CENTER,
        )
        if label:
            self._add_text(slide, x, y + h * 0.68, w, 0.9, label, 26, self.theme.text, align=PP_ALIGN.CENTER)

    def _block_quote(self, slide: Slide, block: dict, x: float, y: float, w: float, h: float) -> None:
        text = str(block.get("text") or block.get("quote") or "")
        author = str(block.get("author") or block.get("attribution") or "")
        bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, self._ix(x), self._iy(y + 0.2), self._ix(0.12), self._iy(h * 0.62))
        bar.fill.solid()
        bar.fill.fore_color.rgb = self.theme.accent
        bar.line.fill.background()
        bar.shadow.inherit = False
        self._add_text(slide, x + 0.6, y + 0.3, w - 0.8, h * 0.62, f"\u201c{text}\u201d", 34, self.theme.text, bold=True)
        if author:
            self._add_text(slide, x + 0.6, y + h * 0.78, w - 0.8, 0.5, author.upper(), 18, self.theme.muted)

    def _block_cards(self, slide: Slide, block: dict, x: float, y: float, w: float, h: float) -> None:
        raw = block.get("items") or block.get("cards") or block.get("columns_items") or []
        items: list[tuple[str, str, str]] = []
        for it in raw:
            if isinstance(it, dict):
                title = clean_inline_text(it.get("title", ""))
                body = clean_inline_text(it.get("body") or it.get("text") or "")
                if title and not body:
                    label, rest = split_label_body(title)
                    title = label
                    body = rest
                items.append((title, body, clean_inline_text(it.get("icon") or title)))
            else:
                title, body = split_label_body(it)
                if not body:
                    body = title
                    title = ""
                items.append((title, body, clean_inline_text(it)))
        if not items:
            return
        cols = int(block.get("columns") or 0) or (3 if len(items) >= 3 else max(len(items), 1))
        cols = max(1, min(cols, 4))
        rows = (len(items) + cols - 1) // cols
        gap = LAYOUT.content_gap
        cw = (w - gap * (cols - 1)) / cols
        ch = min((h - gap * (rows - 1)) / rows, h)
        pad = 0.4
        for i, (title, body, icon) in enumerate(items):
            r, c = divmod(i, cols)
            cx = x + c * (cw + gap)
            cy = y + r * (ch + gap)
            self._add_card(slide, cx, cy, cw, ch)
            self._add_icon_chip(slide, cx + pad, cy + pad, icon=icon)
            ty = cy + pad + 0.85
            if title:
                self._add_text(slide, cx + pad, ty, cw - 2 * pad, 0.45, title, 20, self.theme.text, bold=True)
                ty += 0.55
            if body:
                self._add_card_text(slide, cx + pad, ty, cw - 2 * pad, max(cy + ch - ty - 0.2, 0.4), body, 16, self.theme.muted)

    def _block_process(self, slide: Slide, block: dict, x: float, y: float, w: float, h: float) -> None:
        raw = block.get("steps") or block.get("items") or []
        steps = []
        for it in raw:
            if isinstance(it, dict):
                title = clean_inline_text(it.get("title", ""))
                body = clean_inline_text(it.get("body") or it.get("text") or "")
                if title and not body:
                    label, rest = split_label_body(title)
                    title = label
                    body = rest
                steps.append((title, body))
            else:
                title, body = split_label_body(it)
                steps.append((title if body else "", body or title))
        steps = steps[:4]
        if not steps:
            return
        cols = len(steps)
        gap = LAYOUT.content_gap
        cw = (w - gap * (cols - 1)) / cols
        ch = min(h, 3.8)
        pad = 0.4
        for i, (title, body) in enumerate(steps):
            cx = x + i * (cw + gap)
            self._add_card(slide, cx, y, cw, ch)
            self._add_number_circle(slide, cx + pad, y + pad, i + 1)
            ty = y + pad + 1.05
            if title:
                self._add_text(slide, cx + pad, ty, cw - 2 * pad, 0.45, title, 20, self.theme.text, bold=True)
                ty += 0.55
            self._add_card_text(slide, cx + pad, ty, cw - 2 * pad, max(y + ch - ty - 0.2, 0.4), body, 16, self.theme.muted)

    def _block_table(self, slide: Slide, block: dict, x: float, y: float, w: float, h: float) -> None:
        headers = [clean_inline_text(c) for c in (block.get("headers") or [])]
        rows = []
        for raw_row in block.get("rows") or []:
            if isinstance(raw_row, (list, tuple)):
                row = [clean_inline_text(c) for c in raw_row]
                if len(row) == 1:
                    row = table_row_from_text(row[0])
                elif len(row) >= 2 and not row[1]:
                    split = table_row_from_text(row[0])
                    if split[1]:
                        row = [split[0], split[1], *row[2:]]
                rows.append(row)
        if not headers and not rows:
            return
        ncols = len(headers) or (len(rows[0]) if rows else 1)
        nrows = (1 if headers else 0) + len(rows)
        if nrows == 0 or ncols == 0:
            return
        table_h = min(h, 0.7 + 0.75 * len(rows) + (0.7 if headers else 0))
        gtable = slide.shapes.add_table(nrows, ncols, self._ix(x), self._iy(y), self._ix(w), self._iy(table_h)).table
        gtable.first_row = bool(headers)
        gtable.horz_banding = True
        r0 = 0
        if headers:
            for c, htext in enumerate(headers):
                cell = gtable.cell(0, c)
                cell.text = htext
                p = cell.text_frame.paragraphs[0]
                p.font.name = "Arial"
                p.font.size = Pt(18)
                p.font.bold = True
                p.font.color.rgb = WHITE
                cell.fill.solid()
                cell.fill.fore_color.rgb = self._active_theme().strong
            r0 = 1
        for ri, row in enumerate(rows):
            for c in range(ncols):
                cell = gtable.cell(r0 + ri, c)
                cell.text = row[c] if c < len(row) else ""
                p = cell.text_frame.paragraphs[0]
                p.font.name = "Arial"
                p.font.size = Pt(16)
                p.font.color.rgb = self._map_color(self.theme.text if c == 0 else self.theme.muted)
