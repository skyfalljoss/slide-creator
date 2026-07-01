import re
import textwrap

from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE, PP_ALIGN
from pptx.slide import Slide
from pptx.util import Pt

from app.models.schemas import SlideData
from app.services.presentation.pptx_layout import LAYOUT
from app.services.presentation.variants import register_variant
from app.services.presentation.pptx_text import (
    clean_inline_text,
    is_leaked_instruction as _is_leaked_instruction,
    looks_like_instruction as _looks_like_instruction,
    split_label_body,
    table_row_from_text as _table_row_from_text,
)
from app.services.presentation.pptx_theme import CITI_RED, DISPLAY_FONT, THEMES, WHITE


_PROHIBITED_AGENDA_COPY = "The context and core insight for the discussion."


class PptxLayoutMixin:
    """Slide-specific layout renderers for the PPTX engine."""

    def _apply_title_slide(self, slide: Slide, data: SlideData) -> None:
        image_left = 10.15
        image_top = 0.92
        image_width = 6.8
        image_height = 8.14
        if not self._add_slide_image(slide, getattr(data, "image_b64", None), image_left, image_top, image_width, image_height):
            self._add_cover_visual_panel(slide, image_left, image_top, image_width, image_height)
        kicker = getattr(data, "kicker", None)
        if kicker:
            self._add_eyebrow(slide, LAYOUT.left_margin, 3.05, kicker)

        title_top = 3.45
        title_width = 8.9
        title_height = self._title_slide_title_height(data.title, title_width)
        self._add_text(slide, LAYOUT.left_margin, title_top, title_width, title_height, data.title, 48, self.theme.text, bold=True)

        accent_top = title_top + title_height + 0.22
        self._add_accent_bar(slide, LAYOUT.left_margin, accent_top, LAYOUT.accent_rule_width, height=0.06)

        subtitle = data.subtitle or (data.bullets[0] if data.bullets else "")
        secondary = data.bullets if data.subtitle else data.bullets[1:]
        subtitle_top = accent_top + 0.25
        if subtitle:
            self._add_text(slide, LAYOUT.left_margin, subtitle_top, 9.0, 0.8, subtitle, 22, self.theme.muted)
        if secondary:
            secondary_top = subtitle_top + (0.92 if subtitle else 0)
            self._add_bullets_box(slide, secondary, LAYOUT.left_margin, secondary_top, 9.2, 2.5)

    @staticmethod
    def _title_slide_title_height(title: str, width: float) -> float:
        chars_per_line = max(18, int(width / 0.44))
        lines = textwrap.wrap(title or "", width=chars_per_line) or [""]
        return min(max(1.25, len(lines) * 0.72), 3.25)

    @register_variant("big_statement")
    def _apply_big_statement(self, slide: Slide, data: SlideData) -> None:
        label = getattr(data, "kicker", None) or "CORE INSIGHT"
        self._add_eyebrow(slide, LAYOUT.left_margin, 1.95, label)
        self._add_text(slide, LAYOUT.left_margin, 2.55, 14.9, 2.35, data.title, 48, self.theme.text, bold=True)
        self._add_accent_bar(slide, LAYOUT.left_margin, 5.35, 1.7, height=0.06)
        if data.subtitle:
            self._add_text(slide, LAYOUT.left_margin, 5.7, 13.2, 0.75, data.subtitle, 22, self.theme.muted)
        elif data.bullets:
            bullet_text = "\n".join(data.bullets)
            panel_top = 5.75
            self._add_card(slide, LAYOUT.left_margin, panel_top, 14.9, 2.55, border_color=self._active_theme().panel_border)
            self._add_card_text(slide, LAYOUT.left_margin + 0.45, panel_top + 0.4, 13.9, 1.85, bullet_text, 16, self.theme.muted)

    @register_variant("three_points")
    def _apply_three_points(self, slide: Slide, data: SlideData) -> None:
        top = self._add_content_header(slide, data.title, getattr(data, "kicker", None), getattr(data, "subtitle", None))
        block = self._first_block(data)
        if block and str(block.get("type", "")).lower() in ("cards", "card_grid", "grid", "columns"):
            self._block_cards(slide, block, LAYOUT.left_margin, top, LAYOUT.content_width, LAYOUT.content_bottom - top)
            return
        card_items = data.bullets[:3]
        extra_items = data.bullets[3:]
        items = [{"title": item, "body": ""} for item in card_items]
        self._block_cards(
            slide,
            {"type": "cards", "columns": min(max(len(items), 1), 3), "items": items},
            LAYOUT.left_margin,
            top,
            LAYOUT.content_width,
            LAYOUT.content_bottom - top,
        )
        if extra_items:
            detail_top = top + 4.0
            self._add_bullets_box(slide, extra_items, LAYOUT.left_margin, detail_top, LAYOUT.content_width, 1.5)

    @register_variant("split_image")
    def _apply_split_image(self, slide: Slide, data: SlideData) -> None:
        left = LAYOUT.left_margin
        top = 1.55
        panel = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, self._ix(0), self._iy(0), self._ix(9.05), self._iy(self._canvas_h))
        panel.fill.solid()
        panel.fill.fore_color.rgb = THEMES["dark"].background
        panel.line.fill.background()
        with self._theme_mode(dark=True):
            self._add_eyebrow(slide, left, top, getattr(data, "kicker", None) or "CURRENT STATE")
            self._add_text(slide, left, top + 0.55, 7.9, 1.7, data.title, 34, self.theme.text, bold=True)
            text_top = top + 2.55
            if data.subtitle:
                self._add_text(slide, left, text_top, 7.2, 0.8, data.subtitle, 19, self.theme.muted)
                text_top += 1.0
            if data.bullets:
                self._add_bullets_box(slide, data.bullets, left, text_top, 7.2, 4.0)

        img_left = 9.7
        img_top = 1.6
        img_w = 7.25
        img_h = 6.95
        if not self._add_slide_image(slide, getattr(data, "image_b64", None), img_left, img_top, img_w, img_h):
            self._add_image_placeholder(slide, img_left, img_top, img_w, img_h)

    @register_variant("big_stat")
    def _apply_big_stat(self, slide: Slide, data: SlideData) -> None:
        block = self._first_block(data) or {}
        value = str(block.get("value") or block.get("number") or "")
        label = str(block.get("label") or block.get("caption") or "")
        if not value:
            source = " ".join([data.title, *data.bullets])
            match = re.search(r"[$#]?\d+(?:\.\d+)?\s*(?:%|x|M|B|K)?", source)
            if match:
                value = match.group(0)
            else:
                self._apply_big_statement(slide, data)
                return
        label_from_first_bullet = False
        if not label:
            label = data.subtitle or (data.bullets[0] if data.bullets else "")
            label_from_first_bullet = bool(not data.subtitle and data.bullets)

        kicker = getattr(data, "kicker", None) or "KEY FIGURE"
        has_chapter = bool(getattr(data, "chapter_number", None) is not None and clean_inline_text(getattr(data, "chapter_title", "")))
        if has_chapter:
            pill_width = min(max(1.45, len(kicker.upper()) * 0.105 + 0.55), 3.4)
            self._add_eyebrow(slide, 15.8 - pill_width, 0.23, kicker, width=3.4)
            title_top = 0.92
            rule_top = 1.75
        else:
            self._add_eyebrow(slide, LAYOUT.left_margin, 0.62, kicker)
            title_top = 1.02
            rule_top = 1.78
        if data.title and clean_inline_text(data.title) != clean_inline_text(value):
            self._add_text(slide, LAYOUT.left_margin, title_top, LAYOUT.header_title_width, 0.7, data.title, 28, self.theme.text, bold=True)
            self._add_accent_bar(slide, LAYOUT.left_margin, rule_top, LAYOUT.accent_rule_width, height=LAYOUT.accent_rule_height)

        self._add_text(slide, LAYOUT.left_margin, 2.2, LAYOUT.content_width, 1.55, value, 108, self.theme.accent, bold=True, align=PP_ALIGN.CENTER)
        if label:
            self._add_text(slide, 2.4, 3.83, self._LOGICAL_WIDTH - 4.8, 0.65, label, 26, self.theme.text, bold=True, align=PP_ALIGN.CENTER)
        remaining = data.bullets[1:] if label_from_first_bullet else data.bullets
        if remaining:
            self._add_big_stat_supporting_cards(slide, remaining, 5.05)

    @register_variant("before_after")
    def _apply_before_after(self, slide: Slide, data: SlideData) -> None:
        top = self._add_content_header(slide, data.title, getattr(data, "kicker", None), getattr(data, "subtitle", None))
        block = self._first_block(data) or {}
        before_title = "Before"
        after_title = "After"
        before_items: list[str] = []
        after_items: list[str] = []
        rows = block.get("rows") if isinstance(block, dict) else None
        headers = block.get("headers") if isinstance(block, dict) else None
        if isinstance(headers, list) and len(headers) >= 2:
            before_title = str(headers[-2])
            after_title = str(headers[-1])
        if isinstance(rows, list) and rows:
            for row in rows:
                if isinstance(row, (list, tuple)) and len(row) >= 2:
                    before_items.append(str(row[-2]))
                    after_items.append(str(row[-1]))
        if not before_items and data.bullets:
            mid = max(1, len(data.bullets) // 2)
            before_items = data.bullets[:mid]
            after_items = data.bullets[mid:] or data.bullets[:mid]
        self._comparison_card(slide, LAYOUT.left_margin, top, 7.85, 4.9, before_title, before_items, muted=True)
        self._comparison_card(slide, 9.1, top, 7.85, 4.9, after_title, after_items, muted=False)

    @register_variant("comparison_table")
    def _apply_comparison_table(self, slide: Slide, data: SlideData) -> None:
        top = self._add_content_header(slide, data.title, getattr(data, "kicker", None), getattr(data, "subtitle", None))
        block = self._first_block(data)
        if block and str(block.get("type", "")).lower() in ("table", "comparison"):
            self._block_table(slide, block, LAYOUT.left_margin, top + 0.15, LAYOUT.content_width, 8.6 - top)
            return
        rows = [_table_row_from_text(item) for item in data.bullets]
        self._block_table(slide, {"headers": ["Topic", "Detail"], "rows": rows}, LAYOUT.left_margin, top + 0.15, LAYOUT.content_width, 8.6 - top)

    @register_variant("process")
    def _apply_process_variant(self, slide: Slide, data: SlideData) -> None:
        if clean_inline_text(data.title).lower() == "presentation agenda":
            self._apply_agenda_overview(slide, data)
            return
        top = self._add_content_header(slide, data.title, getattr(data, "kicker", None), getattr(data, "subtitle", None))
        block = self._first_block(data)
        if block and str(block.get("type", "")).lower() in ("process", "steps", "timeline"):
            self._block_process(slide, block, LAYOUT.left_margin, top + 0.35, LAYOUT.content_width, 4.2)
            return
        self._block_process(
            slide,
            {"steps": [{"title": item, "body": ""} for item in data.bullets[:5]]},
            LAYOUT.left_margin,
            top + 0.35,
            LAYOUT.content_width,
            4.2,
        )

    def _apply_agenda_overview(self, slide: Slide, data: SlideData) -> None:
        self._add_text(slide, LAYOUT.left_margin, 0.25, 12.0, 0.75, "Presentation Agenda", 44, self.theme.text, bold=True)
        steps = self._agenda_items(data)
        left = LAYOUT.left_margin
        top = 1.55
        gap_x = 0.55
        gap_y = 0.45
        cw = (LAYOUT.content_width - gap_x) / 2
        ch = 2.75
        number_color = self._active_theme().accent
        for i, (title, body) in enumerate(steps):
            row, col = divmod(i, 2)
            x = left + col * (cw + gap_x)
            y = top + row * (ch + gap_y)
            self._add_card(slide, x, y, cw, ch)
            self._add_text(slide, x + 0.5, y + 0.3, 1.5, 0.58, f"{i + 1:02d}", 48, number_color, bold=True)
            title_box = self._add_text(
                slide,
                x + 0.5,
                y + 1.02,
                cw - 1.0,
                0.72,
                title,
                self._agenda_title_size(title),
                self.theme.text,
                bold=True,
            )
            body_box = self._add_card_text(
                slide,
                x + 0.5,
                y + 1.84,
                cw - 1.0,
                0.58,
                body,
                self._agenda_body_size(body),
                self.theme.muted,
            )
            title_box.text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
            body_box.text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE

    def _agenda_items(self, data: SlideData) -> list[tuple[str, str]]:
        block = self._first_block(data) or {}
        raw = block.get("steps") if isinstance(block, dict) else None
        items: list[tuple[str, str]] = []
        if isinstance(raw, list):
            for item in raw:
                if isinstance(item, dict):
                    title = clean_inline_text(item.get("title", ""))
                    body = clean_inline_text(item.get("body") or item.get("text") or "")
                    self._append_agenda_item(items, title, body)
                else:
                    title, body = split_label_body(item)
                    self._append_agenda_item(items, title, body)
        for item in data.bullets:
            title, body = split_label_body(item)
            self._append_agenda_item(items, title, body)
        defaults = ["The Core Insight", "Strategic Pillars", "Execution Pipeline", "Impact & Outlook"]
        return (items or [(title, self._agenda_card_body()) for title in defaults])[:4]

    def _append_agenda_item(self, items: list[tuple[str, str]], title: str, body: str) -> None:
        title = clean_inline_text(title)
        body = clean_inline_text(body)
        if (
            not title
            or title.casefold() == _PROHIBITED_AGENDA_COPY.casefold()
            or _is_leaked_instruction(title)
            or _looks_like_instruction(title)
        ):
            return
        title_key = title.removesuffix(":").strip().casefold()
        if any(existing.removesuffix(":").strip().casefold() == title_key for existing, _ in items):
            return
        if (
            not body
            or body.casefold() == _PROHIBITED_AGENDA_COPY.casefold()
            or _is_leaked_instruction(body)
            or _looks_like_instruction(body)
        ):
            body = self._agenda_card_body()
        items.append((title, body))

    @staticmethod
    def _agenda_card_body() -> str:
        return "Key context, evidence, and decisions for this chapter."

    @staticmethod
    def _agenda_title_size(title: str) -> int:
        if len(title) <= 34:
            return 23
        if len(title) <= 58:
            return 20
        return 17

    @staticmethod
    def _agenda_body_size(body: str) -> int:
        return 14 if len(body) <= 90 else 12

    @register_variant("quote")
    def _apply_quote_variant(self, slide: Slide, data: SlideData) -> None:
        block = self._first_block(data) or {}
        quote = str(block.get("text") or block.get("quote") or (data.bullets[0] if data.bullets else data.title))
        author = str(block.get("author") or block.get("attribution") or "")
        self._add_eyebrow(slide, LAYOUT.left_margin, 1.75, getattr(data, "kicker", None) or "VISION")
        bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, self._ix(LAYOUT.left_margin), self._iy(2.65), self._ix(0.12), self._iy(3.1))
        bar.fill.solid()
        bar.fill.fore_color.rgb = self._active_theme().accent
        bar.line.fill.background()
        self._add_text(slide, 1.45, 2.55, 13.8, 2.7, f"\"{quote}\"", 34, self.theme.text, bold=True)
        if author:
            self._add_text(slide, 1.45, 5.65, 10.0, 0.45, author.upper(), 16, self.theme.muted)
        if getattr(data, "image_b64", None):
            self._add_slide_image(slide, data.image_b64, 12.9, 5.85, 3.6, 2.15)

    @register_variant("closing")
    def _apply_closing(self, slide: Slide, data: SlideData) -> None:
        label = getattr(data, "kicker", None) or "NEXT STEPS"
        self._add_eyebrow(slide, LAYOUT.left_margin, 1.55, label)
        self._add_text(slide, LAYOUT.left_margin, 2.02, 13.2, 1.05, data.title, 44, self.theme.text, bold=True)
        self._add_accent_bar(slide, LAYOUT.left_margin, 3.45, 1.5, height=0.06)
        if data.subtitle:
            self._add_text(slide, LAYOUT.left_margin, 3.72, 9.5, 0.5, data.subtitle, 20, self.theme.muted)
        block = self._first_block(data)
        if block and str(block.get("type", "")).lower() in ("process", "steps", "timeline"):
            self._block_process(slide, block, LAYOUT.left_margin, 4.25, LAYOUT.content_width, 2.4)
        elif data.bullets:
            self._add_bullets_box(slide, data.bullets, LAYOUT.left_margin, 4.25, 9.7, 3.0)

    def _add_big_stat_supporting_cards(self, slide: Slide, items: list[str], top: float) -> None:
        visible = items[:3]
        if not visible:
            return
        gap = LAYOUT.content_gap
        usable = LAYOUT.content_width
        cw = (usable - gap * (len(visible) - 1)) / len(visible)
        ch = min(LAYOUT.content_bottom - top, 2.55)
        for i, item in enumerate(visible):
            x = LAYOUT.left_margin + i * (cw + gap)
            title, body = split_label_body(item)
            if not body:
                body = title
                title = ""
            self._add_card(slide, x, top, cw, ch)
            self._add_icon_chip(slide, x + 0.35, top + 0.35, size=0.54, icon=title or body)
            text_left = x + 1.05
            text_width = cw - 1.4
            text_top = top + 0.34
            if title:
                self._add_text(slide, text_left, text_top, text_width, 0.62, title, 17, self.theme.text, bold=True)
                text_top += 0.74
            else:
                text_top += 0.18
            self._add_card_text(slide, text_left, text_top, text_width, max(top + ch - text_top - 0.25, 0.5), body, 14, self.theme.muted)

    def _comparison_card(self, slide: Slide, x: float, y: float, w: float, h: float, title: str, items: list[str], *, muted: bool) -> None:
        border = self.theme.muted if muted else self.theme.accent
        self._add_card(slide, x, y, w, h, border_color=border)
        self._add_text(slide, x + 0.45, y + 0.45, w - 0.9, 0.45, title, 24, border, bold=True)
        text = "\n".join(items or ["No details available."])
        self._add_card_text(slide, x + 0.45, y + 1.15, w - 0.9, h - 1.45, text, 18, self.theme.muted if muted else self.theme.text)

    def _add_image_placeholder(self, slide: Slide, left: float, top: float, width: float, height: float) -> None:
        rect = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, self._ix(left), self._iy(top), self._ix(width), self._iy(height))
        rect.fill.solid()
        rect.fill.fore_color.rgb = self._active_theme().surface
        rect.line.color.rgb = self._active_theme().border
        rect.line.dash_style = 4
        rect.shadow.inherit = False
        self._add_text(slide, left + 0.55, top + height / 2 - 0.25, width - 1.1, 0.5, "Image Placeholder", 18, self.theme.muted, align=PP_ALIGN.CENTER)

    def _add_cover_visual_panel(self, slide: Slide, left: float, top: float, width: float, height: float) -> None:
        panel = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, self._ix(left), self._iy(top), self._ix(width), self._iy(height))
        panel.fill.solid()
        panel.fill.fore_color.rgb = self._active_theme().panel_bg
        panel.line.color.rgb = self._active_theme().panel_border
        panel.shadow.inherit = False
        for offset, alpha_width in ((0.55, 0.08), (1.05, 0.04), (1.55, 0.04)):
            stripe = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                self._ix(left + offset),
                self._iy(top + 0.65),
                self._ix(alpha_width),
                self._iy(height - 1.3),
            )
            stripe.fill.solid()
            stripe.fill.fore_color.rgb = self._active_theme().accent
            stripe.line.fill.background()
            stripe.shadow.inherit = False

    def _apply_section_divider(self, slide: Slide, data: SlideData) -> None:
        self._guard_divider_text(data.title)
        subtitle = data.subtitle or (data.bullets[0] if data.bullets else "")
        if not subtitle or _is_leaked_instruction(subtitle):
            subtitle = "Section overview"
        self._guard_divider_text(subtitle)

        image_b64 = getattr(data, "image_b64", None)
        if image_b64 and self._add_fullbleed_image(slide, image_b64):
            self._add_scrim(slide)
            light = RGBColor(0xFF, 0xFF, 0xFF)
            mid = self._LOGICAL_WIDTH / 2
            box_w = self._LOGICAL_WIDTH - 3.0
            self._add_text(slide, 1.5, 3.45, box_w, 1.0, data.title, 40, light, bold=True, align=PP_ALIGN.CENTER)
            self._add_accent_bar(slide, mid - 0.7, 4.65, LAYOUT.accent_rule_width, height=LAYOUT.accent_rule_height)
            self._add_text(slide, 1.5, 4.9, box_w, 0.6, subtitle, 20, light, align=PP_ALIGN.CENTER)
            return

        self._add_text(slide, LAYOUT.left_margin, 3.6, 9.7, 0.76, data.title, 36, self.theme.text, bold=True)
        self._add_accent_bar(slide, LAYOUT.left_margin, 4.55, 1.11, height=0.06)
        self._add_text(slide, LAYOUT.left_margin, 4.85, 8.33, 0.42, subtitle, 20, self.theme.muted)

    def _guard_divider_text(self, text: str) -> None:
        """Raise if divider text reads like a leaked generation instruction."""
        if _looks_like_instruction(text):
            raise ValueError(
                "Section divider text reads like a generation instruction, "
                f"not client-facing content: {text[:80]!r}..."
            )

    def _apply_executive_summary(self, slide: Slide, data: SlideData) -> None:
        top = self._add_content_header(slide, data.title, getattr(data, "kicker", None) or "EXECUTIVE SUMMARY")

        bullets = data.bullets or []
        mid = (len(bullets) + 1) // 2 if bullets else 0
        columns = [
            ("Overview & Context", bullets[:mid] or ["No overview details available."]),
            ("Key Objectives & Focus Areas", bullets[mid:] or ["No focus area details available."]),
        ]
        gap = LAYOUT.content_gap
        usable = LAYOUT.content_width
        cw = (usable - gap) / 2
        ch = LAYOUT.content_bottom - top
        for i, (heading, items) in enumerate(columns):
            x = LAYOUT.left_margin + i * (cw + gap)
            self._add_card(slide, x, top, cw, ch)
            self._add_icon_chip(slide, x + 0.45, top + 0.45)
            self._add_text(slide, x + 0.45, top + 1.3, cw - 0.9, 0.45, heading, 22, self.theme.text, bold=True)
            self._add_card_text(slide, x + 0.45, top + 2.0, cw - 0.9, ch - 2.4, "\n".join(items), 18, self.theme.muted)

    def _apply_next_steps(self, slide: Slide, data: SlideData) -> None:
        top = self._add_content_header(slide, data.title, getattr(data, "kicker", None) or "NEXT STEPS")
        bullets = data.bullets or []
        steps = bullets[:5]
        if not steps:
            self._add_visual_panel(slide, data, LAYOUT.left_margin, top, LAYOUT.content_width, LAYOUT.content_bottom - top)
            return

        gap = LAYOUT.content_gap
        usable = LAYOUT.content_width
        cols = len(steps)
        cw = (usable - gap * (cols - 1)) / cols
        ch = 3.6
        for i, bullet in enumerate(steps):
            x = LAYOUT.left_margin + i * (cw + gap)
            self._add_card(slide, x, top, cw, ch)
            self._add_number_circle(slide, x + 0.45, top + 0.45, i + 1)
            self._add_card_text(slide, x + 0.45, top + 1.5, cw - 0.9, ch - 1.8, bullet, 18, self.theme.text)

        lower = bullets[5:]
        if lower:
            sy = top + ch + 0.5
            self._add_eyebrow(slide, LAYOUT.left_margin, sy, "Recommended Timeline")
            self._add_card_text(slide, LAYOUT.left_margin, sy + 0.5, usable, max(LAYOUT.content_bottom - (sy + 0.5), 0.6), "  •  ".join(lower), 18, self.theme.muted)

    def _apply_chart_slide(self, slide: Slide, data: SlideData) -> None:
        top = self._add_content_header(slide, data.title, getattr(data, "kicker", None))
        self._add_bullets_box(slide, data.bullets, LAYOUT.left_margin, top, 6.0, 8.6 - top)
        self._add_vertical_divider(slide, 7.1, top, 8.4 - top)
        ch_h = 8.0 - top
        if data.chart_data:
            self._add_chart(slide, data.chart_data, left=7.5, top=top, width=9.44, height=ch_h)
            self._add_source_note(slide, data.chart_audit, 7.5, top + ch_h + 0.1, 9.44)
        elif not self._add_slide_image(slide, getattr(data, "image_b64", None), 7.5, top, 9.44, ch_h):
            self._add_visual_panel(slide, data, 7.5, top, 9.44, ch_h)

    def _apply_standard_content(self, slide: Slide, data: SlideData) -> None:
        top = self._add_content_header(slide, data.title, getattr(data, "kicker", None))
        if data.callout:
            self._add_callout_box(slide, data.callout, LAYOUT.left_margin, top, LAYOUT.content_width, 0.6)
            top += 0.85
        if data.bullets:
            self._add_modern_bullet_panel(slide, data.bullets, top)
        else:
            self._add_visual_panel(slide, data, LAYOUT.left_margin, top, LAYOUT.content_width, LAYOUT.content_bottom - top)

    def _add_visual_panel(self, slide: Slide, data: SlideData, left: float, top: float, width: float, height: float) -> None:
        card = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, self._ix(left), self._iy(top), self._ix(width), self._iy(height))
        card.fill.solid()
        card.fill.fore_color.rgb = self._active_theme().panel_bg
        card.line.color.rgb = self._active_theme().panel_bg

        stripe = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, self._ix(left), self._iy(top), self._ix(0.08), self._iy(height))
        stripe.fill.solid()
        stripe.fill.fore_color.rgb = self._active_theme().accent
        stripe.line.color.rgb = self._active_theme().accent

        self._add_text(slide, left + 0.35, top + 0.35, width - 0.7, 0.35, "Visual Direction & Context", 16, self.theme.accent, bold=True)
        text = data.visual_direction or "Use a clean Citi-style callout visual for the slide's core message."
        self._add_card_text(slide, left + 0.35, top + 0.9, width - 0.7, height - 1.2, text, 16, self.theme.text)

    def _add_content_header(self, slide: Slide, title: str, kicker: str | None, subtitle: str | None = None) -> float:
        """Eyebrow + title (+ optional subtitle) + accent rule. Returns content top."""
        data = getattr(self, "_current_slide_data", None)
        has_chapter = bool(
            data is not None
            and getattr(data, "chapter_number", None) is not None
            and clean_inline_text(getattr(data, "chapter_title", ""))
        )
        if kicker and has_chapter:
            # Chapter marker owns the top band; kicker + title sit clearly below it.
            self._add_eyebrow(slide, LAYOUT.left_margin, 1.34, kicker)
            self._add_text(slide, LAYOUT.left_margin, 1.9, LAYOUT.header_title_width, 0.9, title, 30, self.theme.text, bold=True)
            rule_top = 2.66
        elif kicker:
            self._add_eyebrow(slide, LAYOUT.left_margin, 0.62, kicker)
            self._add_text(slide, LAYOUT.left_margin, 1.02, LAYOUT.header_title_width, 0.9, title, 30, self.theme.text, bold=True)
            rule_top = 1.78
        elif has_chapter:
            self._add_text(slide, LAYOUT.left_margin, 1.34, LAYOUT.header_title_width, 0.9, title, 30, self.theme.text, bold=True)
            rule_top = 2.18
        else:
            self._add_text(slide, LAYOUT.left_margin, 0.66, LAYOUT.header_title_width, 0.9, title, 30, self.theme.text, bold=True)
            rule_top = 1.6
        self._add_accent_bar(slide, LAYOUT.left_margin, rule_top, LAYOUT.accent_rule_width, height=LAYOUT.accent_rule_height)
        top = rule_top + 0.45
        if subtitle:
            self._add_text(slide, LAYOUT.left_margin, top, LAYOUT.header_title_width, 0.5, subtitle, 18, self.theme.muted)
            top += 0.65
        return top

    def _add_chapter_marker(self, slide: Slide, data: SlideData) -> None:
        chapter_number = getattr(data, "chapter_number", None)
        chapter_title = clean_inline_text(getattr(data, "chapter_title", ""))
        if chapter_number is None or not chapter_title:
            return
        if (
            data.layout.lower() == "title"
            or clean_inline_text(data.title).lower() == "presentation agenda"
            or (data.variant or "").lower() == "closing"
        ):
            return

        left = LAYOUT.left_margin
        top = 0.18
        badge_width = 0.68
        badge_height = 0.46
        badge = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            self._ix(left),
            self._iy(top),
            self._ix(badge_width),
            self._iy(badge_height),
        )
        badge.fill.solid()
        badge.fill.fore_color.rgb = CITI_RED
        badge.line.fill.background()
        badge.shadow.inherit = False
        try:
            badge.adjustments[0] = 0.35
        except (IndexError, TypeError, ValueError):
            pass

        text_frame = badge.text_frame
        text_frame.clear()
        text_frame.margin_left = 0
        text_frame.margin_right = 0
        text_frame.margin_top = 0
        text_frame.margin_bottom = 0
        text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        paragraph = text_frame.paragraphs[0]
        paragraph.text = f"{chapter_number:02d}"
        paragraph.alignment = PP_ALIGN.CENTER
        paragraph.font.name = DISPLAY_FONT
        paragraph.font.size = Pt(13)
        paragraph.font.bold = True
        paragraph.font.color.rgb = WHITE

        dark_backed = (
            self._active_dark
            or (data.variant or "").lower() == "split_image"
            or (data.layout.lower() == "section_divider" and bool(getattr(data, "image_b64", None)))
        )
        label_color = THEMES["dark"].text if dark_backed else self._active_theme().text
        label = self._add_text(
            slide,
            left + badge_width + 0.2,
            top,
            9.5,
            badge_height,
            chapter_title.upper(),
            16,
            label_color,
            bold=True,
        )
        label.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

    def _add_card_grid(self, slide: Slide, items: list[str], top: float) -> None:
        """Lay bullets out as a responsive grid of icon-chip surface cards."""
        n = len(items)
        if n == 0:
            return
        cols = 1 if n == 1 else 2 if n == 2 else 2 if n == 4 else 3
        rows = (n + cols - 1) // cols
        left = LAYOUT.left_margin
        gap = LAYOUT.content_gap
        usable = LAYOUT.content_width
        cw = (usable - gap * (cols - 1)) / cols
        bottom = LAYOUT.content_bottom
        ch = min((bottom - top - gap * (rows - 1)) / rows, 3.1)
        pad = 0.4
        for i, item in enumerate(items):
            row, col = divmod(i, cols)
            x = left + col * (cw + gap)
            y = top + row * (ch + gap)
            self._add_card(slide, x, y, cw, ch)
            self._add_icon_chip(slide, x + pad, y + pad, icon=item)
            self._add_card_text(slide, x + pad, y + pad + 0.85, cw - pad * 2, ch - pad - 0.95, item, 18, self.theme.text)

    def _add_modern_bullet_panel(self, slide: Slide, items: list[str], top: float) -> None:
        panel_h = min(LAYOUT.content_bottom - top, 5.9)
        self._add_card(slide, LAYOUT.left_margin, top, 10.8, panel_h)
        self._add_accent_bar(slide, LAYOUT.left_margin + 0.48, top + 0.52, 1.0, height=0.05)
        self._add_bullets_box(slide, items, LAYOUT.left_margin + 0.55, top + 0.95, 9.4, panel_h - 1.25)
