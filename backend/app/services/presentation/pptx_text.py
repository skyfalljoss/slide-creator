import re

from pptx.util import Pt


_ICON_SHAPES: list[tuple[tuple[str, ...], str]] = [
    (("speed", "velocity", "fast", "performance", "accelerate", "efficiency", "agile"), "LIGHTNING_BOLT"),
    (("security", "secure", "compliance", "trust", "protect", "risk", "privacy", "governance"), "HEXAGON"),
    (("scale", "growth", "grow", "increase", "expand", "revenue", "value", "roi", "return"), "UP_ARROW"),
    (("global", "network", "reach", "connect", "integration", "ecosystem", "api", "market"), "DONUT"),
    (("process", "system", "automation", "operations", "workflow", "engine", "infrastructure"), "GEAR_6"),
    (("quality", "excellence", "premium", "star", "award", "leader", "best"), "STAR_5_POINT"),
    (("innovation", "idea", "insight", "vision", "future", "strategy", "opportunity"), "SUN"),
    (("client", "customer", "people", "team", "partner", "relationship", "service"), "HEART"),
    (("next", "step", "action", "timeline", "plan", "roadmap", "phase"), "CHEVRON"),
    (("data", "cloud", "platform", "technology", "digital", "analytics"), "CLOUD"),
]


def icon_shape(keyword: str | None):
    from pptx.enum.shapes import MSO_SHAPE

    text = (keyword or "").lower()
    if text:
        for words, name in _ICON_SHAPES:
            if any(w in text for w in words):
                shape = getattr(MSO_SHAPE, name, None)
                if shape is not None:
                    return shape
    return MSO_SHAPE.DIAMOND


def clean_inline_text(text: object) -> str:
    cleaned = re.sub(r"\*\*(.+?)\*\*", r"\1", str(text or ""))
    return re.sub(r"\s+", " ", cleaned).strip()


def split_label_body(text: object) -> tuple[str, str]:
    cleaned = clean_inline_text(text)
    if ":" not in cleaned:
        return cleaned, ""
    label, body = cleaned.split(":", 1)
    if 1 <= len(label.strip()) <= 36 and body.strip():
        return f"{label.strip()}:", body.strip()
    return cleaned, ""


def table_row_from_text(text: object) -> list[str]:
    topic, detail = split_label_body(text)
    return [topic.rstrip(":"), detail]


def add_markdown_paragraph(text_frame, line, font_size, font_name, color_rgb):
    if len(text_frame.paragraphs) == 1 and not text_frame.paragraphs[0].runs:
        paragraph = text_frame.paragraphs[0]
    else:
        paragraph = text_frame.add_paragraph()
    paragraph.space_before = Pt(8)

    segments = re.split(r"\*\*(.+?)\*\*", line)
    for index, segment in enumerate(segments):
        if segment == "":
            continue
        run = paragraph.add_run()
        run.text = segment
        run.font.bold = index % 2 == 1
        run.font.size = Pt(font_size)
        run.font.name = font_name
        run.font.color.rgb = color_rgb
    return paragraph


_INSTRUCTION_VERBS = frozenset(
    {
        "use", "create", "add", "show", "place", "center", "centered", "convey",
        "ensure", "include", "display", "render", "apply", "make", "design",
        "arrange", "position", "highlight", "emphasize", "illustrate", "depict",
        "showcase", "present", "draw", "feature",
    }
)

_LEAK_PREFIXES = ("bold title", "centered on the slide")
_LEAK_PHRASES = ("potentially with", "background image related to")


def is_leaked_instruction(text: str) -> bool:
    lowered = text.strip().lower()
    if lowered.startswith(_LEAK_PREFIXES):
        return True
    return any(phrase in lowered for phrase in _LEAK_PHRASES)


def looks_like_instruction(text: str) -> bool:
    stripped = text.strip()
    if len(stripped) <= 120:
        return False
    if "potentially" in stripped.lower():
        return True
    first_word = re.split(r"[\s,]+", stripped, maxsplit=1)[0]
    return bool(first_word) and first_word[0].islower() and first_word.lower().strip(".,:;") in _INSTRUCTION_VERBS
