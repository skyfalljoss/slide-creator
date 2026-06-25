from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.slide import Slide

from app.models.schemas import ChartAudit, ChartData


class PptxChartMixin:
    """Chart rendering helpers for PPTX slides."""

    def _add_chart(
        self,
        slide: Slide,
        chart_data: ChartData | dict,
        *,
        left: float = 6.7,
        top: float = 1.6,
        width: float = 5.8,
        height: float = 4.6,
    ) -> None:
        if isinstance(chart_data, ChartData):
            chart_data = chart_data.model_dump()
        if not isinstance(chart_data, dict):
            return
        categories = chart_data.get("categories")
        series_items = chart_data.get("series")
        if not isinstance(categories, (list, tuple)) or not isinstance(series_items, list):
            return
        for series in series_items:
            if not isinstance(series, dict) or not isinstance(series.get("values"), (list, tuple)):
                return

        chart = CategoryChartData()
        chart.categories = [str(category) for category in categories]
        for series in series_items:
            chart.add_series(str(series.get("name", "Series")), series.get("values", []))

        chart_type = XL_CHART_TYPE.LINE_MARKERS if chart_data.get("type") == "line" else XL_CHART_TYPE.COLUMN_CLUSTERED
        graphic_frame = slide.shapes.add_chart(
            chart_type,
            self._ix(left),
            self._iy(top),
            self._ix(width),
            self._iy(height),
            chart,
        )
        title = str(chart_data.get("title", ""))
        if title:
            graphic_frame.chart.has_title = True
            graphic_frame.chart.chart_title.text_frame.text = title

    def _add_source_note(self, slide: Slide, chart_audit: ChartAudit | None, left: float, top: float, width: float) -> None:
        if chart_audit is None:
            return
        columns = ", ".join([chart_audit.category_column, *chart_audit.value_columns])
        text = f"Source: {chart_audit.source_filename}; columns: {columns}"
        self._add_text(slide, left, top, width, 0.3, text, 8, self.theme.muted)
