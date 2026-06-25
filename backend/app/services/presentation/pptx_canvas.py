import base64
import io
import logging
from collections.abc import Callable

from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Inches, Pt

from app.services.media.icons import render_icon_png
from app.services.presentation.pptx_layout import LAYOUT
from app.services.presentation.pptx_text import add_markdown_paragraph, clean_inline_text, icon_shape
from app.services.presentation.pptx_theme import BODY_FONT, DISPLAY_FONT, WHITE, Theme

logger = logging.getLogger(__name__)


class PptxCanvas:
    """Low-level drawing wrapper for python-pptx primitives."""

    def __init__(
        self,
        *,
        canvas_width: float,
        canvas_height: float,
        logical_width: float,
        theme: Theme,
        active_theme: Callable[[], Theme],
        map_color: Callable[[RGBColor], RGBColor],
    ) -> None:
        self.canvas_width = canvas_width
        self.canvas_height = canvas_height
        self.logical_width = logical_width
        self.theme = theme
        self._active_theme = active_theme
        self._map_color = map_color
        self._hscale = canvas_width / logical_width

    def ix(self, value: float) -> Inches:
        """Map a logical X position or width onto the target canvas."""
        return Inches(value * self._hscale)

    def iy(self, value: float) -> Inches:
        """Y positions and heights are unscaled (canvas height is fixed at 10in)."""
        return Inches(value)

    def clear_text_frame(self, text_frame) -> None:
        text_frame.clear()
        text_frame.word_wrap = True
        text_frame.margin_left = Inches(0)
        text_frame.margin_right = Inches(0)
        text_frame.margin_top = Inches(0)
        text_frame.margin_bottom = Inches(0)

    def add_text(
        self,
        slide,
        left: float,
        top: float,
        width: float,
        height: float,
        text: str,
        size: int,
        color: RGBColor,
        *,
        bold: bool = False,
        align: "PP_ALIGN | None" = None,
    ):
        text_box = slide.shapes.add_textbox(self.ix(left), self.iy(top), self.ix(width), self.iy(height))
        text_frame = text_box.text_frame
        self.clear_text_frame(text_frame)
        paragraph = text_frame.paragraphs[0]
        paragraph.text = clean_inline_text(text)
        paragraph.font.name = DISPLAY_FONT if bold or size >= 28 else BODY_FONT
        paragraph.font.size = Pt(size)
        paragraph.font.bold = bold
        paragraph.font.color.rgb = self._map_color(color)
        if align is not None:
            paragraph.alignment = align
        return text_box

    def add_card_text(self, slide, x: float, y: float, w: float, h: float, text: str, size: int, color: RGBColor):
        """Markdown-aware text that stays within a fixed box (no autosize growth)."""
        box = slide.shapes.add_textbox(self.ix(x), self.iy(y), self.ix(w), self.iy(h))
        text_frame = box.text_frame
        self.clear_text_frame(text_frame)
        mapped = self._map_color(color)
        for line in text.split("\n"):
            add_markdown_paragraph(text_frame, line, size, BODY_FONT, mapped)
        return box

    def add_bullets_box(self, slide, bullets: list[str], left: float, top: float, width: float, height: float):
        text_box = slide.shapes.add_textbox(self.ix(left), self.iy(top), self.ix(width), self.iy(height))
        text_frame = text_box.text_frame
        self.clear_text_frame(text_frame)
        text_frame.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT
        for bullet in bullets:
            add_markdown_paragraph(text_frame, f"\u2022 {bullet}", 18, BODY_FONT, self._map_color(self.theme.text))
        return text_box

    def add_card(self, slide, x: float, y: float, w: float, h: float, border_color: RGBColor | None = None):
        """A surface-filled rounded rectangle with a hairline border."""
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, self.ix(x), self.iy(y), self.ix(w), self.iy(h))
        card.fill.solid()
        card.fill.fore_color.rgb = self._active_theme().surface
        card.line.color.rgb = self._map_color(border_color) if border_color else self._active_theme().border
        card.line.width = Pt(1)
        card.shadow.inherit = False
        try:
            card.adjustments[0] = 0.05
        except (IndexError, TypeError, ValueError):
            logger.debug("Failed to adjust card corner radius", exc_info=True)
        return card

    def add_icon_chip(self, slide, x: float, y: float, size: float = 0.62, icon: str | None = None) -> None:
        """Accent-tinted rounded square with a centered Font Awesome icon."""
        chip = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, self.ix(x), self.iy(y), self.iy(size), self.iy(size))
        chip.fill.solid()
        chip.fill.fore_color.rgb = self._active_theme().accent_soft
        chip.line.fill.background()
        chip.shadow.inherit = False
        try:
            chip.adjustments[0] = 0.3
        except (IndexError, TypeError, ValueError):
            logger.debug("Failed to adjust icon chip corner radius", exc_info=True)
        inner = size * 0.52
        off = (size - inner) / 2
        png = render_icon_png(icon, str(self._active_theme().accent))
        if png:
            slide.shapes.add_picture(io.BytesIO(png), self.ix(x + off), self.iy(y + off), self.iy(inner), self.iy(inner))
            return
        mark = slide.shapes.add_shape(icon_shape(icon), self.ix(x + off), self.iy(y + off), self.iy(inner), self.iy(inner))
        mark.fill.solid()
        mark.fill.fore_color.rgb = self._active_theme().accent
        mark.line.fill.background()
        mark.shadow.inherit = False

    def add_accent_bar(self, slide, left: float, top: float, width: float, height: float = LAYOUT.accent_rule_height) -> None:
        bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            self.ix(left),
            self.iy(top),
            self.ix(width),
            self.iy(height * self.theme.accent_weight),
        )
        bar.fill.solid()
        bar.fill.fore_color.rgb = self._active_theme().accent
        bar.line.color.rgb = self._active_theme().accent

    def add_vertical_divider(self, slide, left: float, top: float, height: float):
        line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, self.ix(left), self.iy(top), self.ix(0.01), self.iy(height))
        line.fill.solid()
        line.fill.fore_color.rgb = self._active_theme().panel_border
        line.line.color.rgb = self._active_theme().panel_border
        return line

    def add_number_circle(self, slide, x: float, y: float, number: int, size: float = 0.72) -> None:
        circ = slide.shapes.add_shape(MSO_SHAPE.OVAL, self.ix(x), self.iy(y), self.iy(size), self.iy(size))
        circ.fill.solid()
        circ.fill.fore_color.rgb = self._active_theme().accent
        circ.line.fill.background()
        circ.shadow.inherit = False
        text_frame = circ.text_frame
        text_frame.word_wrap = False
        text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        paragraph = text_frame.paragraphs[0]
        paragraph.text = str(number)
        paragraph.alignment = PP_ALIGN.CENTER
        paragraph.font.name = DISPLAY_FONT
        paragraph.font.size = Pt(22)
        paragraph.font.bold = True
        paragraph.font.color.rgb = WHITE

    def add_scrim(self, slide, color: RGBColor = RGBColor(0x0A, 0x16, 0x28), opacity: int = 58) -> None:
        """Overlay a semi-transparent dark rectangle over the whole slide for contrast."""
        rect = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            self.ix(0),
            self.iy(0),
            self.ix(self.logical_width),
            self.iy(self.canvas_height),
        )
        rect.line.fill.background()
        rect.fill.solid()
        rect.fill.fore_color.rgb = color
        srgb = rect.fill.fore_color._xFill.find(qn("a:srgbClr"))
        if srgb is not None:
            alpha = srgb.makeelement(qn("a:alpha"), {"val": str(int(opacity * 1000))})
            srgb.append(alpha)

    def add_slide_image(self, slide, image_b64, left: float, top: float, width: float, height: float) -> bool:
        """Insert a base64 image at the given rect. Returns True on success."""
        if not image_b64:
            return False
        try:
            image_data = base64.b64decode(image_b64)
            slide.shapes.add_picture(io.BytesIO(image_data), self.ix(left), self.iy(top), self.ix(width), self.iy(height))
            return True
        except Exception:
            logger.warning("Failed to insert slide image", exc_info=True)
            return False
