from collections.abc import Iterator
from contextlib import contextmanager
import io
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.presentation import Presentation as PresentationType
from pptx.slide import Slide, SlideLayout
from pptx.util import Inches, Pt

from app.config import settings
from app.models.schemas import SlideData
from app.services.presentation.pptx_blocks import PptxBlockMixin
from app.services.presentation.pptx_canvas import PptxCanvas
from app.services.presentation.pptx_charts import PptxChartMixin
from app.services.presentation.pptx_layout import CANVAS_DIMS, LAYOUT
from app.services.presentation.pptx_layouts import PptxLayoutMixin
from app.services.presentation.pptx_theme import DISPLAY_FONT, THEMES, Theme, resolve_theme


class PptxEngine(PptxLayoutMixin, PptxChartMixin, PptxBlockMixin):
    # Logical authoring width (16:9). All layout coordinates are written against
    # this width and scaled horizontally to the target canvas.
    _LOGICAL_WIDTH = LAYOUT.logical_width
    _CANVAS_DIMS = CANVAS_DIMS
    _DARK_VARIANTS = frozenset({"big_statement", "big_stat", "quote", "closing"})

    def __init__(self, template_path: str | None = None, theme: str | Theme | None = None, aspect_ratio: str = "16:9"):
        self.template_path = template_path
        self.theme = theme if isinstance(theme, Theme) else resolve_theme(theme)
        self.aspect_ratio = aspect_ratio if aspect_ratio in self._CANVAS_DIMS else "16:9"
        self._canvas_w, self._canvas_h = self._CANVAS_DIMS[self.aspect_ratio]
        self._hscale = self._canvas_w / self._LOGICAL_WIDTH
        # The Citi template is 16:9; only use it for the 16:9 minimalist path.
        self._use_template = bool(template_path) and self.theme.use_template and self.aspect_ratio == "16:9"
        self._active_dark = self.theme.background == THEMES["dark"].background
        self.canvas = PptxCanvas(
            canvas_width=self._canvas_w,
            canvas_height=self._canvas_h,
            logical_width=self._LOGICAL_WIDTH,
            theme=self.theme,
            active_theme=self._active_theme,
            map_color=self._map_color,
        )

    def _ix(self, value: float) -> Inches:
        """Map a logical X position or width onto the target canvas."""
        return self.canvas.ix(value)

    def _iy(self, value: float) -> Inches:
        """Y positions and heights are unscaled (canvas height is fixed at 10in)."""
        return self.canvas.iy(value)

    def render(self, slides: list[SlideData]) -> bytes:
        prs = Presentation(self.template_path) if self._use_template else Presentation()

        if self._use_template:
            # Delete all original slides in the template
            id_list = prs.slides._sldIdLst
            for i in range(len(prs.slides) - 1, -1, -1):
                rId = id_list[i].rId
                prs.part.drop_rel(rId)
                id_list.remove(id_list[i])
        else:
            prs.slide_width = Inches(self._canvas_w)
            prs.slide_height = Inches(self._canvas_h)

        for slide_data in slides:
            layout = self._find_layout(prs, slide_data.layout) if self._use_template else self._blank_layout(prs)
            slide = prs.slides.add_slide(layout)
            self._set_slide_style(slide_data)
            self._apply_background(slide)

            if slide_data.index == 1 or slide_data.layout == "title":
                self._apply_title_slide(slide, slide_data)
            else:
                self._apply_content_slide(slide, slide_data)

            self._add_speaker_notes(slide, slide_data)
            self._add_brand_header(slide, prs.slide_width)

        buf = io.BytesIO()
        prs.save(buf)
        buf.seek(0)
        return buf.read()

    def _set_slide_style(self, data: SlideData) -> None:
        self._active_dark = self.theme.background == THEMES["dark"].background or self._uses_dark_variant(data)

    def _uses_dark_variant(self, data: SlideData) -> bool:
        layout = data.layout.lower()
        if layout == "title":
            return True
        return (self._variant_for(data) or "") in self._DARK_VARIANTS

    def _active_theme(self) -> Theme:
        return THEMES["dark"] if self._active_dark else self.theme

    @contextmanager
    def _theme_mode(self, *, dark: bool) -> Iterator[None]:
        previous = self._active_dark
        self._active_dark = dark
        try:
            yield
        finally:
            self._active_dark = previous

    def _map_color(self, color: RGBColor) -> RGBColor:
        if not self._active_dark:
            return color
        dark = THEMES["dark"]
        if color == self.theme.text:
            return dark.text
        if color == self.theme.muted:
            return dark.muted
        if color == self.theme.strong:
            return dark.strong
        if color == self.theme.panel_bg:
            return dark.panel_bg
        if color == self.theme.panel_border:
            return dark.panel_border
        if color == self.theme.surface:
            return dark.surface
        if color == self.theme.border:
            return dark.border
        if color == self.theme.accent_soft:
            return dark.accent_soft
        return color

    def _apply_background(self, slide: Slide) -> None:
        background = self._active_theme().background
        if background is None:
            return
        fill = slide.background.fill
        fill.solid()
        fill.fore_color.rgb = background
        self._add_background_depth(slide)

    def _add_background_depth(self, slide: Slide) -> None:
        theme = self._active_theme()
        if self._active_dark:
            panels = [
                (13.25, 0, 4.53, self._canvas_h, theme.panel_bg),
                (0, 8.95, self._LOGICAL_WIDTH, 1.05, RGBColor(0x0B, 0x16, 0x26)),
            ]
        else:
            panels = [
                (0, 0, self._LOGICAL_WIDTH, 1.1, theme.panel_bg),
                (0, 8.85, self._LOGICAL_WIDTH, 1.15, RGBColor(0xF1, 0xF5, 0xF9)),
            ]
        for x, y, w, h, color in panels:
            panel = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, self._ix(x), self._iy(y), self._ix(w), self._iy(h))
            panel.fill.solid()
            panel.fill.fore_color.rgb = color
            panel.line.fill.background()
            panel.shadow.inherit = False

    def _blank_layout(self, prs: PresentationType) -> SlideLayout:
        return prs.slide_layouts[6] if len(prs.slide_layouts) > 6 else prs.slide_layouts[0]

    def _find_layout(self, prs: PresentationType, layout_name: str) -> SlideLayout:
        for layout in prs.slide_layouts:
            if layout_name.lower() in layout.name.lower():
                return layout
        return prs.slide_layouts[0]

    def _apply_content_slide(self, slide: Slide, data: SlideData) -> None:
        self._current_slide_data = data
        self._apply_content_body(slide, data)
        self._add_chapter_marker(slide, data)

    def _apply_content_body(self, slide: Slide, data: SlideData) -> None:
        layout = data.layout.lower()
        variant = self._variant_for(data)
        layout_handlers = self._content_layout_handlers()
        variant_handlers = self._framework_variant_handlers()

        if layout == "section_divider":
            layout_handlers[layout](slide, data)
            return
        if data.chart_data:
            self._apply_chart_slide(slide, data)
            return
        if variant in variant_handlers:
            variant_handlers[variant](slide, data)
            return
        if getattr(data, "image_b64", None):
            self._apply_split_image(slide, data)
            return
        if getattr(data, "blocks", None):
            self._apply_blocks(slide, data)
            return

        handler = layout_handlers.get(layout, self._apply_standard_content)
        handler(slide, data)

    def _content_layout_handlers(self):
        return {
            "section_divider": self._apply_section_divider,
            "executive_summary": self._apply_executive_summary,
            "next_steps": self._apply_next_steps,
            "chart": self._apply_chart_slide,
        }

    def _framework_variant_handlers(self):
        from app.services.presentation.variants import _VARIANTS

        return {name: getattr(self, func.__name__) for name, func in _VARIANTS.items()}

    def _variant_for(self, data: SlideData) -> str | None:
        variant = (getattr(data, "variant", None) or "").lower().strip()
        if variant:
            return variant
        blocks = [block for block in (getattr(data, "blocks", None) or []) if isinstance(block, dict)]
        if len(blocks) > 1:
            return None
        block = self._first_block(data)
        btype = str(block.get("type", "")).lower() if block else ""
        if btype == "stat":
            return "big_stat"
        if btype == "quote":
            return "quote"
        if btype in ("table", "comparison"):
            return "comparison_table"
        if btype in ("process", "steps", "timeline"):
            return "process"
        if btype in ("cards", "card_grid", "grid", "columns"):
            return "three_points"
        if getattr(data, "image_b64", None):
            return "split_image"
        return None

    def _first_block(self, data: SlideData) -> dict | None:
        blocks = getattr(data, "blocks", None) or []
        for block in blocks:
            if isinstance(block, dict):
                return block
        return None

    def _apply_framework_variant(self, slide: Slide, data: SlideData, variant: str) -> None:
        handler = self._framework_variant_handlers().get(variant, self._apply_standard_content)
        handler(slide, data)

    def _add_speaker_notes(self, slide: Slide, data: SlideData) -> None:
        if not data.notes:
            return
        notes_frame = slide.notes_slide.notes_text_frame
        if notes_frame is None:
            return
        notes_frame.text = data.notes
        if getattr(data, "narrative_context", None):
            notes_frame.text += "\n\n---\n\n" + data.narrative_context

    def _add_bullets_box(self, slide: Slide, bullets: list[str], left: float, top: float, width: float, height: float) -> None:
        return self.canvas.add_bullets_box(slide, bullets, left, top, width, height)

    def _add_number_circle(self, slide: Slide, x: float, y: float, number: int, size: float = 0.72) -> None:
        return self.canvas.add_number_circle(slide, x, y, number, size)

    def _add_slide_image(self, slide: Slide, image_b64, left: float, top: float, width: float, height: float) -> bool:
        """Insert a base64 image at the given rect. Returns True on success."""
        return self.canvas.add_slide_image(slide, image_b64, left, top, width, height)

    def _add_fullbleed_image(self, slide: Slide, image_b64) -> bool:
        """Insert a base64 image covering the entire slide. Returns True on success."""
        return self._add_slide_image(slide, image_b64, 0, 0, self._LOGICAL_WIDTH, self._canvas_h)

    def _add_scrim(self, slide: Slide, color: RGBColor = RGBColor(0x0A, 0x16, 0x28), opacity: int = 58) -> None:
        """Overlay a semi-transparent dark rectangle over the whole slide for contrast."""
        return self.canvas.add_scrim(slide, color, opacity)

    def _add_card(self, slide: Slide, x: float, y: float, w: float, h: float, border_color: RGBColor | None = None):
        """A surface-filled rounded rectangle with a hairline border."""
        return self.canvas.add_card(slide, x, y, w, h, border_color=border_color)

    def _add_icon_chip(self, slide: Slide, x: float, y: float, size: float = 0.62, icon: str | None = None) -> None:
        """Accent-tinted rounded square with a centered Font Awesome icon (accent color)."""
        return self.canvas.add_icon_chip(slide, x, y, size, icon)

    def _add_eyebrow(self, slide: Slide, x: float, y: float, text: str, width: float = 12.0) -> None:
        """Small uppercase accent-colored label (the template's red kicker)."""
        label = text.upper()
        pill_w = min(max(1.45, len(label) * 0.105 + 0.55), width)
        pill = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, self._ix(x), self._iy(y - 0.05), self._ix(pill_w), self._iy(0.42))
        pill.fill.solid()
        pill.fill.fore_color.rgb = self._active_theme().accent_soft
        pill.line.color.rgb = self._active_theme().accent
        pill.shadow.inherit = False
        try:
            pill.adjustments[0] = 0.35
        except (IndexError, TypeError, ValueError):
            pass
        self._add_text(slide, x + 0.18, y + 0.03, max(pill_w - 0.36, 0.5), 0.24, label, 12, self._active_theme().accent, bold=True)

    def _add_card_text(self, slide: Slide, x: float, y: float, w: float, h: float, text: str, size: int, color: RGBColor):
        """Markdown-aware text that stays within a fixed box (no autosize growth)."""
        return self.canvas.add_card_text(slide, x, y, w, h, text, size, color)

    def _add_accent_bar(self, slide: Slide, left: float, top: float, width: float, height: float = LAYOUT.accent_rule_height) -> None:
        return self.canvas.add_accent_bar(slide, left, top, width, height)

    def _add_vertical_divider(self, slide: Slide, left: float, top: float, height: float) -> None:
        return self.canvas.add_vertical_divider(slide, left, top, height)

    def _add_text(
        self,
        slide: Slide,
        left: float,
        top: float,
        width: float,
        height: float,
        text: str,
        size: int,
        color: RGBColor,
        *,
        bold: bool = False,
        align: "PP_ALIGN | None" = None,
    ):
        return self.canvas.add_text(
            slide,
            left,
            top,
            width,
            height,
            text,
            size,
            color,
            bold=bold,
            align=align,
        )

    def _add_brand_header(self, slide: Slide, slide_width: int) -> None:
        width = Inches(1.0)
        height = Inches(0.65)
        left = max(Inches(0.5), slide_width - Inches(0.5) - width)
        top = Inches(0.18)

        if settings.citi_logo_path and Path(settings.citi_logo_path).exists():
            slide.shapes.add_picture(settings.citi_logo_path, left, top, width=width)
            return

        text_box = slide.shapes.add_textbox(left, top, width, height)
        paragraph = text_box.text_frame.paragraphs[0]
        paragraph.text = "citi"
        paragraph.font.name = DISPLAY_FONT
        paragraph.font.size = Pt(16)
        paragraph.font.bold = True
        paragraph.font.color.rgb = self._active_theme().accent
        paragraph.alignment = PP_ALIGN.RIGHT
