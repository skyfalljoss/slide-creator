import re

from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.slide import Slide

from app.models.schemas import SlideData
from app.services.pptx_layout import LAYOUT
from app.services.pptx_text import (
    is_leaked_instruction as _is_leaked_instruction,
    looks_like_instruction as _looks_like_instruction,
    table_row_from_text as _table_row_from_text,
)
from app.services.pptx_theme import THEMES


class PptxLayoutMixin:
    """Slide-specific layout renderers for the PPTX engine."""

    def _apply_title_slide(self, slide: Slide, data: SlideData) -> None:
        self._add_slide_image(slide, getattr(data, "image_b64", None), 10.9, 0.95, 6.88, 8.1)
        kicker = getattr(data, "kicker", None)
        if kicker:
            self._add_eyebrow(slide, LAYOUT.left_margin, 3.45, kicker)
        self._add_text(slide, LAYOUT.left_margin, 3.9, 9.7, 1.8, data.title, 48, self.theme.text, bold=True)
        self._add_accent_bar(slide, LAYOUT.left_margin, 5.85, LAYOUT.accent_rule_width, height=0.06)
        subtitle = data.subtitle or (data.bullets[0] if data.bullets else "")
        secondary = data.bullets if data.subtitle else data.bullets[1:]
        if subtitle:
            self._add_text(slide, LAYOUT.left_margin, 6.1, 9.0, 0.8, subtitle, 22, self.theme.muted)
        if secondary:
            self._add_text(slide, LAYOUT.left_margin, 7.0, 9.0, 0.5, "  |  ".join(secondary), 15, self.theme.muted)

    def _apply_big_statement(self, slide: Slide, data: SlideData) -> None:
        label = getattr(data, "kicker", None) or "CORE INSIGHT"
        self._add_eyebrow(slide, LAYOUT.left_margin, 2.25, label)
        self._add_text(slide, LAYOUT.left_margin, 2.9, 14.8, 2.7, data.title, 50, self.theme.text, bold=True)
        self._add_accent_bar(slide, LAYOUT.left_margin, 5.95, 1.7, height=0.06)
        if data.subtitle:
            self._add_text(slide, LAYOUT.left_margin, 6.25, 11.5, 0.75, data.subtitle, 22, self.theme.muted)
        elif data.bullets:
            self._add_text(slide, LAYOUT.left_margin, 6.25, 11.5, 0.75, data.bullets[0], 22, self.theme.muted)

    def _apply_three_points(self, slide: Slide, data: SlideData) -> None:
        top = self._add_content_header(slide, data.title, getattr(data, "kicker", None), getattr(data, "subtitle", None))
        block = self._first_block(data)
        if block and str(block.get("type", "")).lower() in ("cards", "card_grid", "grid", "columns"):
            self._block_cards(slide, block, LAYOUT.left_margin, top, LAYOUT.content_width, LAYOUT.content_bottom - top)
            return
        items = [{"title": item, "body": ""} for item in data.bullets[:3]]
        self._block_cards(
            slide,
            {"type": "cards", "columns": min(max(len(items), 1), 3), "items": items},
            LAYOUT.left_margin,
            top,
            LAYOUT.content_width,
            LAYOUT.content_bottom - top,
        )

    def _apply_split_image(self, slide: Slide, data: SlideData) -> None:
        left = LAYOUT.left_margin
        top = 1.55
        panel = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, self._ix(0), self._iy(0), self._ix(9.05), self._iy(self._canvas_h))
        panel.fill.solid()
        panel.fill.fore_color.rgb = THEMES["dark"].background
        panel.line.fill.background()
        with self._theme_mode(dark=True):
            self._add_eyebrow(slide, left, top, getattr(data, "kicker", None) or "CURRENT STATE")
            self._add_text(slide, left, top + 0.55, 7.9, 1.7, data.title, 34, self.theme.text, bold=True)
            text_top = top + 2.55
            if data.subtitle:
                self._add_text(slide, left, text_top, 7.2, 0.8, data.subtitle, 19, self.theme.muted)
                text_top += 1.0
            if data.bullets:
                self._add_bullets_box(slide, data.bullets[:4], left, text_top, 7.2, 3.0)

        img_left = 9.7
        img_top = 1.6
        img_w = 7.25
        img_h = 6.95
        if not self._add_slide_image(slide, getattr(data, "image_b64", None), img_left, img_top, img_w, img_h):
            self._add_image_placeholder(slide, img_left, img_top, img_w, img_h)

    def _apply_big_stat(self, slide: Slide, data: SlideData) -> None:
        block = self._first_block(data) or {}
        value = str(block.get("value") or block.get("number") or "")
        label = str(block.get("label") or block.get("caption") or "")
        if not value:
            source = " ".join([data.title, *data.bullets])
            match = re.search(r"[$#]?\d+(?:\.\d+)?\s*(?:%|x|M|B|K)?", source)
            if match:
                value = match.group(0)
            else:
                self._apply_big_statement(slide, data)
                return
        if not label:
            label = data.subtitle or (data.bullets[0] if data.bullets else "")
        self._add_text(slide, LAYOUT.left_margin, 2.35, LAYOUT.content_width, 1.9, value, 130, self.theme.accent, bold=True, align=PP_ALIGN.CENTER)
        if label:
            self._add_text(slide, 2.4, 4.65, self._LOGICAL_WIDTH - 4.8, 0.9, label, 28, self.theme.text, bold=True, align=PP_ALIGN.CENTER)
        if data.title and data.title != value:
            self._add_text(slide, 2.8, 5.75, self._LOGICAL_WIDTH - 5.6, 0.7, data.title, 20, self.theme.muted, align=PP_ALIGN.CENTER)

    def _apply_before_after(self, slide: Slide, data: SlideData) -> None:
        top = self._add_content_header(slide, data.title, getattr(data, "kicker", None), getattr(data, "subtitle", None))
        block = self._first_block(data) or {}
        before_title = "Before"
        after_title = "After"
        before_items: list[str] = []
        after_items: list[str] = []
        rows = block.get("rows") if isinstance(block, dict) else None
        headers = block.get("headers") if isinstance(block, dict) else None
        if isinstance(headers, list) and len(headers) >= 2:
            before_title = str(headers[-2])
            after_title = str(headers[-1])
        if isinstance(rows, list) and rows:
            for row in rows:
                if isinstance(row, (list, tuple)) and len(row) >= 2:
                    before_items.append(str(row[-2]))
                    after_items.append(str(row[-1]))
        if not before_items and data.bullets:
            mid = max(1, len(data.bullets) // 2)
            before_items = data.bullets[:mid]
            after_items = data.bullets[mid:] or data.bullets[:mid]
        self._comparison_card(slide, LAYOUT.left_margin, top, 7.85, 4.9, before_title, before_items, muted=True)
        self._comparison_card(slide, 9.1, top, 7.85, 4.9, after_title, after_items, muted=False)

    def _apply_comparison_table(self, slide: Slide, data: SlideData) -> None:
        top = self._add_content_header(slide, data.title, getattr(data, "kicker", None), getattr(data, "subtitle", None))
        block = self._first_block(data)
        if block and str(block.get("type", "")).lower() in ("table", "comparison"):
            self._block_table(slide, block, LAYOUT.left_margin, top + 0.15, LAYOUT.content_width, 8.6 - top)
            return
        rows = [_table_row_from_text(item) for item in data.bullets]
        self._block_table(slide, {"headers": ["Topic", "Detail"], "rows": rows}, LAYOUT.left_margin, top + 0.15, LAYOUT.content_width, 8.6 - top)

    def _apply_process_variant(self, slide: Slide, data: SlideData) -> None:
        top = self._add_content_header(slide, data.title, getattr(data, "kicker", None), getattr(data, "subtitle", None))
        block = self._first_block(data)
        if block and str(block.get("type", "")).lower() in ("process", "steps", "timeline"):
            self._block_process(slide, block, LAYOUT.left_margin, top + 0.35, LAYOUT.content_width, 4.2)
            return
        self._block_process(
            slide,
            {"steps": [{"title": item, "body": ""} for item in data.bullets[:4]]},
            LAYOUT.left_margin,
            top + 0.35,
            LAYOUT.content_width,
            4.2,
        )

    def _apply_quote_variant(self, slide: Slide, data: SlideData) -> None:
        block = self._first_block(data) or {}
        quote = str(block.get("text") or block.get("quote") or (data.bullets[0] if data.bullets else data.title))
        author = str(block.get("author") or block.get("attribution") or "")
        self._add_eyebrow(slide, LAYOUT.left_margin, 1.75, getattr(data, "kicker", None) or "VISION")
        bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, self._ix(LAYOUT.left_margin), self._iy(2.65), self._ix(0.12), self._iy(3.1))
        bar.fill.solid()
        bar.fill.fore_color.rgb = self._active_theme().accent
        bar.line.fill.background()
        self._add_text(slide, 1.45, 2.55, 13.8, 2.7, f"\"{quote}\"", 34, self.theme.text, bold=True)
        if author:
            self._add_text(slide, 1.45, 5.65, 10.0, 0.45, author.upper(), 16, self.theme.muted)
        if getattr(data, "image_b64", None):
            self._add_slide_image(slide, data.image_b64, 12.9, 5.85, 3.6, 2.15)

    def _apply_closing(self, slide: Slide, data: SlideData) -> None:
        label = getattr(data, "kicker", None) or "NEXT STEPS"
        self._add_eyebrow(slide, LAYOUT.left_margin, 2.6, label)
        self._add_text(slide, LAYOUT.left_margin, 3.1, 12.8, 2.1, data.title, 48, self.theme.text, bold=True)
        self._add_accent_bar(slide, LAYOUT.left_margin, 5.55, 1.5, height=0.06)
        if data.subtitle:
            self._add_text(slide, LAYOUT.left_margin, 5.85, 9.5, 0.7, data.subtitle, 22, self.theme.muted)
        block = self._first_block(data)
        if block and str(block.get("type", "")).lower() in ("process", "steps", "timeline"):
            self._block_process(slide, block, LAYOUT.left_margin, 6.45, LAYOUT.content_width, 2.0)
        elif data.bullets:
            self._add_bullets_box(slide, data.bullets[:3], LAYOUT.left_margin, 6.35, 9.0, 1.6)

    def _comparison_card(self, slide: Slide, x: float, y: float, w: float, h: float, title: str, items: list[str], *, muted: bool) -> None:
        border = self.theme.muted if muted else self.theme.accent
        self._add_card(slide, x, y, w, h, border_color=border)
        self._add_text(slide, x + 0.45, y + 0.45, w - 0.9, 0.45, title, 24, border, bold=True)
        text = "\n".join(items or ["No details available."])
        self._add_card_text(slide, x + 0.45, y + 1.15, w - 0.9, h - 1.45, text, 18, self.theme.muted if muted else self.theme.text)

    def _add_image_placeholder(self, slide: Slide, left: float, top: float, width: float, height: float) -> None:
        rect = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, self._ix(left), self._iy(top), self._ix(width), self._iy(height))
        rect.fill.solid()
        rect.fill.fore_color.rgb = self._active_theme().surface
        rect.line.color.rgb = self._active_theme().border
        rect.line.dash_style = 4
        rect.shadow.inherit = False
        self._add_text(slide, left + 0.55, top + height / 2 - 0.25, width - 1.1, 0.5, "Image Placeholder", 18, self.theme.muted, align=PP_ALIGN.CENTER)

    def _apply_section_divider(self, slide: Slide, data: SlideData) -> None:
        number = str(max(data.index - 4, 1)) + "."
        self._guard_divider_text(data.title)
        subtitle = data.subtitle or (data.bullets[0] if data.bullets else "")
        if not subtitle or _is_leaked_instruction(subtitle):
            subtitle = "Section overview"
        self._guard_divider_text(subtitle)

        image_b64 = getattr(data, "image_b64", None)
        if image_b64 and self._add_fullbleed_image(slide, image_b64):
            self._add_scrim(slide)
            light = RGBColor(0xFF, 0xFF, 0xFF)
            accent_light = RGBColor(0x8A, 0xC4, 0xEC)
            mid = self._LOGICAL_WIDTH / 2
            box_w = self._LOGICAL_WIDTH - 3.0
            self._add_text(slide, 1.5, 3.0, box_w, 1.0, number, 40, accent_light, bold=True, align=PP_ALIGN.CENTER)
            self._add_text(slide, 1.5, 3.9, box_w, 1.0, data.title, 40, light, bold=True, align=PP_ALIGN.CENTER)
            self._add_accent_bar(slide, mid - 0.7, 5.05, LAYOUT.accent_rule_width, height=LAYOUT.accent_rule_height)
            self._add_text(slide, 1.5, 5.25, box_w, 0.6, subtitle, 20, light, align=PP_ALIGN.CENTER)
            return

        self._add_text(slide, LAYOUT.left_margin, 2.78, 2.78, 0.97, number, 44, self.theme.strong, bold=True)
        self._add_accent_bar(slide, LAYOUT.left_margin, 3.89, 1.11, height=0.06)
        self._add_text(slide, LAYOUT.left_margin, 4.17, 9.7, 0.76, data.title, 36, self.theme.text, bold=True)
        self._add_text(slide, LAYOUT.left_margin, 5.14, 8.33, 0.42, subtitle, 20, self.theme.muted)

    def _guard_divider_text(self, text: str) -> None:
        """Raise if divider text reads like a leaked generation instruction."""
        if _looks_like_instruction(text):
            raise ValueError(
                "Section divider text reads like a generation instruction, "
                f"not client-facing content: {text[:80]!r}..."
            )

    def _apply_executive_summary(self, slide: Slide, data: SlideData) -> None:
        top = self._add_content_header(slide, data.title, getattr(data, "kicker", None) or "EXECUTIVE SUMMARY")

        bullets = data.bullets or []
        mid = (len(bullets) + 1) // 2 if bullets else 0
        columns = [
            ("Overview & Context", bullets[:mid] or ["No overview details available."]),
            ("Key Objectives & Focus Areas", bullets[mid:] or ["No focus area details available."]),
        ]
        gap = LAYOUT.content_gap
        usable = LAYOUT.content_width
        cw = (usable - gap) / 2
        ch = LAYOUT.content_bottom - top
        for i, (heading, items) in enumerate(columns):
            x = LAYOUT.left_margin + i * (cw + gap)
            self._add_card(slide, x, top, cw, ch)
            self._add_icon_chip(slide, x + 0.45, top + 0.45)
            self._add_text(slide, x + 0.45, top + 1.3, cw - 0.9, 0.45, heading, 22, self.theme.text, bold=True)
            self._add_card_text(slide, x + 0.45, top + 2.0, cw - 0.9, ch - 2.4, "\n".join(items), 18, self.theme.muted)

    def _apply_next_steps(self, slide: Slide, data: SlideData) -> None:
        top = self._add_content_header(slide, data.title, getattr(data, "kicker", None) or "NEXT STEPS")
        bullets = data.bullets or []
        steps = bullets[:3]
        if not steps:
            self._add_visual_panel(slide, data, LAYOUT.left_margin, top, LAYOUT.content_width, LAYOUT.content_bottom - top)
            return

        gap = LAYOUT.content_gap
        usable = LAYOUT.content_width
        cols = len(steps)
        cw = (usable - gap * (cols - 1)) / cols
        ch = 3.6
        for i, bullet in enumerate(steps):
            x = LAYOUT.left_margin + i * (cw + gap)
            self._add_card(slide, x, top, cw, ch)
            self._add_number_circle(slide, x + 0.45, top + 0.45, i + 1)
            self._add_card_text(slide, x + 0.45, top + 1.5, cw - 0.9, ch - 1.8, bullet, 18, self.theme.text)

        lower = bullets[3:]
        if lower:
            sy = top + ch + 0.5
            self._add_eyebrow(slide, LAYOUT.left_margin, sy, "Recommended Timeline")
            self._add_card_text(slide, LAYOUT.left_margin, sy + 0.5, usable, max(LAYOUT.content_bottom - (sy + 0.5), 0.6), "  •  ".join(lower), 18, self.theme.muted)

    def _apply_chart_slide(self, slide: Slide, data: SlideData) -> None:
        top = self._add_content_header(slide, data.title, getattr(data, "kicker", None))
        self._add_bullets_box(slide, data.bullets, LAYOUT.left_margin, top, 6.0, 8.6 - top)
        self._add_vertical_divider(slide, 7.1, top, 8.4 - top)
        ch_h = 8.0 - top
        if data.chart_data:
            self._add_chart(slide, data.chart_data, left=7.5, top=top, width=9.44, height=ch_h)
            self._add_source_note(slide, data.chart_audit, 7.5, top + ch_h + 0.1, 9.44)
        elif not self._add_slide_image(slide, getattr(data, "image_b64", None), 7.5, top, 9.44, ch_h):
            self._add_visual_panel(slide, data, 7.5, top, 9.44, ch_h)

    def _apply_standard_content(self, slide: Slide, data: SlideData) -> None:
        top = self._add_content_header(slide, data.title, getattr(data, "kicker", None))
        if data.bullets:
            self._add_card_grid(slide, data.bullets, top)
        else:
            self._add_visual_panel(slide, data, LAYOUT.left_margin, top, LAYOUT.content_width, LAYOUT.content_bottom - top)

    def _add_visual_panel(self, slide: Slide, data: SlideData, left: float, top: float, width: float, height: float) -> None:
        card = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, self._ix(left), self._iy(top), self._ix(width), self._iy(height))
        card.fill.solid()
        card.fill.fore_color.rgb = self._active_theme().panel_bg
        card.line.color.rgb = self._active_theme().panel_bg

        stripe = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, self._ix(left), self._iy(top), self._ix(0.08), self._iy(height))
        stripe.fill.solid()
        stripe.fill.fore_color.rgb = self._active_theme().accent
        stripe.line.color.rgb = self._active_theme().accent

        self._add_text(slide, left + 0.35, top + 0.35, width - 0.7, 0.35, "Visual Direction & Context", 16, self.theme.accent, bold=True)
        text = data.visual_direction or "Use a clean Citi-style callout visual for the slide's core message."
        self._add_card_text(slide, left + 0.35, top + 0.9, width - 0.7, height - 1.2, text, 16, self.theme.text)

    def _add_content_header(self, slide: Slide, title: str, kicker: str | None, subtitle: str | None = None) -> float:
        """Eyebrow + title (+ optional subtitle) + accent rule. Returns content top."""
        if kicker:
            self._add_eyebrow(slide, LAYOUT.left_margin, 0.62, kicker)
            self._add_text(slide, LAYOUT.left_margin, 1.02, LAYOUT.header_title_width, 0.9, title, 30, self.theme.text, bold=True)
        else:
            self._add_text(slide, LAYOUT.left_margin, 0.66, LAYOUT.header_title_width, 0.9, title, 30, self.theme.text, bold=True)
        rule_top = 1.78 if kicker else 1.6
        self._add_accent_bar(slide, LAYOUT.left_margin, rule_top, LAYOUT.accent_rule_width, height=LAYOUT.accent_rule_height)
        top = rule_top + 0.45
        if subtitle:
            self._add_text(slide, LAYOUT.left_margin, top, LAYOUT.header_title_width, 0.5, subtitle, 18, self.theme.muted)
            top += 0.65
        return top

    def _add_card_grid(self, slide: Slide, items: list[str], top: float) -> None:
        """Lay bullets out as a responsive grid of icon-chip surface cards."""
        n = len(items)
        if n == 0:
            return
        cols = 1 if n == 1 else 2 if n == 2 else 2 if n == 4 else 3
        rows = (n + cols - 1) // cols
        left = LAYOUT.left_margin
        gap = LAYOUT.content_gap
        usable = LAYOUT.content_width
        cw = (usable - gap * (cols - 1)) / cols
        bottom = LAYOUT.content_bottom
        ch = min((bottom - top - gap * (rows - 1)) / rows, 3.1)
        pad = 0.4
        for i, item in enumerate(items):
            row, col = divmod(i, cols)
            x = left + col * (cw + gap)
            y = top + row * (ch + gap)
            self._add_card(slide, x, y, cw, ch)
            self._add_icon_chip(slide, x + pad, y + pad, icon=item)
            self._add_card_text(slide, x + pad, y + pad + 0.85, cw - pad * 2, ch - pad - 0.95, item, 18, self.theme.text)
